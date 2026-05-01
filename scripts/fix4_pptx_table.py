"""
Debug and fix the table row addition for Slide 18.
"""
import sys
import copy
from pptx import Presentation
from pptx.oxml.ns import qn, nsmap
from lxml import etree

PPTX_PATH = "publication/thesis/defense_presentation.pptx"

def safe_print(s):
    sys.stdout.buffer.write((s + "\n").encode("utf-8", errors="replace"))
    sys.stdout.buffer.flush()

def find_tbl(shape):
    """Find the a:tbl element inside a table shape (graphic frame)."""
    # Try direct
    elem = shape._element
    tbl = elem.find(qn("a:tbl"))
    if tbl is not None:
        return tbl
    # Try through graphic data
    # The structure is: p:graphicFrame > a:graphic > a:graphicData > a:tbl
    for child in elem.iter():
        if child.tag == qn("a:tbl"):
            return child
    return None

def set_cell_text(row_elem, col_idx, text):
    cells = row_elem.findall(qn("a:tc"))
    if col_idx >= len(cells):
        return
    tc = cells[col_idx]
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return
    paras = txBody.findall(qn("a:p"))
    if paras:
        p = paras[0]
        # Remove all runs
        for r in p.findall(qn("a:r")):
            p.remove(r)
        # Add a fresh run with the text
        r_new = etree.SubElement(p, qn("a:r"))
        t_new = etree.SubElement(r_new, qn("a:t"))
        t_new.text = text

def main():
    prs = Presentation(PPTX_PATH)
    slide18 = prs.slides[17]

    # Find the table shape
    tbl_shape = None
    for shape in slide18.shapes:
        if shape.shape_type == 19:  # TABLE
            tbl_shape = shape
            break

    if tbl_shape is None:
        safe_print("ERROR: No table found on slide 18")
        return

    # Debug: print element tag and first children
    safe_print(f"Table shape element tag: {tbl_shape._element.tag}")
    for i, child in enumerate(tbl_shape._element):
        safe_print(f"  child[{i}]: {child.tag}")

    tbl = find_tbl(tbl_shape)
    if tbl is None:
        safe_print("ERROR: Could not find a:tbl element")
        # Print full XML for debug
        safe_print(etree.tostring(tbl_shape._element, pretty_print=True).decode("utf-8", errors="replace")[:3000])
        return

    safe_print(f"Found a:tbl: {tbl.tag}")

    # Check existing rows
    rows = tbl.findall(qn("a:tr"))
    safe_print(f"Existing rows: {len(rows)}")
    row_first_cells = []
    for row in rows:
        cells = row.findall(qn("a:tc"))
        if cells:
            txBody = cells[0].find(qn("a:txBody"))
            if txBody is not None:
                paras = txBody.findall(qn("a:p"))
                if paras:
                    t = paras[0].find(qn("a:r"))
                    row_first_cells.append(t.text if t is not None else "")

    safe_print(f"First cell texts: {row_first_cells}")

    # Add Phase 7 if missing
    if "7" not in row_first_cells:
        last_row = rows[-1]
        new_row = copy.deepcopy(last_row)
        set_cell_text(new_row, 0, "7")
        set_cell_text(new_row, 1, "DSPy Optimization")
        set_cell_text(new_row, 2, "+Automated prompt search")
        set_cell_text(new_row, 3, "BootstrapFewShot compiled program")
        set_cell_text(new_row, 4, "vs Phase 2")
        tbl.append(new_row)
        safe_print("Added Phase 7 row")
        rows = tbl.findall(qn("a:tr"))  # Refresh

    # Add Phase 8 if missing
    row_first_cells_after = []
    for row in tbl.findall(qn("a:tr")):
        cells = row.findall(qn("a:tc"))
        if cells:
            txBody = cells[0].find(qn("a:txBody"))
            if txBody is not None:
                paras = txBody.findall(qn("a:p"))
                if paras:
                    runs = paras[0].findall(qn("a:r"))
                    row_first_cells_after.append(runs[0].text if runs else "")

    if "8" not in row_first_cells_after:
        rows = tbl.findall(qn("a:tr"))
        last_row = rows[-1]
        new_row = copy.deepcopy(last_row)
        set_cell_text(new_row, 0, "8")
        set_cell_text(new_row, 1, "RAG (BM25)")
        set_cell_text(new_row, 2, "+Retrieved examples")
        set_cell_text(new_row, 3, "Top-5 similar (OCR,GT) pairs from training")
        set_cell_text(new_row, 4, "vs Phase 2")
        tbl.append(new_row)
        safe_print("Added Phase 8 row")

    # Verify
    safe_print("\nTable AFTER:")
    for row in tbl_shape.table.rows:
        row_texts = [cell.text_frame.text.strip() for cell in row.cells]
        safe_print(f"  {row_texts}")

    prs.save(PPTX_PATH)
    safe_print(f"\nSaved: {PPTX_PATH}")

if __name__ == "__main__":
    main()
