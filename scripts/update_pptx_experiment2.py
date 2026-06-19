#!/usr/bin/env python3
"""Update MS_Masters_Progress_Report_v3.pptx with Experiment 2 results.

Changes:
  - Slide 11: Replace 'Next Step 1' with 'Experiment 2 Results'
  - Slide 14: Update Summary bullet points to include Exp2 finding
  - Slide 1 / footer: Update date to June 2026

Saves output as MS_Masters_Progress_Report_v4.pptx.
"""
from __future__ import annotations

import copy
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"D:\Masters\Arabic-Post-OCR-Correction\publication\thesis\MS_Masters_Progress_Report_v3.pptx"
DST = r"D:\Masters\Arabic-Post-OCR-Correction\publication\thesis\MS_Masters_Progress_Report_v4.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_tf_text(tf, text: str, font_size_pt: float = 14, bold: bool = False,
                color: tuple | None = None, font_size: float | None = None) -> None:
    if font_size is not None:
        font_size_pt = font_size
    """Replace all text in a text frame with new text + optional formatting."""
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)


def add_textbox(slide, left_in, top_in, width_in, height_in,
                text: str, size_pt: float = 14, bold: bool = False,
                color: tuple = (0x33, 0x33, 0x33), italic: bool = False) -> None:
    txb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor(*color)


def add_table(slide, left_in, top_in, width_in, rows, cols, data,
              header_bg=(0x1F, 0x49, 0x7D), header_fg=(0xFF, 0xFF, 0xFF),
              row_bg_a=(0xEE, 0xF2, 0xFF), row_bg_b=(0xFF, 0xFF, 0xFF),
              font_size=11.0) -> None:
    """Add a simple table with header row."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    table = slide.shapes.add_table(
        rows, cols,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches((rows) * 0.28),
    ).table

    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = data[0][c]
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(*header_fg)
        _set_cell_bg(cell, header_bg)
        para.alignment = PP_ALIGN.CENTER

    for r in range(1, rows):
        bg = row_bg_a if r % 2 == 1 else row_bg_b
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            _set_cell_bg(cell, bg)
            para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT


def _set_cell_bg(cell, rgb_tuple: tuple) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "{:02X}{:02X}{:02X}".format(*rgb_tuple))


def delete_all_shapes(slide) -> None:
    """Remove every shape from a slide."""
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            sp_tree.remove(sp)


# ---------------------------------------------------------------------------
# Slide 11 — Experiment 2 Results
# ---------------------------------------------------------------------------

DARK_BLUE = (0x1F, 0x49, 0x7D)
ACCENT    = (0x2E, 0x75, 0xB6)
GREEN     = (0x10, 0x73, 0x36)
RED       = (0xC0, 0x00, 0x00)
WHITE     = (0xFF, 0xFF, 0xFF)
GREY80    = (0x33, 0x33, 0x33)


def build_slide11(slide) -> None:
    """Replace slide 11 content with Experiment 2 Results."""
    delete_all_shapes(slide)

    # Header bar
    from pptx.util import Inches
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.AUTO_SHAPE
        Inches(0), Inches(0), Inches(13.33), Inches(1.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*DARK_BLUE)
    bar.line.fill.background()

    # Title
    add_textbox(slide, 0.15, 0.05, 11.0, 0.8,
                "Experiment 2 — Prompt Design Study (8 Trials)",
                size_pt=28, bold=True, color=WHITE)

    add_textbox(slide, 0.15, 0.75, 11.0, 0.35,
                "11 / 15  ·  Testing 4 prompt styles × 2 strategies (zero-shot vs RAG)  ·  13 datasets  ·  4,664 samples  ·  Qwen3-4B",
                size_pt=11, color=(0xCC, 0xCC, 0xFF))

    # ---- Left column: trial table ----
    add_textbox(slide, 0.15, 1.15, 6.5, 0.35,
                "Trial Configurations",
                size_pt=14, bold=True, color=DARK_BLUE)

    trial_data = [
        ["Trial", "Prompt Style", "Strategy", "Status"],
        ["T1  base_zs",      "Base (structured)",          "Zero-shot", "Slight regression"],
        ["T2  base_rag",     "Base (structured)",          "RAG (BM25)", "≈ Phase 8"],
        ["T3  cons_zs",      "Conservative (copy-first)",  "Zero-shot", "✓ BEST"],
        ["T4  cons_rag",     "Conservative (copy-first)",  "RAG (BM25)", "Partial failures"],
        ["T5  hp_zs",        "HP + error patterns",        "Zero-shot", "Similar to T1"],
        ["T6  hp_rag",       "HP + error patterns",        "RAG (BM25)", "Good KHATT"],
        ["T7  hp_cons_zs",   "HP + conservative",          "Zero-shot", "⚠ FAILED"],
        ["T8  hp_cons_rag",  "HP + conservative",          "RAG (BM25)", "Partial failures"],
    ]
    add_table(slide, 0.15, 1.50, 6.5, 9, 4, trial_data,
              font_size=9.5)

    # T7 note
    add_textbox(slide, 0.15, 4.20, 6.5, 0.4,
                "⚠  T7 failed: prompt-template tags leaked into LLM output (100-300%+ CER on PATS).",
                size_pt=9, italic=True, color=RED)

    # ---- Right column: results table ----
    add_textbox(slide, 6.8, 1.15, 6.3, 0.35,
                "Results Summary (Normal samples, no diacritics)",
                size_pt=14, bold=True, color=DARK_BLUE)

    results_data = [
        ["Config",       "PATS avg CER", "KHATT CER"],
        ["OCR baseline", "5.57%",        "34.24%"],
        ["Phase 8 RAG",  "5.45%",        "33.83%"],
        ["T1 base_zs",   "5.84%",        "33.25%"],
        ["T2 base_rag",  "5.44%",        "33.82%"],
        ["T3 cons_zs *", "5.27% ↓",      "32.93% ↓"],
        ["T4 cons_rag",  "27.02% ⚠",     "32.69%"],
        ["T5 hp_zs",     "5.85%",        "32.90%"],
        ["T6 hp_rag",    "19.64% ⚠",     "32.64%"],
        ["T7 hp_cons_zs","190.6% ✗",     "51.64% ✗"],
        ["T8 hp_cons_rag","46.6% ⚠",     "32.54% ↓"],
    ]
    add_table(slide, 6.8, 1.50, 6.3, 11, 3, results_data,
              font_size=10.0)

    add_textbox(slide, 6.8, 4.70, 6.3, 0.4,
                "* T3 Conservative zero-shot: best on PATS (5.27% avg) — beats OCR baseline AND Phase 8 RAG without retrieval.",
                size_pt=9, italic=True, color=GREEN)

    # ---- Bottom: key takeaways ----
    add_textbox(slide, 0.15, 5.00, 12.8, 0.35,
                "Key Takeaways",
                size_pt=13, bold=True, color=DARK_BLUE)

    takeaways = [
        ("★  Conservative zero-shot (T3) BEST for PATS:  5.57% → 5.27% avg CER  (−5.4% vs OCR baseline; −3.3% vs Phase 8 RAG).  No retrieval needed.",
         GREEN),
        ("★  For KHATT (handwritten), T8/T6/T4 are best:  34.24% → 32.54-32.69%  (−5% vs OCR baseline; slight gain over Phase 8 33.83%).",
         ACCENT),
        ("★  HP error-pattern prompts cause catastrophic failures when combined with conservative mode (T7) or RAG (T4/T8).  Prompt overload risk.",
         RED),
        ("◉  New datasets (KHATT-Paragraph 256/440 runaway, Yarmouk 69%+, Muharaf/Historical 230%+ CER) are unsolvable with current approach.",
         GREY80),
    ]
    for i, (text, color) in enumerate(takeaways):
        add_textbox(slide, 0.30, 5.40 + i * 0.43, 12.5, 0.38,
                    text, size_pt=10, color=color)


# ---------------------------------------------------------------------------
# Slide 14 — Update Summary
# ---------------------------------------------------------------------------

def update_slide14(slide) -> None:
    """Update Summary slide bullet points."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if "Next: conservative prompt + Qwen3-14B" in t:
            set_tf_text(
                shape.text_frame,
                "Exp 2: Conservative prompt beats Phase 8 RAG on PATS",
                font_size=15, bold=True, color=DARK_BLUE,
            )
        elif "p2v3 to reduce false positives" in t:
            set_tf_text(
                shape.text_frame,
                "T3 (cons_zs): 5.27% avg PATS CER — best configuration tested. "
                "No RAG needed; simpler and faster than Phase 8.",
                font_size=13,
            )
        elif "Mohamed Sabry" in t and "May 2026" in t:
            set_tf_text(
                shape.text_frame,
                "Mohamed Sabry · Master's Thesis · June 2026",
                font_size=12,
            )


# ---------------------------------------------------------------------------
# Slide 1 — Update date
# ---------------------------------------------------------------------------

def update_slide1(slide) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame and "May 2026" in shape.text_frame.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "May 2026" in run.text:
                        run.text = run.text.replace("May 2026", "June 2026")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation(SRC)

    update_slide1(prs.slides[0])      # Slide 1
    build_slide11(prs.slides[10])     # Slide 11
    update_slide14(prs.slides[13])    # Slide 14

    prs.save(DST)
    print(f"Saved: {DST}")


if __name__ == "__main__":
    main()
