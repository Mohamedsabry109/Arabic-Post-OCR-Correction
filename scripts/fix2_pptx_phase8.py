"""
Second fix pass:
 - Remove erroneously-added bullets from title text boxes on slides 38, 40, 41, 42
 - Inspect slide 18 fully (all shape types)
 - Verify final state
"""
import sys
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = "publication/thesis/defense_presentation.pptx"

def safe_print(s):
    sys.stdout.buffer.write((s + "\n").encode("utf-8", errors="replace"))
    sys.stdout.buffer.flush()

def print_shape(slide_idx, shape):
    safe_print(f"  Slide {slide_idx+1} '{shape.name}' (type={shape.shape_type}):")
    if shape.has_text_frame:
        for i, para in enumerate(shape.text_frame.paragraphs):
            t = para.text
            if t.strip():
                safe_print(f"    para[{i}]: {t[:130]!r}")

def remove_paragraphs_containing(text_frame, substring):
    txBody = text_frame._txBody
    to_remove = []
    for p in txBody.findall(qn("a:p")):
        texts = "".join(r.text or "" for r in p.findall(qn("a:r")))
        if substring in texts:
            to_remove.append(p)
    for p in to_remove:
        txBody.remove(p)
    return len(to_remove)

def add_bullet_paragraph(text_frame, text, copy_from_para=None):
    if copy_from_para is not None:
        new_para_xml = copy.deepcopy(copy_from_para._p)
        for r in new_para_xml.findall(qn("a:r")):
            new_para_xml.remove(r)
        for br in new_para_xml.findall(qn("a:br")):
            new_para_xml.remove(br)
        ref_runs = copy_from_para._p.findall(qn("a:r"))
        if ref_runs:
            new_r = copy.deepcopy(ref_runs[0])
            t_elem = new_r.find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = text
            new_para_xml.append(new_r)
        text_frame._txBody.append(new_para_xml)
    else:
        para = text_frame.add_paragraph()
        para.text = text

def main():
    prs = Presentation(PPTX_PATH)
    safe_print(f"Loaded: {len(prs.slides)} slides")

    # ----------------------------------------------------------------
    # Fix slides 38, 40, 41, 42: remove bullets from TITLE text boxes
    # Title text boxes are identifiable as having a single short para
    # (just the slide title like "Future Work", "Conclusion", etc.)
    # ----------------------------------------------------------------
    TITLE_FIXES = {
        37: [  # Slide 38
            ("Future Work", "Improve RAG"),
        ],
        39: [  # Slide 40
            ("Final Ranking", "Phase 8 (RAG, BM25)"),
        ],
        40: [  # Slide 41
            ("Conclusion", "Domain-matched RAG"),
        ],
        41: [  # Slide 42
            ("Thesis Contributions", "BM25 RAG using domain-matched"),
        ],
    }

    for slide_idx, fixes in TITLE_FIXES.items():
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full = tf.text
            for title_keyword, bad_phrase in fixes:
                if title_keyword in full and bad_phrase in full:
                    n = remove_paragraphs_containing(tf, bad_phrase)
                    safe_print(
                        f"Slide {slide_idx+1} '{shape.name}': removed {n} erroneous bullet(s) "
                        f"containing '{bad_phrase[:40]}'"
                    )

    # ----------------------------------------------------------------
    # Fix slide 38 Future Work: also remove duplicates from content box
    # The RAG bullet may appear twice in TextBox 5
    # ----------------------------------------------------------------
    slide38 = prs.slides[37]
    for shape in slide38.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text
        if "4. Confidence-guided" in full and full.count("Improve RAG") > 1:
            # keep only first occurrence
            txBody = shape.text_frame._txBody
            paras = txBody.findall(qn("a:p"))
            seen = False
            for p in paras:
                texts = "".join(r.text or "" for r in p.findall(qn("a:r")))
                if "Improve RAG" in texts:
                    if seen:
                        txBody.remove(p)
                        safe_print(f"Slide 38 '{shape.name}': removed duplicate RAG bullet")
                    else:
                        seen = True

    # ----------------------------------------------------------------
    # Slide 18 (idx=17): Inspect ALL shape types, then add Phase 8 text
    # ----------------------------------------------------------------
    safe_print("\n--- Slide 18: full shape inspection ---")
    slide18 = prs.slides[17]
    for i, shape in enumerate(slide18.shapes):
        safe_print(f"  shape[{i}] '{shape.name}' type={shape.shape_type} has_tf={shape.has_text_frame}")
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    safe_print(f"    para[{j}]: {t[:130]!r}")

    # The phases list on slide 18 seems to be in a shape that has "Phase 1".
    # Let's try the shape that contains the key design note and add Phase 8 there
    # only if it's really the right shape.
    # Actually, from inspection slide 18 has:
    #   TextBox 5: "Key design: Phases 3-5 are ISOLATED..." — that's the design note
    # The actual phase list may be in a table or placeholder. Check for table shapes.
    for i, shape in enumerate(slide18.shapes):
        if shape.shape_type == 19:  # TABLE
            safe_print(f"  Slide 18 has TABLE at shape[{i}]")
            tbl = shape.table
            for row_idx, row in enumerate(tbl.rows):
                for col_idx, cell in enumerate(row.cells):
                    ct = cell.text_frame.text.strip()
                    if ct:
                        safe_print(f"    row[{row_idx}][{col_idx}]: {ct[:80]!r}")

    # Add Phase 8 as a new line in the design note textbox if Phase 8 not already there
    for shape in slide18.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text
        if "Key design" in full and "Phase 8" not in full:
            # Append a brief Phase 8 note
            tf = shape.text_frame
            paras = tf.paragraphs
            ref_para = None
            for p in paras:
                if p.text.strip():
                    ref_para = p
            add_bullet_paragraph(
                tf,
                "Phase 8: RAG -- BM25 char n-gram retrieval of similar training corrections as few-shot context.",
                copy_from_para=ref_para,
            )
            safe_print(f"Slide 18 '{shape.name}': appended Phase 8 note")

    # ----------------------------------------------------------------
    # Final verification: print key slides
    # ----------------------------------------------------------------
    safe_print("\n=== FINAL STATE ===")
    for idx in [6, 11, 12, 17, 18, 37, 39, 40, 41]:
        slide = prs.slides[idx]
        safe_print(f"\n--- Slide {idx+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print_shape(idx, shape)

    prs.save(PPTX_PATH)
    safe_print(f"\nSaved: {PPTX_PATH}")

if __name__ == "__main__":
    main()
