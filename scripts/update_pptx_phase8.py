"""
Script to add Phase 8 (RAG) content to defense_presentation.pptx
Run from project root:
  python scripts/update_pptx_phase8.py
"""
import copy
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

PPTX_PATH = "publication/thesis/defense_presentation.pptx"


def print_slide_shapes(prs, slide_idx):
    slide = prs.slides[slide_idx]
    print(f"\n--- Slide {slide_idx + 1} ---")
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            text = shape.text_frame.text[:120].replace("\n", " | ")
            print(f"  Shape {i} ({shape.shape_type}, name={shape.name!r}): {text!r}")


def replace_text_in_shape(shape, old, new):
    """Replace text while preserving formatting as much as possible."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    # Also check full text for cross-run matches — rebuild if needed
    if not changed:
        full = shape.text_frame.text
        if old in full:
            # Fallback: set text on first paragraph, first run
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run.text = run.text.replace(old, new)
                        changed = True
                        break
    return changed


def get_run_font_props(run):
    """Return a dict of key font properties from a run."""
    f = run.font
    return {
        "bold": f.bold,
        "italic": f.italic,
        "size": f.size,
        "color": f.color.rgb if f.color and f.color.type else None,
        "name": f.name,
    }


def apply_font_props(run, props):
    f = run.font
    if props.get("bold") is not None:
        f.bold = props["bold"]
    if props.get("italic") is not None:
        f.italic = props["italic"]
    if props.get("size") is not None:
        f.size = props["size"]
    if props.get("color") is not None:
        f.color.rgb = props["color"]
    if props.get("name") is not None:
        f.name = props["name"]


def add_bullet_paragraph(text_frame, text, level=0, copy_from_para=None):
    """Add a new paragraph with bullet text, copying format from copy_from_para."""
    from pptx.oxml.ns import qn
    from lxml import etree

    if copy_from_para is not None:
        # Deep-copy the reference paragraph XML
        new_para_xml = copy.deepcopy(copy_from_para._p)
        # Clear all runs
        for r in new_para_xml.findall(qn("a:r")):
            new_para_xml.remove(r)
        for br in new_para_xml.findall(qn("a:br")):
            new_para_xml.remove(br)
        # Copy one run for text
        ref_runs = copy_from_para._p.findall(qn("a:r"))
        if ref_runs:
            new_r = copy.deepcopy(ref_runs[0])
            # Set text
            t_elem = new_r.find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = text
            new_para_xml.append(new_r)
        text_frame._txBody.append(new_para_xml)
    else:
        from pptx.util import Pt
        para = text_frame.add_paragraph()
        para.text = text
        para.level = level


def main():
    prs = Presentation(PPTX_PATH)
    total_slides = len(prs.slides)
    print(f"Loaded PPTX: {total_slides} slides")

    # ----------------------------------------------------------------
    # DIAGNOSTIC: print shapes for slides we'll edit
    # ----------------------------------------------------------------
    for idx in [6, 11, 12, 17, 18, 37, 39, 40, 41]:
        if idx < total_slides:
            print_slide_shapes(prs, idx)

    # ----------------------------------------------------------------
    # Slide 7 (idx=6): Contributions — "7-phase" → "8-phase", add RAG
    # ----------------------------------------------------------------
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full = tf.text
        # Replace "7-phase" with "8-phase"
        for para in tf.paragraphs:
            for run in para.runs:
                if "7-phase" in run.text:
                    run.text = run.text.replace("7-phase", "8-phase")
                    print("Slide 7: replaced '7-phase' -> '8-phase'")
                if "five" in run.text.lower() and "strateg" in full.lower():
                    run.text = run.text.replace("five", "six").replace("Five", "Six")
                    print("Slide 7: replaced 'five strategies' → 'six strategies'")
        # If this shape contains the contributions list, add RAG item
        if "Systematic" in full or "experimental framework" in full.lower() or "contribution" in full.lower():
            # Check if RAG is already there
            if "RAG" not in full and "retrieval" not in full.lower():
                # Find a good reference paragraph (last bullet)
                paras = tf.paragraphs
                ref_para = None
                for p in paras:
                    if p.text.strip() and len(p.text.strip()) > 5:
                        ref_para = p
                add_bullet_paragraph(
                    tf,
                    "7. BM25 RAG from training corrections: domain-matched retrieval improves robustness",
                    copy_from_para=ref_para,
                )
                print("Slide 7: added RAG contribution bullet")

    # ----------------------------------------------------------------
    # Slide 12 (idx=11): "7 Research Questions Across 6 Experimental Phases"
    # ----------------------------------------------------------------
    slide12 = prs.slides[11]
    for shape in slide12.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if "7 Research Questions" in run.text:
                    run.text = run.text.replace(
                        "7 Research Questions", "8 Research Questions"
                    )
                    print("Slide 12: updated RQ count")
                if "6 Experimental Phases" in run.text:
                    run.text = run.text.replace(
                        "6 Experimental Phases", "7 Experimental Phases"
                    )
                    print("Slide 12: updated phase count")

    # ----------------------------------------------------------------
    # Slide 13 (idx=12): RQ list — append RQ8
    # ----------------------------------------------------------------
    slide13 = prs.slides[12]
    for shape in slide13.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full = tf.text
        if "RQ7" in full or "Overarching" in full:
            if "RQ8" not in full:
                # Find last meaningful paragraph as reference
                paras = tf.paragraphs
                ref_para = None
                for p in paras:
                    if "RQ7" in p.text or "Overarching" in p.text:
                        ref_para = p
                if ref_para is None:
                    for p in paras:
                        if p.text.strip():
                            ref_para = p
                add_bullet_paragraph(
                    tf,
                    "RQ8 (Phase 8): Does retrieval-augmented generation with domain-matched training corrections improve post-OCR correction?",
                    copy_from_para=ref_para,
                )
                print("Slide 13: added RQ8")

    # ----------------------------------------------------------------
    # Slide 18 (idx=17): Experimental Phases Overview — add Phase 8
    # ----------------------------------------------------------------
    slide18 = prs.slides[17]
    for shape in slide18.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full = tf.text
        if "Phase" in full and ("Baseline" in full or "Zero-Shot" in full or "CAMeL" in full):
            if "Phase 8" not in full and "RAG" not in full:
                paras = tf.paragraphs
                ref_para = None
                for p in paras:
                    if "Phase 7" in p.text or "DSPy" in p.text:
                        ref_para = p
                if ref_para is None:
                    for p in paras:
                        if p.text.strip():
                            ref_para = p
                add_bullet_paragraph(
                    tf,
                    "Phase 8: RAG — BM25 retrieval of similar training corrections as few-shot context",
                    copy_from_para=ref_para,
                )
                print("Slide 18: added Phase 8 to phases overview")

    # ----------------------------------------------------------------
    # Slide 19 (idx=18): Prompt Engineering Strategies — add Phase 8
    # ----------------------------------------------------------------
    slide19 = prs.slides[18]
    for shape in slide19.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full = tf.text
        # Look for shape with phase descriptions (Phase 5, 6 etc.)
        if ("Phase 5" in full or "Phase 6" in full or "CAMeL" in full) and "Phase 8" not in full:
            paras = tf.paragraphs
            ref_para = None
            for p in paras:
                if "Phase 6" in p.text or "Phase 7" in p.text or "CAMeL" in p.text or "DSPy" in p.text:
                    ref_para = p
            if ref_para is None:
                for p in paras:
                    if p.text.strip():
                        ref_para = p
            add_bullet_paragraph(
                tf,
                "Phase 8 (RAG): BM25 char n-gram index over Phase 2 training corrections; top-5 (OCR→GT) pairs injected as few-shot context",
                copy_from_para=ref_para,
            )
            print("Slide 19: added Phase 8 RAG strategy")

    # ----------------------------------------------------------------
    # Slide 38 (idx=37): Future Work — add RAG future work item
    # ----------------------------------------------------------------
    if total_slides > 37:
        slide38 = prs.slides[37]
        for shape in slide38.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full = tf.text
            if "future" in full.lower() or "limitation" in full.lower() or "selective" in full.lower():
                if "dense retrieval" not in full.lower() and "dense embed" not in full.lower():
                    paras = tf.paragraphs
                    ref_para = None
                    for p in paras:
                        if p.text.strip() and len(p.text.strip()) > 5:
                            ref_para = p
                    add_bullet_paragraph(
                        tf,
                        "Improve RAG: replace BM25 with dense Arabic embeddings; explore cross-dataset retrieval",
                        copy_from_para=ref_para,
                    )
                    print("Slide 38: added RAG future work")

    # ----------------------------------------------------------------
    # Slide 40 (idx=39): Final Ranking — add Phase 8
    # ----------------------------------------------------------------
    if total_slides > 39:
        slide40 = prs.slides[39]
        for shape in slide40.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full = tf.text
            if "ranking" in full.lower() or "rank" in full.lower() or "Phase 7" in full or "DSPy" in full:
                if "Phase 8" not in full and "RAG" not in full:
                    paras = tf.paragraphs
                    ref_para = None
                    for p in paras:
                        if "Phase 7" in p.text or "DSPy" in p.text or "rank" in p.text.lower():
                            ref_para = p
                    if ref_para is None:
                        for p in paras:
                            if p.text.strip():
                                ref_para = p
                    add_bullet_paragraph(
                        tf,
                        "Phase 8 (RAG, BM25): PATS-A01 CER 12.1% | KHATT CER 42.9% — modest improvement via domain-matched retrieval",
                        copy_from_para=ref_para,
                    )
                    print("Slide 40: added Phase 8 to ranking")

    # ----------------------------------------------------------------
    # Slide 41 (idx=40): Conclusion — add Phase 8 finding
    # ----------------------------------------------------------------
    if total_slides > 40:
        slide41 = prs.slides[40]
        for shape in slide41.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full = tf.text
            if "conclusion" in full.lower() or "finding" in full.lower() or "LLM" in full:
                if "RAG" not in full and "Phase 8" not in full:
                    paras = tf.paragraphs
                    ref_para = None
                    for p in paras:
                        if p.text.strip() and len(p.text.strip()) > 10:
                            ref_para = p
                    add_bullet_paragraph(
                        tf,
                        "Domain-matched RAG (Phase 8) provides additional modest gains, validating corpus alignment as a prerequisite for effective retrieval",
                        copy_from_para=ref_para,
                    )
                    print("Slide 41: added Phase 8 conclusion point")

    # ----------------------------------------------------------------
    # Slide 42 (idx=41): Contributions — add RAG contribution
    # ----------------------------------------------------------------
    if total_slides > 41:
        slide42 = prs.slides[41]
        for shape in slide42.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full = tf.text
            if ("contribution" in full.lower() or "over-correction" in full.lower()
                    or "complementarity" in full.lower() or "pipeline" in full.lower()):
                if "RAG" not in full and "BM25" not in full:
                    paras = tf.paragraphs
                    ref_para = None
                    for p in paras:
                        if p.text.strip() and len(p.text.strip()) > 5:
                            ref_para = p
                    add_bullet_paragraph(
                        tf,
                        "7. BM25 RAG using domain-matched training corrections as few-shot context for per-sample adaptive correction",
                        copy_from_para=ref_para,
                    )
                    print("Slide 42: added RAG contribution")

    # ----------------------------------------------------------------
    # Save
    # ----------------------------------------------------------------
    prs.save(PPTX_PATH)
    print(f"\nSaved updated PPTX to {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
