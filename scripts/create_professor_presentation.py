"""
Create professor meeting presentation: Problem + Results + Analysis + Next Steps
Run from project root:
  python scripts/create_professor_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy, os, datetime

OUTPUT_PATH = "publication/thesis/professor_meeting.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x38, 0x64)
BLUE    = RGBColor(0x27, 0x74, 0xAE)
LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)
GREEN   = RGBColor(0x37, 0x86, 0x36)
RED     = RGBColor(0xC0, 0x00, 0x00)
ORANGE  = RGBColor(0xED, 0x7D, 0x31)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY   = RGBColor(0x40, 0x40, 0x40)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
GOLD    = RGBColor(0xFF, 0xC0, 0x00)
TEAL    = RGBColor(0x00, 0x70, 0xC0)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE
    fill_shape(shape, color)
    shape.line.fill.background()
    return shape


def tf_para(tf, text, bold=False, size=18, color=WHITE, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return p


def add_textbox(slide, left, top, width, height,
                text="", bold=False, size=18, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return txb, tf


def add_table(slide, left, top, width, height, rows, cols):
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height)
    return tbl.table


def cell_text(cell, text, bold=False, size=13, color=DGRAY,
              align=PP_ALIGN.CENTER, bg: RGBColor = None):
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    # clear any existing runs
    for r in p.runs:
        r.text = ""
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color


def header_bar(slide, title: str, subtitle: str = ""):
    """Dark navy header bar across the top."""
    bar = add_rect(slide, 0, 0, W, Inches(1.25), NAVY)
    add_textbox(slide, Inches(0.35), Inches(0.10), Inches(12.0), Inches(0.65),
                text=title, bold=True, size=28, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(0.72), Inches(12.0), Inches(0.45),
                    text=subtitle, bold=False, size=16, color=LBLUE)


def slide_number(slide, num: int, total: int = 18):
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
                text=f"{num} / {total}", size=11, color=DGRAY, align=PP_ALIGN.RIGHT)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs, total):
    """Slide 1 – Title"""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, NAVY)
    add_rect(sl, 0, Inches(5.9), W, Inches(1.6), BLUE)

    add_textbox(sl, Inches(0.8), Inches(0.9), Inches(11.7), Inches(1.0),
                text="Arabic Post-OCR Correction with LLMs",
                bold=True, size=36, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.7),
                text="Can LLMs bridge the performance gap between open-source and closed-source Arabic OCR?",
                bold=False, size=20, color=LBLUE, align=PP_ALIGN.CENTER)

    add_textbox(sl, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.6),
                text="Progress Report — Professor Meeting",
                bold=False, size=18, color=GOLD, align=PP_ALIGN.CENTER)

    add_textbox(sl, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.5),
                text="Mohamed Sabry · Master's Thesis · May 2026",
                bold=False, size=15, color=LBLUE, align=PP_ALIGN.CENTER)

    add_textbox(sl, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.5),
                text="Model: Qwen3-4B-Instruct-2507  ·  Datasets: PATS-A01 (typewritten) + KHATT (handwritten)",
                bold=False, size=13, color=WHITE, align=PP_ALIGN.CENTER)


def slide_research_question(prs, n, total):
    """Slide 2 – Research Question & Motivation"""
    sl = blank_slide(prs)
    header_bar(sl, "Research Question & Motivation")
    slide_number(sl, n, total)

    # Problem box
    add_rect(sl, Inches(0.35), Inches(1.45), Inches(6.0), Inches(2.5), LGRAY)
    add_textbox(sl, Inches(0.5), Inches(1.5), Inches(5.7), Inches(0.4),
                text="The Problem", bold=True, size=17, color=NAVY)
    lines = [
        "• Arabic OCR (Qaari) is open-source & accessible",
        "• But makes systematic character-level errors:",
        "   - Similar-shape confusions (ب / ت / ث)",
        "   - Missing/extra diacritics",
        "   - Punctuation & spacing artefacts",
        "• Closed-source VLMs are far more accurate",
        "  but expensive and API-dependent",
    ]
    y = Inches(1.95)
    for line in lines:
        add_textbox(sl, Inches(0.5), y, Inches(5.7), Inches(0.33),
                    text=line, size=14, color=DGRAY)
        y += Inches(0.31)

    # RQ box
    add_rect(sl, Inches(6.6), Inches(1.45), Inches(6.38), Inches(1.3), NAVY)
    add_textbox(sl, Inches(6.75), Inches(1.55), Inches(6.1), Inches(1.1),
                text="Research Question\n"
                     "Can a small open-source LLM (4B) post-process Qaari OCR output "
                     "to approach the accuracy of closed-source VLMs?",
                bold=False, size=15, color=WHITE)

    # Hypothesis
    add_rect(sl, Inches(6.6), Inches(2.9), Inches(6.38), Inches(1.1), BLUE)
    add_textbox(sl, Inches(6.75), Inches(3.0), Inches(6.1), Inches(0.9),
                text="Hypothesis\n"
                     "LLMs can exploit linguistic context to fix OCR errors cheaply — "
                     "making open-source OCR competitive.",
                bold=False, size=14, color=WHITE)

    # Why it matters
    add_rect(sl, Inches(0.35), Inches(4.1), Inches(12.63), Inches(1.0), LBLUE)
    add_textbox(sl, Inches(0.5), Inches(4.15), Inches(12.4), Inches(0.8),
                text="Why It Matters:  Digitising Arabic manuscripts, newspapers, and historical archives at low cost "
                     "requires accurate open-source OCR — LLM post-correction is a practical bridge.",
                size=14, color=NAVY)

    # 8-phase answer
    add_textbox(sl, Inches(0.35), Inches(5.3), Inches(12.63), Inches(0.4),
                text="Approach:  8 experiments test progressively richer prompting strategies "
                     "(zero-shot → prompt engineering → RAG → DSPy optimisation)",
                bold=True, size=14, color=NAVY)


def slide_datasets(prs, n, total):
    """Slide 3 – Datasets & OCR Baseline"""
    sl = blank_slide(prs)
    header_bar(sl, "Datasets & OCR Baseline (Phase 1)", "Two domains: typewritten (PATS-A01) vs handwritten (KHATT)")
    slide_number(sl, n, total)

    # PATS block
    add_rect(sl, Inches(0.35), Inches(1.4), Inches(6.2), Inches(5.7), LGRAY)
    add_textbox(sl, Inches(0.5), Inches(1.45), Inches(5.9), Inches(0.4),
                text="PATS-A01  (Typewritten / Synthetic)", bold=True, size=16, color=NAVY)
    add_textbox(sl, Inches(0.5), Inches(1.85), Inches(5.9), Inches(0.35),
                text="8 Arabic fonts  ·  ~2766 pages/font  ·  85/15 train/val split",
                size=13, color=DGRAY)

    # PATS table
    tbl = add_table(sl, Inches(0.42), Inches(2.25), Inches(6.05), Inches(4.6), 10, 3)
    headers = ["Font", "Val CER", "Val WER"]
    col_w = [Inches(2.5), Inches(1.75), Inches(1.75)]
    for c, (h, w) in enumerate(zip(headers, col_w)):
        tbl.columns[c].width = w
        cell_text(tbl.cell(0, c), h, bold=True, size=13, color=WHITE, bg=NAVY)

    fonts = [
        ("Naskh",       "0.86%", "0.92%"),
        ("Thuluth",     "1.02%", "1.07%"),
        ("Arial",       "1.06%", "1.08%"),
        ("Akhbar",      "1.14%", "1.11%"),
        ("Simplified",  "1.32%", "1.33%"),
        ("Traditional", "1.18%", "1.22%"),
        ("Andalus",     "1.54%", "1.66%"),
        ("Tahoma",      "2.39%", "2.35%"),
        ("PATS avg",    "1.31%", "1.34%"),
    ]
    for r, (font, cer, wer) in enumerate(fonts, 1):
        bg = LBLUE if font == "PATS avg" else (LGRAY if r % 2 == 0 else WHITE)
        bold = font == "PATS avg"
        cell_text(tbl.cell(r, 0), font,  bold=bold, size=12, color=DGRAY, align=PP_ALIGN.LEFT, bg=bg)
        cell_text(tbl.cell(r, 1), cer,   bold=bold, size=12, color=DGRAY, bg=bg)
        cell_text(tbl.cell(r, 2), wer,   bold=bold, size=12, color=DGRAY, bg=bg)

    # KHATT block
    add_rect(sl, Inches(6.7), Inches(1.4), Inches(6.28), Inches(5.7), LGRAY)
    add_textbox(sl, Inches(6.85), Inches(1.45), Inches(6.0), Inches(0.4),
                text="KHATT  (Handwritten / Real)", bold=True, size=16, color=NAVY)
    add_textbox(sl, Inches(6.85), Inches(1.85), Inches(6.0), Inches(0.35),
                text="Writers of varied styles  ·  226 validation pages",
                size=13, color=DGRAY)

    khatt_data = [
        ("KHATT-validation", "5.66%", "7.25%", "17.7%"),
    ]
    add_textbox(sl, Inches(6.85), Inches(2.3), Inches(6.0), Inches(0.4),
                text="CER: 5.66%   |   WER: 7.25%   |   Runaway: 17.7%",
                bold=True, size=19, color=RED)

    add_textbox(sl, Inches(6.85), Inches(3.0), Inches(6.0), Inches(0.4),
                text="Key Observations:", bold=True, size=14, color=NAVY)
    obs = [
        "KHATT CER is 4.3x higher than PATS average",
        "Runaway samples (Qaari repetition bug):",
        "  KHATT: 17.7%,  Tahoma: 10.4%",
        "  Most PATS fonts: 3.5 – 6%",
        "Tahoma is hardest typewritten font (CER 2.39%)",
        "Naskh is easiest typewritten font (CER 0.86%)",
        "",
        "Two very different correction challenges:",
        "  PATS — already accurate, LLM risks over-correction",
        "  KHATT — high noise, LLM has more to gain",
    ]
    y = Inches(3.4)
    for o in obs:
        add_textbox(sl, Inches(6.85), y, Inches(6.0), Inches(0.31),
                    text=o, size=13, color=DGRAY)
        y += Inches(0.29)


def slide_experimental_design(prs, n, total):
    """Slide 4 – 8-Phase Experimental Design"""
    sl = blank_slide(prs)
    header_bar(sl, "8-Phase Experimental Design",
               "Each phase isolates one variable; all compare to Phase 2 (zero-shot) baseline")
    slide_number(sl, n, total)

    tbl = add_table(sl, Inches(0.35), Inches(1.35), Inches(12.63), Inches(5.85), 9, 4)
    col_ws = [Inches(1.0), Inches(2.5), Inches(4.5), Inches(4.63)]
    for c, w in enumerate(col_ws):
        tbl.columns[c].width = w

    headers = ["Phase", "Name", "What it tests", "Key signal added"]
    for c, h in enumerate(headers):
        cell_text(tbl.cell(0, c), h, bold=True, size=13, color=WHITE, bg=NAVY)

    rows = [
        ("1", "Baseline",          "Raw Qaari OCR quality, error taxonomy",               "None (reference)"),
        ("2", "Zero-Shot LLM",     "Can vanilla LLM fix OCR without any extra context?",  "None (base prompt p2v2)"),
        ("3", "OCR-Aware",         "Does confusion-matrix knowledge help?",                "Top-10 char confusions + LLM failure filter"),
        ("4", "Self-Reflective",   "Do error-pattern warnings reduce over-correction?",    "Training failure pairs + overcorrection alerts"),
        ("5", "CAMeL Validation",  "Does morphological revert save over-corrected words?", "CAMeL morph check + known-OC revert"),
        ("6", "Combinations",      "Do Phase 3+4 signals synergize?",                     "Confusion + self-reflective fused prompt"),
        ("7", "RAG (BM25)",        "Does per-sample retrieval of similar corrections help?","Top-5 (OCR→GT) pairs from training index"),
        ("8", "Error-Sig RAG",     "Does structural similarity retrieval beat surface?",   "CAMeL + confusion-matrix error signature"),
    ]
    row_colors = [LGRAY, WHITE] * 5
    for r, (ph, name, what, signal) in enumerate(rows, 1):
        bg = LBLUE if ph in ("8",) else (LGRAY if r % 2 == 1 else WHITE)
        bold = ph == "8"
        cell_text(tbl.cell(r, 0), ph,     bold=bold, size=12, color=NAVY,  bg=bg)
        cell_text(tbl.cell(r, 1), name,   bold=bold, size=12, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)
        cell_text(tbl.cell(r, 2), what,   bold=False, size=11, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)
        cell_text(tbl.cell(r, 3), signal, bold=False, size=11, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)


def slide_main_results(prs, n, total):
    """Slide 5 – Main Results Table (all phases)"""
    sl = blank_slide(prs)
    header_bar(sl, "Main Results — All Phases",
               "CER = Character Error Rate  ·  WER = Word Error Rate  ·  Lower is better")
    slide_number(sl, n, total)

    add_textbox(sl, Inches(0.35), Inches(1.3), Inches(12.63), Inches(0.32),
                text="Validation set  ·  Diacritics stripped before evaluation  ·  "
                     "Model: Qwen3-4B-Instruct-2507",
                size=12, color=DGRAY)

    tbl = add_table(sl, Inches(0.35), Inches(1.65), Inches(12.63), Inches(5.5), 9, 5)
    col_ws = [Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.93)]
    for c, w in enumerate(col_ws):
        tbl.columns[c].width = w

    headers = ["Phase", "PATS avg CER", "PATS avg WER", "KHATT CER", "KHATT WER"]
    for c, h in enumerate(headers):
        cell_text(tbl.cell(0, c), h, bold=True, size=13, color=WHITE, bg=NAVY)

    data = [
        ("Phase 1  (OCR Baseline)",    "1.31%", "1.34%", "5.66%", "7.25%",  False, False),
        ("Phase 2  (Zero-Shot)",        "1.32%", "1.38%", "3.21%", "4.53%",  False, False),
        ("Phase 3  (OCR-Aware)",        "1.47%", "1.51%", "3.33%", "3.59%",  False, False),
        ("Phase 4  (Self-Reflective)",  "1.45%", "1.49%", "3.02%", "3.44%",  False, False),
        ("Phase 6  (Conf+Self)",        "1.53%", "1.57%", "3.30%", "3.87%",  False, False),
        ("Phase 7  (RAG BM25)",         "1.29%", "1.32%", "2.77%", "3.28%",  True,  True),
        ("Phase 8  (Error-Sig RAG)",    "1.44%", "1.48%", "3.18%", "3.65%",  False, False),
        ("Phase 5  (CAMeL revert)",     "see note*", "—",  "3.26%", "4.64%", False, False),
    ]

    for r, (phase, pc, pw, kc, kw, best_p, best_k) in enumerate(data, 1):
        alt = LGRAY if r % 2 == 1 else WHITE
        bg = LBLUE if (best_p or best_k) else alt
        bold = best_p or best_k
        cell_text(tbl.cell(r, 0), phase, bold=bold, size=11, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)
        cell_text(tbl.cell(r, 1), pc,    bold=bold, size=12, color=(GREEN if best_p else DGRAY), bg=bg)
        cell_text(tbl.cell(r, 2), pw,    bold=bold, size=12, color=(GREEN if best_p else DGRAY), bg=bg)
        cell_text(tbl.cell(r, 3), kc,    bold=bold, size=12, color=(GREEN if best_k else DGRAY), bg=bg)
        cell_text(tbl.cell(r, 4), kw,    bold=bold, size=12, color=(GREEN if best_k else DGRAY), bg=bg)

    add_textbox(sl, Inches(0.35), Inches(7.15), Inches(12.63), Inches(0.28),
                text="* Phase 5 (CAMeL): only normal-subset metrics reported (excludes runaway samples).",
                size=10, color=DGRAY)


def slide_finding_rag_wins(prs, n, total):
    """Slide 6 – Key Finding: RAG Wins"""
    sl = blank_slide(prs)
    header_bar(sl, "Key Finding 1 — RAG (BM25) Is the Best Overall Strategy")
    slide_number(sl, n, total)

    # Big number callouts
    for i, (label, val, sub, col) in enumerate([
        ("KHATT CER reduction\nvs OCR baseline", "−51%", "5.66% → 2.77%", GREEN),
        ("KHATT CER reduction\nvs Zero-Shot (P2)", "−14%", "3.21% → 2.77%", BLUE),
        ("KHATT WER reduction\nvs OCR baseline", "−55%", "7.25% → 3.28%", GREEN),
        ("PATS avg CER\n(only phase to beat P2)", "1.29%", "P1=1.31%, P2=1.32%", TEAL),
    ]):
        x = Inches(0.35 + i * 3.25)
        add_rect(sl, x, Inches(1.35), Inches(3.1), Inches(2.2), NAVY)
        add_textbox(sl, x + Inches(0.1), Inches(1.45), Inches(2.9), Inches(0.5),
                    text=label, size=12, color=LBLUE, align=PP_ALIGN.CENTER)
        add_textbox(sl, x + Inches(0.1), Inches(1.95), Inches(2.9), Inches(0.8),
                    text=val, bold=True, size=34, color=col, align=PP_ALIGN.CENTER)
        add_textbox(sl, x + Inches(0.1), Inches(2.75), Inches(2.9), Inches(0.4),
                    text=sub, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Why RAG works
    add_rect(sl, Inches(0.35), Inches(3.75), Inches(12.63), Inches(0.4), BLUE)
    add_textbox(sl, Inches(0.5), Inches(3.8), Inches(12.4), Inches(0.35),
                text="Why RAG (BM25) Works", bold=True, size=16, color=WHITE)

    reasons = [
        ("Domain-matched examples", "Top-5 (OCR→GT) pairs from Phase 2 training corrections —\n"
                                    "same domain, same font family, same error patterns."),
        ("Grounded few-shot",       "Instead of abstract prompting, the model sees concrete\n"
                                    "examples of what the correction should look like."),
        ("Reduces hallucination",   "Ground-truth anchors prevent the model from inventing\n"
                                    "changes unsupported by training evidence."),
        ("BM25 char n-grams",       "Character-level matching captures OCR-specific confusable\n"
                                    "character patterns rather than word-level semantics."),
    ]
    for i, (title, body) in enumerate(reasons):
        x = Inches(0.4 + (i % 2) * 6.3)
        y = Inches(4.3 + (i // 2) * 1.3)
        add_rect(sl, x, y, Inches(6.1), Inches(1.15), LGRAY)
        add_textbox(sl, x + Inches(0.1), y + Inches(0.05), Inches(5.9), Inches(0.35),
                    text=title, bold=True, size=13, color=NAVY)
        add_textbox(sl, x + Inches(0.1), y + Inches(0.4), Inches(5.9), Inches(0.65),
                    text=body, size=12, color=DGRAY)


def slide_finding_prompt_fails(prs, n, total):
    """Slide 7 – Key Finding: Prompt Engineering Regresses on PATS"""
    sl = blank_slide(prs)
    header_bar(sl, "Key Finding 2 — Prompt Engineering Regresses on Typewritten Text (PATS)")
    slide_number(sl, n, total)

    add_textbox(sl, Inches(0.35), Inches(1.3), Inches(12.63), Inches(0.4),
                text="Phases 3, 4, 6 all increase PATS CER compared to zero-shot (Phase 2). "
                     "More prompt context = more incorrect changes on already-accurate text.",
                size=15, color=DGRAY)

    # Bar-style comparison
    tbl = add_table(sl, Inches(0.35), Inches(1.8), Inches(12.63), Inches(2.8), 5, 4)
    for c, w in enumerate([Inches(3.5), Inches(3.0), Inches(3.0), Inches(3.13)]):
        tbl.columns[c].width = w
    for c, h in enumerate(["Phase", "PATS avg CER", "vs Zero-Shot", "Interpretation"]):
        cell_text(tbl.cell(0, c), h, bold=True, size=13, color=WHITE, bg=NAVY)

    rows = [
        ("Phase 2 — Zero-Shot (baseline)", "1.32%", "—",       "Reference",                        False),
        ("Phase 3 — OCR-Aware",            "1.47%", "+0.15pp", "Confusion matrix → over-corrects",  True),
        ("Phase 4 — Self-Reflective",      "1.45%", "+0.13pp", "Failure warnings → more edits",     True),
        ("Phase 6 — Conf+Self",            "1.53%", "+0.21pp", "Combined signals → worst of both",  True),
    ]
    for r, (ph, cer, delta, interp, bad) in enumerate(rows, 1):
        bg = LGRAY if r % 2 == 1 else WHITE
        col_d = RED if bad else DGRAY
        cell_text(tbl.cell(r, 0), ph,     size=12, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)
        cell_text(tbl.cell(r, 1), cer,    size=12, color=DGRAY, bg=bg)
        cell_text(tbl.cell(r, 2), delta,  size=12, color=col_d, bg=bg, bold=bad)
        cell_text(tbl.cell(r, 3), interp, size=11, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)

    # Root cause
    add_rect(sl, Inches(0.35), Inches(4.75), Inches(12.63), Inches(0.4), ORANGE)
    add_textbox(sl, Inches(0.5), Inches(4.8), Inches(12.4), Inches(0.35),
                text="Root Cause: Over-Correction at Low Error Rate", bold=True, size=16, color=WHITE)

    causes = [
        ("The floor effect",       "PATS CER is already 0.86%–1.54%. At this level, any\n"
                                   "non-trivial correction strategy introduces more errors\n"
                                   "than it fixes. Signal-to-noise ratio is very low."),
        ("More context = more edits", "Confusion matrices and failure examples implicitly\n"
                                      "encourage the model to make changes — even when\n"
                                      "the OCR text is already correct."),
        ("Exception: Tahoma",      "CER 2.39% is high enough for LLM to help (Phase 2:\n"
                                   "2.70%), suggesting a per-dataset difficulty threshold\n"
                                   "below which LLM correction is net-harmful."),
    ]
    for i, (title, body) in enumerate(causes):
        x = Inches(0.4 + i * 4.3)
        add_rect(sl, x, Inches(5.3), Inches(4.1), Inches(1.9), LGRAY)
        add_textbox(sl, x + Inches(0.1), Inches(5.35), Inches(3.9), Inches(0.35),
                    text=title, bold=True, size=13, color=NAVY)
        add_textbox(sl, x + Inches(0.1), Inches(5.7), Inches(3.9), Inches(1.4),
                    text=body, size=12, color=DGRAY)


def slide_finding_khatt(prs, n, total):
    """Slide 8 – Key Finding: LLM Significantly Helps Handwritten Text"""
    sl = blank_slide(prs)
    header_bar(sl, "Key Finding 3 — LLM Post-Correction Strongly Benefits Handwritten Text (KHATT)")
    slide_number(sl, n, total)

    # Timeline of improvement
    phases_khatt = [
        ("OCR Baseline\n(Phase 1)", "5.66%", ORANGE),
        ("Zero-Shot\n(Phase 2)",    "3.21%", BLUE),
        ("Self-Reflect\n(Phase 4)", "3.02%", BLUE),
        ("RAG BM25\n(Phase 8)",     "2.77%", GREEN),
    ]

    add_textbox(sl, Inches(0.35), Inches(1.3), Inches(12.63), Inches(0.4),
                text="KHATT CER Progression   (WER follows similar pattern: 7.25% → 3.28%)",
                bold=True, size=16, color=NAVY)

    bar_max_h = Inches(3.2)
    max_val = 5.66
    for i, (label, cer, col) in enumerate(phases_khatt):
        val = float(cer.replace("%", ""))
        bar_h = bar_max_h * (val / max_val)
        x = Inches(1.3 + i * 2.9)
        y_bottom = Inches(5.1)
        add_rect(sl, x, y_bottom - bar_h, Inches(2.2), bar_h, col)
        add_textbox(sl, x, y_bottom - bar_h - Inches(0.4), Inches(2.2), Inches(0.35),
                    text=cer, bold=True, size=20, color=col, align=PP_ALIGN.CENTER)
        add_textbox(sl, x, y_bottom + Inches(0.05), Inches(2.2), Inches(0.55),
                    text=label, size=12, color=DGRAY, align=PP_ALIGN.CENTER)

    # Reduction callout
    add_rect(sl, Inches(0.35), Inches(5.75), Inches(12.63), Inches(0.5), NAVY)
    add_textbox(sl, Inches(0.5), Inches(5.8), Inches(12.4), Inches(0.4),
                text="Total CER reduction:  OCR → RAG =  5.66% → 2.77%  =  51% relative improvement  "
                     "(WER: 7.25% → 3.28% = 55% relative improvement)",
                bold=True, size=15, color=GOLD, align=PP_ALIGN.CENTER)

    note = ("KHATT is the real target domain: high error rate gives LLM enough signal to correct "
            "without excessive false positives. RAG amplifies this further via domain-matched examples.")
    add_textbox(sl, Inches(0.35), Inches(6.35), Inches(12.63), Inches(0.45),
                text=note, size=13, color=DGRAY)


def slide_per_font(prs, n, total):
    """Slide 9 – Per-Font Breakdown"""
    sl = blank_slide(prs)
    header_bar(sl, "Per-Font Breakdown — PATS-A01",
               "Not all fonts benefit equally; Tahoma uniquely resists LLM correction")
    slide_number(sl, n, total)

    tbl = add_table(sl, Inches(0.35), Inches(1.35), Inches(12.63), Inches(5.85), 10, 6)
    for c, w in enumerate([Inches(2.0), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(3.43)]):
        tbl.columns[c].width = w

    for c, h in enumerate(["Font", "P1 CER\n(OCR)", "P2 CER\n(ZeroShot)", "P8 CER\n(RAG)", "Best Phase", "Observation"]):
        cell_text(tbl.cell(0, c), h, bold=True, size=12, color=WHITE, bg=NAVY)

    fonts_data = [
        ("Naskh",       "0.86%", "1.01%", "0.90%", "P1",   "Already near-perfect; LLM hurts"),
        ("Thuluth",     "1.02%", "0.83%", "0.75%", "P8 ★", "RAG most effective here"),
        ("Arial",       "1.06%", "1.21%", "1.10%", "P1",   "LLM degrades accuracy"),
        ("Akhbar",      "1.14%", "0.87%", "0.89%", "P2",   "ZeroShot sufficient"),
        ("Traditional", "1.18%", "1.27%", "1.30%", "P1",   "Slight LLM regression"),
        ("Simplified",  "1.32%", "1.61%", "1.26%", "P8",   "RAG recovers from P2 regression"),
        ("Andalus",     "1.54%", "1.08%", "1.28%", "P2",   "ZeroShot helps most"),
        ("Tahoma",      "2.39%", "2.70%", "2.81%", "P1",   "No phase beats OCR baseline!"),
        ("PATS avg",    "1.31%", "1.32%", "1.29%", "P8",   "RAG is only phase to beat P1"),
    ]
    for r, (font, p1, p2, p8, best, obs) in enumerate(fonts_data, 1):
        is_avg = font == "PATS avg"
        is_bad = font == "Tahoma"
        bg = LBLUE if is_avg else (RGBColor(0xFF, 0xE6, 0xE6) if is_bad else (LGRAY if r % 2 == 1 else WHITE))
        cell_text(tbl.cell(r, 0), font,  bold=is_avg, size=12, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)
        cell_text(tbl.cell(r, 1), p1,   size=12, color=DGRAY, bg=bg)
        cell_text(tbl.cell(r, 2), p2,   size=12, color=DGRAY, bg=bg)
        cell_text(tbl.cell(r, 3), p8,   size=12, color=(GREEN if "★" in best or best=="P8" else DGRAY), bg=bg)
        cell_text(tbl.cell(r, 4), best, bold=True, size=12,
                  color=(GREEN if "P8" in best else (RED if best=="P1" else BLUE)), bg=bg)
        cell_text(tbl.cell(r, 5), obs,  size=11, color=DGRAY, bg=bg, align=PP_ALIGN.LEFT)


def slide_analysis_threshold(prs, n, total):
    """Slide 10 – Analysis: Error-Rate Threshold"""
    sl = blank_slide(prs)
    header_bar(sl, "Analysis — The Error-Rate Threshold Effect",
               "LLM correction is beneficial only when OCR error rate exceeds a critical threshold")
    slide_number(sl, n, total)

    add_textbox(sl, Inches(0.35), Inches(1.3), Inches(12.63), Inches(0.4),
                text="Hypothesis: below ~2% CER, LLM correction introduces more errors than it fixes.",
                bold=True, size=16, color=NAVY)

    # Zones
    for x, w, col, label, sub in [
        (Inches(0.35),  Inches(4.0),  RGBColor(0xFF,0xEB,0xEB), "HARMFUL ZONE\nCER < ~1.5%",
         "LLM makes text worse\nNaskh, Arial, Traditional\nBest strategy: return unchanged"),
        (Inches(4.5),   Inches(4.0),  RGBColor(0xFF,0xF5,0xE0), "MIXED ZONE\nCER 1.5 – 2.5%",
         "Depends on strategy\nAndalus, Simplified, Tahoma\nRAG and ZeroShot can help"),
        (Inches(9.0),   Inches(4.0),  RGBColor(0xE8,0xF5,0xE9), "BENEFICIAL ZONE\nCER > ~2.5%",
         "LLM significantly improves\nKHATT (5.66%), Tahoma (2.39%)\nRAG gives best results"),
    ]:
        add_rect(sl, x, Inches(1.85), w, Inches(3.5), col)
        add_textbox(sl, x+Inches(0.1), Inches(1.95), w-Inches(0.2), Inches(0.8),
                    text=label, bold=True, size=15, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, x+Inches(0.1), Inches(2.8), w-Inches(0.2), Inches(2.4),
                    text=sub, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

    add_textbox(sl, Inches(0.35), Inches(5.55), Inches(12.63), Inches(0.4),
                text="Implication for thesis:", bold=True, size=15, color=NAVY)
    impl = [
        "1. Dataset difficulty must be reported alongside CER/WER improvements.",
        "2. A selective-correction strategy (detect high-error samples, then apply LLM) could unlock gains on PATS.",
        "3. Phase 2 (zero-shot) already achieves diminishing returns — advanced methods mainly help KHATT.",
    ]
    y = Inches(6.05)
    for imp in impl:
        add_textbox(sl, Inches(0.5), y, Inches(12.4), Inches(0.3),
                    text=imp, size=13, color=DGRAY)
        y += Inches(0.3)


def slide_next_conservative_prompt(prs, n, total):
    """Slide 11 – Next: Conservative Prompt (p2v3)"""
    sl = blank_slide(prs)
    header_bar(sl, "Next Step 1 — Conservative Prompt (p2v3)",
               "Hypothesis: an even stricter 'do-nothing-when-unsure' policy will reduce PATS over-correction")
    slide_number(sl, n, total)

    # Current prompts
    add_rect(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.85), LGRAY)
    add_textbox(sl, Inches(0.5), Inches(1.4), Inches(5.7), Inches(0.35),
                text="Current prompts", bold=True, size=16, color=NAVY)

    add_textbox(sl, Inches(0.5), Inches(1.85), Inches(5.7), Inches(0.3),
                text="p2v1 (aggressive):", bold=True, size=13, color=RED)
    add_textbox(sl, Inches(0.5), Inches(2.15), Inches(5.7), Inches(0.65),
                text='"Fix OCR errors in this Arabic text. Return only the corrected text."',
                size=12, color=DGRAY)

    add_textbox(sl, Inches(0.5), Inches(2.9), Inches(5.7), Inches(0.3),
                text="p2v2 (conservative, current):", bold=True, size=13, color=ORANGE)
    add_textbox(sl, Inches(0.5), Inches(3.2), Inches(5.7), Inches(1.2),
                text='"You are an Arabic OCR corrector. If the text looks correct, return it unchanged. '
                     'Only fix clear OCR errors such as wrong characters, missing characters, '
                     'or obvious character confusions. Do not rephrase or paraphrase."',
                size=12, color=DGRAY)

    add_textbox(sl, Inches(0.5), Inches(4.5), Inches(5.7), Inches(0.35),
                text="Problem with p2v2:", bold=True, size=13, color=RED)
    add_textbox(sl, Inches(0.5), Inches(4.85), Inches(5.7), Inches(0.7),
                text="Still triggers too many edits on PATS fonts (Arial, Traditional, Naskh) "
                     "where OCR is already highly accurate. CER increases vs baseline.",
                size=12, color=DGRAY)

    # Proposed p2v3
    add_rect(sl, Inches(6.5), Inches(1.35), Inches(6.48), Inches(5.85), NAVY)
    add_textbox(sl, Inches(6.65), Inches(1.4), Inches(6.2), Inches(0.35),
                text="Proposed: p2v3 (ultra-conservative)", bold=True, size=16, color=GOLD)

    add_textbox(sl, Inches(6.65), Inches(1.85), Inches(6.2), Inches(2.2),
                text='"You are an Arabic OCR post-processor. Your ONLY job is to fix '
                     'unambiguous character-substitution errors caused by the OCR scanner — '
                     'where a character was clearly misread due to visual similarity. '
                     'Do NOT change: word choice, spelling variants, diacritics, punctuation, '
                     'spacing, or any word that could be a valid Arabic word. '
                     'If you are unsure whether something is an OCR error or correct Arabic, '
                     'leave it unchanged. Return ONLY the corrected text, nothing else."',
                size=12.5, color=WHITE)

    add_textbox(sl, Inches(6.65), Inches(4.15), Inches(6.2), Inches(0.35),
                text="Key differences from p2v2:", bold=True, size=13, color=GOLD)
    diffs = [
        "• Explicit list of what NOT to change",
        "• Uncertainty → do nothing (conservative default)",
        "• 'Valid Arabic word' test blocks grammar 'fixes'",
        "• Targets only visual confusion errors",
    ]
    y = Inches(4.55)
    for d in diffs:
        add_textbox(sl, Inches(6.65), y, Inches(6.2), Inches(0.28),
                    text=d, size=13, color=LBLUE)
        y += Inches(0.28)

    add_textbox(sl, Inches(6.65), Inches(5.9), Inches(6.2), Inches(0.4),
                text="Expected: PATS CER holds at baseline; KHATT may trade recall for precision",
                bold=True, size=12, color=GOLD)

    add_textbox(sl, Inches(0.35), Inches(7.12), Inches(12.63), Inches(0.3),
                text="Action: implement p2v3, run on all 9 datasets, compare to p2v2 and OCR baseline",
                size=12, color=DGRAY, bold=True)


def slide_next_bigger_model(prs, n, total):
    """Slide 12 – Next: Bigger LLM"""
    sl = blank_slide(prs)
    header_bar(sl, "Next Step 2 — Larger LLM Model",
               "Hypothesis: a larger model has better Arabic orthographic knowledge → fewer over-corrections")
    slide_number(sl, n, total)

    # Current
    add_rect(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(3.5), LGRAY)
    add_textbox(sl, Inches(0.5), Inches(1.4), Inches(5.7), Inches(0.35),
                text="Current Model", bold=True, size=16, color=NAVY)
    for label, val in [
        ("Model",        "Qwen3-4B-Instruct-2507"),
        ("Parameters",   "4 Billion"),
        ("Architecture", "Dense transformer"),
        ("Platform",     "Kaggle (T4 GPU, ~15 GB VRAM)"),
        ("Speed",        "~3-5 samples/second"),
        ("Arabic skill", "Good — trained on multilingual data"),
        ("Weakness",     "Limited context retention on long Arabic\n"
                         "lines; borderline word-validity decisions"),
    ]:
        add_textbox(sl, Inches(0.5), Inches(1.85 + list(zip(*[["Model","Parameters","Architecture","Platform","Speed","Arabic skill","Weakness"]]))[0].index(label) * 0.38),
                    Inches(1.8), Inches(0.35), text=label+":", bold=True, size=12, color=DGRAY)
        add_textbox(sl, Inches(2.35), Inches(1.85 + list(zip(*[["Model","Parameters","Architecture","Platform","Speed","Arabic skill","Weakness"]]))[0].index(label) * 0.38),
                    Inches(3.9), Inches(0.35), text=val, size=12, color=DGRAY)

    # Proposed models
    add_rect(sl, Inches(6.5), Inches(1.35), Inches(6.48), Inches(5.85), NAVY)
    add_textbox(sl, Inches(6.65), Inches(1.4), Inches(6.2), Inches(0.35),
                text="Candidate Larger Models", bold=True, size=16, color=GOLD)

    models = [
        ("Qwen3-14B",       "14B dense",    "Kaggle (A100)",  "Direct scale-up; 3.5x more params"),
        ("Qwen3-30B-A3B",   "30B MoE",      "Kaggle (A100)",  "Efficient: only 3B active params"),
        ("Qwen3-32B",       "32B dense",    "Kaggle 2xA100",  "Largest Qwen3 accessible"),
        ("Llama-3.3-70B-Q4","70B quant",    "Colab Pro+",     "Strong Arabic; quantized to 4-bit"),
    ]

    y = Inches(1.9)
    for model, size, hw, note in models:
        add_rect(sl, Inches(6.6), y, Inches(6.28), Inches(1.1), BLUE)
        add_textbox(sl, Inches(6.75), y+Inches(0.05), Inches(6.0), Inches(0.35),
                    text=model, bold=True, size=14, color=WHITE)
        add_textbox(sl, Inches(6.75), y+Inches(0.4), Inches(6.0), Inches(0.28),
                    text=f"Size: {size}  ·  Platform: {hw}", size=12, color=LBLUE)
        add_textbox(sl, Inches(6.75), y+Inches(0.68), Inches(6.0), Inches(0.28),
                    text=note, size=11, color=LGRAY)
        y += Inches(1.22)

    add_textbox(sl, Inches(6.65), Inches(6.75), Inches(6.2), Inches(0.3),
                text="Recommendation: Qwen3-14B first (same family, same tokenizer, easy comparison)",
                bold=True, size=12, color=GOLD)

    # What we expect
    add_rect(sl, Inches(0.35), Inches(5.0), Inches(5.9), Inches(2.2), LGRAY)
    add_textbox(sl, Inches(0.5), Inches(5.05), Inches(5.7), Inches(0.35),
                text="Expected Impact of Larger Model", bold=True, size=14, color=NAVY)
    exps = [
        ("PATS", "Better morphological knowledge → fewer false corrections → CER may drop toward P1"),
        ("KHATT", "Stronger context modelling → more accurate error detection"),
        ("RAG",  "Higher reasoning capacity → better use of retrieved examples"),
    ]
    y = Inches(5.5)
    for group, exp in exps:
        add_textbox(sl, Inches(0.5), y, Inches(0.8), Inches(0.35),
                    text=group+":", bold=True, size=12, color=NAVY)
        add_textbox(sl, Inches(1.35), y, Inches(4.8), Inches(0.35),
                    text=exp, size=12, color=DGRAY)
        y += Inches(0.38)


def slide_next_bigger_model_fixed(prs, n, total):
    """Slide 12 – Next: Bigger LLM (simplified, no comprehension issues)"""
    sl = blank_slide(prs)
    header_bar(sl, "Next Step 2 — Larger LLM Model",
               "Hypothesis: a larger model has better Arabic orthographic knowledge → fewer over-corrections")
    slide_number(sl, n, total)

    # Current model box
    add_rect(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.85), LGRAY)
    add_textbox(sl, Inches(0.5), Inches(1.4), Inches(5.7), Inches(0.35),
                text="Current Model", bold=True, size=16, color=NAVY)

    current_rows = [
        ("Model",        "Qwen3-4B-Instruct-2507"),
        ("Parameters",   "4 Billion (dense)"),
        ("Platform",     "Kaggle T4 GPU (~15 GB VRAM)"),
        ("Speed",        "~3–5 samples / second"),
        ("Arabic skill", "Good — multilingual training data"),
        ("Limitation",   "Borderline word-validity decisions;\n"
                         "limited orthographic certainty → over-corrects"),
    ]
    y = Inches(1.9)
    for label, val in current_rows:
        add_textbox(sl, Inches(0.5), y, Inches(1.8), Inches(0.38),
                    text=label + ":", bold=True, size=12, color=DGRAY)
        add_textbox(sl, Inches(2.4), y, Inches(3.6), Inches(0.38),
                    text=val, size=12, color=DGRAY)
        y += Inches(0.42)

    # Expected impact box
    add_rect(sl, Inches(0.35), Inches(4.95), Inches(5.9), Inches(2.25), RGBColor(0xE8,0xF5,0xE9))
    add_textbox(sl, Inches(0.5), Inches(5.0), Inches(5.7), Inches(0.35),
                text="Expected Impact of Larger Model", bold=True, size=14, color=NAVY)
    for group, exp in [
        ("PATS:  ", "Better morphological knowledge → fewer false positives → CER closer to baseline"),
        ("KHATT: ", "Stronger context modelling → more accurate error detection → lower CER"),
        ("RAG:   ", "Higher reasoning capacity → better use of retrieved correction examples"),
    ]:
        y += Inches(0.0)
    y = Inches(5.45)
    for group, exp in [
        ("PATS:  ", "Better morphological knowledge → fewer false\npositives → CER closer to baseline"),
        ("KHATT: ", "Stronger context modelling → better error detection"),
        ("RAG:   ", "Higher reasoning capacity → better use of examples"),
    ]:
        add_textbox(sl, Inches(0.5), y, Inches(1.1), Inches(0.5),
                    text=group, bold=True, size=12, color=NAVY)
        add_textbox(sl, Inches(1.6), y, Inches(4.5), Inches(0.5),
                    text=exp, size=12, color=DGRAY)
        y += Inches(0.55)

    # Candidate models box
    add_rect(sl, Inches(6.5), Inches(1.35), Inches(6.48), Inches(5.85), NAVY)
    add_textbox(sl, Inches(6.65), Inches(1.4), Inches(6.2), Inches(0.35),
                text="Candidate Larger Models", bold=True, size=16, color=GOLD)

    candidates = [
        ("Qwen3-14B",         "14B dense",   "Kaggle A100",
         "Recommended: same family as current;\n3.5x parameters, direct comparison"),
        ("Qwen3-30B-A3B (MoE)","30B / 3B active","Kaggle A100",
         "Efficient MoE; near 14B quality\nat fraction of compute"),
        ("Qwen3-32B",         "32B dense",   "2x Kaggle A100",
         "Largest accessible Qwen3;\nmay need quantization"),
        ("Llama-3.3-70B-Q4",  "70B (4-bit)", "Colab Pro+",
         "Strong Arabic; higher cost;\ndifferent architecture"),
    ]

    y = Inches(1.88)
    for model, size, hw, note in candidates:
        add_rect(sl, Inches(6.6), y, Inches(6.28), Inches(1.25), BLUE)
        add_textbox(sl, Inches(6.75), y+Inches(0.06), Inches(6.0), Inches(0.35),
                    text=model, bold=True, size=14, color=WHITE)
        add_textbox(sl, Inches(6.75), y+Inches(0.43), Inches(6.0), Inches(0.28),
                    text=f"Size: {size}  ·  Platform: {hw}", size=11, color=LBLUE)
        add_textbox(sl, Inches(6.75), y+Inches(0.73), Inches(6.0), Inches(0.38),
                    text=note, size=11, color=LGRAY)
        y += Inches(1.35)

    add_textbox(sl, Inches(6.65), Inches(7.12), Inches(6.2), Inches(0.3),
                text="Start with Qwen3-14B — same tokenizer, minimal code changes",
                bold=True, size=12, color=GOLD)


def slide_next_steps_plan(prs, n, total):
    """Slide 13 – Plan & Timeline"""
    sl = blank_slide(prs)
    header_bar(sl, "Next Steps & Timeline",
               "Two targeted experiments before thesis finalisation")
    slide_number(sl, n, total)

    steps = [
        ("1", "Implement p2v3 (ultra-conservative prompt)",
         "Add prompt version to PromptBuilder.build_zero_shot(version='v3')\n"
         "Update config to allow phase2.prompt_version: p2v3\n"
         "Run export → infer → analyze on all 9 datasets",
         "~1 week", GREEN),
        ("2", "Run p2v2 + p2v3 on Qwen3-14B",
         "Use same 3-stage pipeline (export → infer on Kaggle A100 → analyze)\n"
         "Compare: 4B-p2v2, 4B-p2v3, 14B-p2v2, 14B-p2v3 (2x2 design)\n"
         "Also run Phase 8 (RAG) with 14B to measure retrieval interaction",
         "~2 weeks", BLUE),
        ("3", "Statistical significance testing",
         "Run bootstrap CIs + paired t-test for best results vs Phase 2 baseline\n"
         "Bonferroni correction for multiple comparisons across 9 datasets\n"
         "Also: Wilcoxon signed-rank for non-normal per-sample CER distributions",
         "~3 days", TEAL),
        ("4", "Write up & finalise thesis",
         "Integrate final numbers into Chapter 4 (Results) and Chapter 5 (Discussion)\n"
         "Draft conclusion: RAG + larger model as recommended pipeline\n"
         "Address research question: yes, LLMs can bridge gap — conditions + caveats",
         "~3 weeks", ORANGE),
    ]

    for i, (num, title, body, time, col) in enumerate(steps):
        y = Inches(1.4 + i * 1.5)
        add_rect(sl, Inches(0.35), y, Inches(0.55), Inches(1.3), col)
        add_textbox(sl, Inches(0.35), y+Inches(0.35), Inches(0.55), Inches(0.55),
                    text=num, bold=True, size=28, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(0.95), y, Inches(10.85), Inches(1.3), LGRAY)
        add_textbox(sl, Inches(1.05), y+Inches(0.05), Inches(8.5), Inches(0.38),
                    text=title, bold=True, size=14, color=NAVY)
        add_textbox(sl, Inches(1.05), y+Inches(0.45), Inches(8.5), Inches(0.75),
                    text=body, size=11, color=DGRAY)
        add_rect(sl, Inches(11.85), y, Inches(1.08), Inches(1.3), col)
        add_textbox(sl, Inches(11.85), y+Inches(0.35), Inches(1.08), Inches(0.6),
                    text=time, bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_discussion_questions(prs, n, total):
    """Slide 14 – Discussion / Open Questions for Professor"""
    sl = blank_slide(prs)
    header_bar(sl, "Discussion Points for Professor Meeting")
    slide_number(sl, n, total)

    questions = [
        ("1. Scope of next experiments",
         "Should p2v3 and Qwen3-14B runs be included in the thesis as separate phases, "
         "or positioned as ablations/sensitivity analyses within existing chapters?"),
        ("2. Tahoma anomaly",
         "Tahoma is the only font where NO phase beats the OCR baseline. "
         "Should we investigate why (font-specific confusion patterns, runaway rate 10.4%) "
         "or treat it as a limitation?"),
        ("3. KHATT vs PATS reporting",
         "Given the two domains behave so differently, should we report them separately "
         "throughout the thesis rather than averaging across all datasets?"),
        ("4. Statistical testing",
         "Bootstrap CI + paired t-test on Phase 7 (RAG) vs Phase 2: is this sufficient, "
         "or should we use a non-parametric test (Wilcoxon signed-rank) "
         "given the non-normal distribution of per-sample CER?"),
        ("5. Selective correction strategy",
         "Results show LLM correction is harmful below ~1.5% CER. "
         "Is it worth adding a sample-level confidence gate (detect high-error samples first, "
         "then apply LLM only where beneficial) before finalising the thesis?"),
    ]

    y = Inches(1.35)
    for q_title, q_body in questions:
        add_rect(sl, Inches(0.35), y, Inches(12.63), Inches(1.0), LGRAY)
        add_textbox(sl, Inches(0.5), y+Inches(0.04), Inches(12.4), Inches(0.35),
                    text=q_title, bold=True, size=14, color=NAVY)
        add_textbox(sl, Inches(0.5), y+Inches(0.42), Inches(12.4), Inches(0.5),
                    text=q_body, size=12, color=DGRAY)
        y += Inches(1.1)


def slide_summary(prs, n, total):
    """Slide 15 – Summary"""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, NAVY)
    add_rect(sl, 0, Inches(2.8), W, Inches(0.08), GOLD)

    add_textbox(sl, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9),
                text="Summary", bold=True, size=36, color=WHITE, align=PP_ALIGN.CENTER)

    takeaways = [
        (GREEN,  "LLMs work well on high-error Arabic OCR",
                 "KHATT: 51% CER reduction (5.66% → 2.77%) using RAG + Qwen3-4B"),
        (ORANGE, "Prompt engineering regresses on low-error text",
                 "PATS phases 3/4/6 increase CER vs zero-shot — over-correction at low noise floor"),
        (BLUE,   "RAG (BM25) is the best strategy overall",
                 "Domain-matched training examples outperform all pure-prompt approaches"),
        (TEAL,   "Next: conservative prompt + Qwen3-14B",
                 "p2v3 to reduce false positives; 14B for stronger Arabic orthographic knowledge"),
    ]

    y = Inches(3.3)
    for col, title, body in takeaways:
        add_rect(sl, Inches(0.6), y, Inches(0.5), Inches(0.9), col)
        add_rect(sl, Inches(1.15), y, Inches(11.5), Inches(0.9), RGBColor(0x28,0x45,0x72))
        add_textbox(sl, Inches(1.3), y+Inches(0.04), Inches(11.3), Inches(0.35),
                    text=title, bold=True, size=15, color=WHITE)
        add_textbox(sl, Inches(1.3), y+Inches(0.45), Inches(11.3), Inches(0.38),
                    text=body, size=13, color=LBLUE)
        y += Inches(1.05)

    add_textbox(sl, Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3),
                text="Mohamed Sabry · Master's Thesis · May 2026",
                size=12, color=LBLUE, align=PP_ALIGN.CENTER)

    slide_number(sl, n, total)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    TOTAL = 15

    slide_title(prs, TOTAL)                          # 1
    slide_research_question(prs, 2, TOTAL)           # 2
    slide_datasets(prs, 3, TOTAL)                    # 3
    slide_experimental_design(prs, 4, TOTAL)         # 4
    slide_main_results(prs, 5, TOTAL)                # 5
    slide_finding_rag_wins(prs, 6, TOTAL)            # 6
    slide_finding_prompt_fails(prs, 7, TOTAL)        # 7
    slide_finding_khatt(prs, 8, TOTAL)               # 8
    slide_per_font(prs, 9, TOTAL)                    # 9
    slide_analysis_threshold(prs, 10, TOTAL)         # 10
    slide_next_conservative_prompt(prs, 11, TOTAL)   # 11
    slide_next_bigger_model_fixed(prs, 12, TOTAL)    # 12
    slide_next_steps_plan(prs, 13, TOTAL)            # 13
    slide_discussion_questions(prs, 14, TOTAL)       # 14
    slide_summary(prs, 15, TOTAL)                    # 15

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}  ({TOTAL} slides)")


if __name__ == "__main__":
    main()
