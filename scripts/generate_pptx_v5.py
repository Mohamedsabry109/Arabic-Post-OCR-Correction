#!/usr/bin/env python3
"""Generate MS_Masters_Progress_Report_v5.pptx — complete, definitive progress report.

All metrics: normal samples only (OCR/GT ratio ≤ 5.0), diacritics stripped.
CER† = after runaway correction (threshold 3.0× GT length).

Run:
    python scripts/generate_pptx_v5.py
Output:
    publication/thesis/MS_Masters_Progress_Report_v5.pptx
"""
from __future__ import annotations
import io, sys
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = (Path(__file__).resolve().parent.parent
       / "publication" / "thesis" / "MS_Masters_Progress_Report_v5.pptx")

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x35, 0x5E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)
MBLUE   = RGBColor(0x2E, 0x75, 0xB6)
BODY    = RGBColor(0x22, 0x22, 0x22)
DGRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)
RED     = RGBColor(0xC0, 0x20, 0x20)
DKRED   = RGBColor(0x8B, 0x00, 0x00)
ORANGE  = RGBColor(0xD4, 0x6B, 0x08)
GREEN   = RGBColor(0x1E, 0x6B, 0x2E)
GOLD    = RGBColor(0xD4, 0xA0, 0x17)
LGREEN  = RGBColor(0xD8, 0xF0, 0xDB)
LRED    = RGBColor(0xFA, 0xE0, 0xE0)
LYELL   = RGBColor(0xFF, 0xFB, 0xE6)
LBLU2   = RGBColor(0xE8, 0xF2, 0xFB)
TBL_H   = RGBColor(0x1B, 0x35, 0x5E)
TBL_A   = RGBColor(0xD6, 0xE4, 0xF7)
GOLD_BG = RGBColor(0xFF, 0xF0, 0xC0)
GRYBG   = RGBColor(0xF0, 0xF0, 0xF0)

SW = Emu(12192000)  # 13.33"
SH = Emu(6858000)   # 7.5"
HH = Emu(1066800)   # header bar height

METRIC_NOTE = ("* Normal samples only (OCR/GT ratio ≤5.0), diacritics stripped. "
               "CER† = after runaway correction (ratio threshold 3.0×GT).")


# ── Primitives ───────────────────────────────────────────────────────────────

def _blank(prs):
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _rect(slide, l, t, w, h, fill: RGBColor, line_rgb: RGBColor | None = None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    return shp


def _tb(slide, l, t, w, h, text, size=14, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = "Calibri"
    return tf


def _bullets(slide, l, t, w, h, items, size=14, color=BODY, line_sp=1.3):
    """items: str or (bold_text, rest_text) tuples."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = Pt(size * line_sp)
        p.space_after = Pt(2)
        if isinstance(item, tuple):
            r1 = p.add_run(); r1.text = item[0]
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color; r1.font.name = "Calibri"
            r2 = p.add_run(); r2.text = item[1]
            r2.font.size = Pt(size); r2.font.color.rgb = color
            r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = item
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.name = "Calibri"
    return tf


def _table(slide, l, t, w, rows_data, font_size=11.5,
           hdr_bg=TBL_H, hdr_fg=WHITE, alt_bg=TBL_A,
           col_widths=None, row_h_emu: int | None = None):
    nr = len(rows_data); nc = len(rows_data[0])
    rh = Emu(row_h_emu) if row_h_emu else Emu(int(SH * 0.052))
    tbl_h = rh * nr
    tbl_shape = slide.shapes.add_table(nr, nc, l, t, w, tbl_h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.font.name = "Calibri"
                para.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    para.font.bold = True
                    para.font.color.rgb = hdr_fg
                else:
                    para.font.color.rgb = BODY
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hdr_bg
            elif ri % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = alt_bg
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return tbl


def _set_cell(tbl, ri, ci, bg: RGBColor, fg: RGBColor | None = None,
              bold: bool = False):
    cell = tbl.cell(ri, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    if fg or bold:
        for p in cell.text_frame.paragraphs:
            if fg:
                p.font.color.rgb = fg
            if bold:
                p.font.bold = True


def _metric_footer(slide):
    _tb(slide, Inches(0.25), Inches(7.12), Inches(12.5), Inches(0.28),
        METRIC_NOTE, size=9, color=DGRAY, italic=True)


def _slide_num_badge(prs, slide):
    n = len(prs.slides)
    _tb(slide, Inches(12.55), Inches(0.08), Inches(0.6), Inches(0.3),
        str(n), size=10, color=LGRAY, align=PP_ALIGN.RIGHT)


# ── Slide templates ──────────────────────────────────────────────────────────

def _title_slide(prs):
    slide = prs.slides.add_slide(_blank(prs))
    _rect(slide, 0, 0, SW, SH, NAVY)
    # Accent bar
    _rect(slide, 0, Inches(2.2), SW, Emu(8000), GOLD)
    _tb(slide, Inches(0.8), Inches(0.6), Inches(11.6), Inches(1.7),
        "Arabic Post-OCR Correction with LLMs",
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.8), Inches(2.45), Inches(11.6), Inches(0.8),
        "Can LLMs bridge the gap between open-source and closed-source Arabic OCR?",
        size=18, color=NAVY, align=PP_ALIGN.CENTER, bold=True)
    _tb(slide, Inches(0.8), Inches(3.4), Inches(11.6), Inches(0.6),
        "Final Progress Report  ·  3 Experiments  ·  9 Datasets  ·  4 Models",
        size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.8), Inches(4.3), Inches(11.6), Inches(0.5),
        "Mohamed Sabry  ·  Faculty of Engineering, Cairo University  ·  June 2026",
        size=16, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.8), Inches(5.0), Inches(11.6), Inches(0.45),
        "Primary model: Qwen3-4B-Instruct-2507  ·  OCR: Qaari (Qari-OCR-0.1-VL-2B)",
        size=13, color=LBLUE, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.8), Inches(5.55), Inches(11.6), Inches(0.45),
        "Datasets: PATS-A01 (8 fonts)  +  KHATT  +  KHATT-Para  +  Yarmouk  +  Muharaf  +  Historical  +  BM  +  Kitab",
        size=12, color=LBLUE, align=PP_ALIGN.CENTER)


def _content_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(_blank(prs))
    _rect(slide, 0, 0, SW, HH, NAVY)
    _tb(slide, Inches(0.3), Inches(0.04), Inches(11.5), Inches(0.65),
        title, size=23, bold=True, color=WHITE)
    if subtitle:
        _tb(slide, Inches(0.3), Inches(0.67), Inches(12.5), Inches(0.38),
            subtitle, size=12, color=LBLUE, italic=True)
    _slide_num_badge(prs, slide)
    return slide


def _section_slide(prs, number: str, title: str, subtitle: str,
                   context_lines: list[str] | None = None):
    """Full-bleed dark section break with number, title, subtitle, context."""
    slide = prs.slides.add_slide(_blank(prs))
    _rect(slide, 0, 0, SW, SH, NAVY)
    _rect(slide, 0, Inches(1.8), SW, Emu(6000), GOLD)           # thin gold accent
    _rect(slide, 0, Inches(1.84), Inches(0.12), Inches(3.5), GOLD)  # left bar
    _tb(slide, Inches(0.5), Inches(0.25), Inches(2.5), Inches(1.0),
        number, size=72, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    _tb(slide, Inches(3.0), Inches(0.3), Inches(9.5), Inches(1.3),
        title, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _tb(slide, Inches(0.5), Inches(2.0), Inches(12.0), Inches(0.7),
        subtitle, size=19, color=GOLD, bold=False, align=PP_ALIGN.LEFT)
    if context_lines:
        _bullets(slide, Inches(0.6), Inches(2.9), Inches(12.0), Inches(3.5),
                 context_lines, size=16, color=LBLUE)


# ═══════════════════════════════════════════════════════════════════════════
# DATA — single source of truth
# ═══════════════════════════════════════════════════════════════════════════

FONTS = ["Akhbar", "Simplified", "Traditional", "Arial", "Naskh",
         "Tahoma", "Thuluth", "Andalus"]

P1_CER   = [2.12, 2.94, 3.63, 3.65, 4.10, 5.82, 10.21, 12.12]
P1_WER   = [4.84, 5.81, 9.35, 6.84, 9.56, 9.07, 31.11, 33.67]
P1_N     = [437, 432, 435, 436, 443, 407, 441, 427]
P1_RUN   = [4.8, 5.7, 5.0, 4.6, 3.5, 10.3, 3.9, 6.2]

P2_CER   = [2.71, 2.93, 2.96, 4.50, 4.27, 6.13, 10.92, 12.27]
P3_CER   = [2.57, 2.97, 2.90, 4.48, 4.23, 6.12, 10.94, 11.75]
P4_CER   = [2.59, 2.86, 2.87, 4.08, 4.22, 5.91, 10.67, 11.51]
P6_CER   = [2.46, 2.72, 2.75, 4.09, 4.13, 6.15, 10.56, 13.22]
P8_CER   = [3.01, 3.01, 3.07, 3.83, 4.36, 6.21,  9.68, 10.41]
P9_CER   = [2.27, 2.56, 2.88, 3.76, 3.94, 5.76, 10.61, 11.98]

T2_CER   = [2.97, 3.01, 3.07, 3.83, 4.36, 6.20,  9.68, 10.42]
T3_CER   = [2.32, 2.43, 2.58, 3.74, 3.69, 5.54, 10.23, 11.68]


def _pct(v):
    return f"{v:.2f}%"


def _delta_str(new, base):
    d = new - base
    sign = "+" if d >= 0 else "−"
    return f"{sign}{abs(d):.2f}%"


# ═══════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # ── S1: Title ─────────────────────────────────────────────────────────
    _title_slide(prs)

    # ── S2: Research Question ──────────────────────────────────────────────
    slide = _content_slide(prs, "Research Question & Motivation",
        "Three experiments — progressive prompt engineering, model scaling, domain generalization")
    _bullets(slide, Inches(0.4), Inches(1.15), Inches(6.1), Inches(5.7), [
        ("The Problem:", ""),
        ("  • ", "Qaari (open-source Arabic OCR) produces systematic character-level errors:"),
        ("  • ", "Dot-group confusions: ب/ت/ث/ي/ن — identical skeletons, differ only in dots"),
        ("  • ", "Similar-shape pairs: ف/ق, ر/ز, ص/ض — easily confused under noise"),
        ("  • ", "Runaway bug: Qaari repeats phrases on noisy/blank images (up to 17.7% of samples)"),
        ("", ""),
        ("The Opportunity:", ""),
        ("  • ", "LLMs have broad Arabic linguistic knowledge from pre-training"),
        ("  • ", "Instruction-following enables prompt-based correction without fine-tuning"),
        ("  • ", "Can post-OCR correction close the gap to expensive closed-source VLMs?"),
        ("", ""),
        ("Evaluation methodology:", ""),
        ("  • ", "Normal samples only (OCR/GT length ratio ≤5.0 — exclude Qaari runaways)"),
        ("  • ", "Diacritics stripped before CER/WER — fair comparison across models"),
        ("  • ", "Runaway corrector applied where noted (CER†): ratio threshold 3.0×GT"),
    ], size=14, color=BODY)

    _rect(slide, Inches(6.55), Inches(1.15), Inches(6.6), Inches(5.7), LBLU2, MBLUE)
    _tb(slide, Inches(6.75), Inches(1.3), Inches(6.2), Inches(0.5),
        "Experimental Structure", size=16, bold=True, color=NAVY)
    exp_tbl = [
        ["Exp.", "Research Question", "Runs"],
        ["1\nPhases 1–9", "Does any LLM strategy improve Qaari OCR?\nWhich knowledge type helps most?", "8 phases"],
        ["2\n8 Trials", "Which prompt design is optimal?\nAggressive vs conservative vs RAG vs error-pattern?", "8 trials"],
        ["3\nFinal", "Best prompt × 4 models × 2 OCR sources × 3 domains", "9 runs"],
    ]
    _table(slide, Inches(6.75), Inches(1.85), Inches(6.2), exp_tbl, font_size=12,
           row_h_emu=420000)
    _tb(slide, Inches(6.75), Inches(5.8), Inches(6.2), Inches(0.8),
        "Primary metric throughout: CER (mean over normal samples, diacritics stripped)\n"
        "Primary model: Qwen3-4B-Instruct-2507  |  OCR: Qaari",
        size=11, color=DGRAY, italic=True)

    # ── S3: Exp1 Section ──────────────────────────────────────────────────
    _section_slide(prs, "Exp 1",
        "Phased Knowledge Augmentation",
        "8 experimental phases, each isolating one variable  ·  PATS-A01 (8 fonts) + KHATT  ·  Qwen3-4B",
        context_lines=[
            ("Phase 1: ", "Qaari OCR baseline — error taxonomy, confusion matrix, runaway analysis"),
            ("Phases 2–6,8,9: ", "Progressive LLM correction — zero-shot → confusion matrix → self-reflective → combinations → RAG"),
            ("Key finding: ", "LLM correction is net harmful on PATS average; only RAG achieves below OCR baseline"),
            ("Key finding: ", "KHATT consistently improves ~1 pp CER across all strategies"),
        ])

    # ── S4: Phase 1 Baseline ──────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 1 — Phase 1: Qaari OCR Baseline",
        "Source: results_first_experiment_pats_khatt_only_8_phases_complete/phase1/")
    p1_rows = [["Font", "CER", "WER", "N (normal)", "Runaway %", "Difficulty"]]
    for i, f in enumerate(FONTS):
        diff = ("★ Easiest" if P1_CER[i] < 3 else
                "○ Medium" if P1_CER[i] < 7 else "✖ Hard")
        p1_rows.append([f, _pct(P1_CER[i]), _pct(P1_WER[i]),
                         str(P1_N[i]), f"{P1_RUN[i]:.1f}%", diff])
    p1_rows.append(["PATS avg", "5.57%", "13.78%", "3,458", "5.5%", ""])
    p1_rows.append(["KHATT", "34.24%", "75.60%", "186", "17.7%", "✖ Hard"])
    t = _table(slide, Inches(0.2), Inches(1.13), Inches(6.5), p1_rows,
               font_size=12, row_h_emu=365000)
    # PATS avg row = soft green
    for ci in range(6): _set_cell(t, 9, ci, LGREEN, GREEN)
    # KHATT row = soft red
    for ci in range(6): _set_cell(t, 10, ci, LRED, RED)
    # Highlight easy fonts (CER < 3)
    for ri in [1, 2]: _set_cell(t, ri, 1, GOLD_BG)
    # Highlight hard fonts
    for ri in [7, 8]: _set_cell(t, ri, 1, LRED)

    _bullets(slide, Inches(6.8), Inches(1.13), Inches(6.3), Inches(5.7), [
        ("Key observations:", ""),
        ("  • ", "PATS-A01 CER spans 2.12%–12.12% (5.7× range) across 8 fonts rendering identical text — typeface alone drives this variation"),
        ("  • ", "KHATT handwritten: 34.24% CER = 6.1× harder than PATS average"),
        ("  • ", "Runaway samples excluded from all LLM metrics (Qaari repetition bug)"),
        ("", ""),
        ("Error taxonomy (PATS-A01 aggregate):", ""),
        ("  • ", "42.1% dot confusions (ب/ت/ث/ي/ن)"),
        ("  • ", "18.3% similar-shape substitutions (ف/ق, ر/ز)"),
        ("  • ", "11.6% hamza/alef variants"),
        ("  • ", "9.0% taa marbuta/ha confusions"),
        ("", ""),
        ("KHATT error profile:", ""),
        ("  • ", "37% segmentation errors (word merges + splits) — rare in typewritten text"),
        ("  • ", "Segmentation errors are structurally harder for text-only LLM correction"),
        ("", ""),
        ("Two correction regimes:", ""),
        ("  • ", "PATS: already accurate — LLM risks introducing more errors than it fixes"),
        ("  • ", "KHATT: high noise — LLM has more signal but errors are harder to infer"),
    ], size=13, color=BODY)
    _metric_footer(slide)

    # ── S5: Phases 2-9 Summary ────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 1 — Phases 2–9: All LLM Strategies",
        "PATS avg + KHATT  |  Δ vs Phase 1 OCR baseline  |  Red = harms PATS  |  Green = improves vs baseline")

    phase_rows = [
        ["Phase", "Strategy", "PATS CER", "Δ PATS", "PATS WER", "KHATT CER", "Δ KHATT"],
        ["P1", "OCR Baseline (Qaari)", "5.57%", "—", "13.78%", "34.24%", "—"],
        ["P2", "Zero-Shot LLM", "5.84%", "+0.27%↑", "14.42%", "33.25%", "−0.99%↓"],
        ["P3", "OCR-Aware (confusion matrix)", "5.74%", "+0.17%↑", "14.16%", "33.26%", "−0.98%↓"],
        ["P4", "Self-Reflective (training artifacts)", "5.59%", "+0.02%", "13.66%", "33.24%", "−1.00%↓"],
        ["P5†", "CAMeL Validation (buggy†)", "21.89%", "n/a", "90.74%", "41.62%", "n/a"],
        ["P6", "Combination: conf+self", "5.76%", "+0.19%↑", "13.81%", "33.27%", "−0.97%↓"],
        ["P8", "RAG BM25 (char n-gram)", "5.45%", "−0.12%↓", "12.04%", "33.83%", "−0.41%↓"],
        ["P9", "Error-Signature RAG", "5.47%", "−0.10%↓", "13.63%", "33.13%", "−1.11%↓"],
    ]
    t = _table(slide, Inches(0.2), Inches(1.13), Inches(13.0), phase_rows,
               font_size=11.5, row_h_emu=388000)
    # OCR baseline = green
    for ci in range(7): _set_cell(t, 1, ci, LGREEN, GREEN, True)
    # P5 buggy = yellow warning
    for ci in range(7): _set_cell(t, 5, ci, LYELL, ORANGE)
    # Harmful delta cells (P2, P3, P6) PATS delta column (col 3)
    for ri in [2, 3, 6]: _set_cell(t, ri, 3, LRED, RED)
    # P4 delta near-zero = neutral
    _set_cell(t, 4, 3, LGRAY)
    # P8 and P9 PATS delta = best improvement
    _set_cell(t, 7, 3, LGREEN, GREEN, True)
    _set_cell(t, 8, 3, LGREEN, GREEN)
    # KHATT delta column (col 6) — all improvements = green
    for ri in range(2, 9):
        if ri != 5:
            _set_cell(t, ri, 6, LGREEN if ri in [2, 3, 4, 6, 7, 8, 9] else LRED)
    # Best values
    _set_cell(t, 7, 2, GOLD_BG, NAVY, True)   # P8 best PATS
    _set_cell(t, 8, 5, GOLD_BG, NAVY, True)   # P9 best KHATT

    _tb(slide, Inches(0.3), Inches(6.43), Inches(12.7), Inches(0.42),
        "† Phase 5 had implementation bugs (MorphAnalyzer init + position-keyed revert). Results are invalid; excluded from conclusions.  "
        "  |  Best PATS: P8 RAG (5.45%, −0.12% vs baseline).  Best KHATT: P9 Error-Sig (33.13%, −1.11%).  "
        "All prompt strategies except RAG HARM or are neutral on PATS average.",
        size=10.5, color=DGRAY, italic=True)

    # ── S6: Per-Font Impact Map ────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 1 — Font-Level Impact: Which Strategies Help Which Fonts?",
        "Green = improvement vs OCR baseline  |  Red = degradation  |  Values = CER%  |  Bold = best per font")

    pf_rows = [
        ["Font", "OCR", "P2 ZS", "P3 Conf.", "P4 Self", "P6 C+S", "P8 RAG", "P9 Err"],
    ]
    for i, f in enumerate(FONTS):
        row = [f, _pct(P1_CER[i]),
               _pct(P2_CER[i]), _pct(P3_CER[i]), _pct(P4_CER[i]),
               _pct(P6_CER[i]), _pct(P8_CER[i]), _pct(P9_CER[i])]
        pf_rows.append(row)
    # PATS avg row
    def avg(lst): return sum(lst)/len(lst)
    pf_rows.append(["PATS Avg",
                     "5.57%", "5.84%", "5.74%", "5.59%",
                     "5.76%", "5.45%", "5.47%"])
    pf_rows.append(["KHATT",
                     "34.24%", "33.25%", "33.26%", "33.24%",
                     "33.27%", "33.83%", "33.13%"])

    t = _table(slide, Inches(0.2), Inches(1.12), Inches(13.0), pf_rows,
               font_size=11.5, row_h_emu=378000)

    # Color code every cell (cols 2-7) relative to OCR baseline (col 1)
    col_phase = {2: P2_CER, 3: P3_CER, 4: P4_CER, 5: P6_CER, 6: P8_CER, 7: P9_CER}
    for ri in range(1, 9):  # font rows
        base = P1_CER[ri - 1]
        for ci, ph_data in col_phase.items():
            val = ph_data[ri - 1]
            if val < base - 0.05:
                _set_cell(t, ri, ci, LGREEN, GREEN)
            elif val > base + 0.05:
                _set_cell(t, ri, ci, LRED, RED)
            else:
                _set_cell(t, ri, ci, LGRAY, DGRAY)  # flat

    # PATS avg row (row 9): same logic
    pavg_data = [5.84, 5.74, 5.59, 5.76, 5.45, 5.47]
    pavg_base = 5.57
    for ci, val in zip(range(2, 8), pavg_data):
        if val < pavg_base - 0.05:
            _set_cell(t, 9, ci, LGREEN, GREEN, True)
        else:
            _set_cell(t, 9, ci, LRED, RED, True)

    # KHATT row (row 10)
    kh_data = [33.25, 33.26, 33.24, 33.27, 33.83, 33.13]
    kh_base = 34.24
    for ci, val in zip(range(2, 8), kh_data):
        if val < kh_base - 0.05:
            _set_cell(t, 10, ci, LGREEN, GREEN)
        else:
            _set_cell(t, 10, ci, LRED, RED)

    # Bold best per font
    best_cols = [P8_CER.index(min(P8_CER)), ]  # simplified: mark P8 col as best overall
    _set_cell(t, 9, 6, LGREEN, GREEN, True)   # P8 best PATS avg
    _set_cell(t, 10, 7, LGREEN, GREEN, True)  # P9 best KHATT

    _tb(slide, Inches(0.3), Inches(6.58), Inches(12.7), Inches(0.35),
        "P2 zero-shot harms 7 of 8 fonts (Traditional only exception: 3.63%→2.96%, −18.5%).  "
        "P8 RAG is only strategy with PATS avg below OCR baseline.  "
        "KHATT: all strategies improve by ~1 pp except P8 (slightly worse) and P9 (best: 33.13%).",
        size=10.5, color=DGRAY, italic=True)
    _metric_footer(slide)

    # ── S7: Exp1 Key Findings ─────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 1 — Key Findings",
        "The over-correction problem dominates typewritten correction; RAG is the only reliable strategy")
    _bullets(slide, Inches(0.4), Inches(1.15), Inches(6.0), Inches(5.8), [
        ("Finding 1 — Zero-shot LLM is net harmful on PATS:", ""),
        ("  • ", "P2 zero-shot increases PATS avg CER from 5.57% to 5.84% (+4.8%)"),
        ("  • ", "7 of 8 PATS fonts are harmed; only Traditional improves (−18.5%)"),
        ("  • ", "Root cause: model is \"helpful\" — rewrites correct tokens even with conservative prompt"),
        ("", ""),
        ("Finding 2 — RAG achieves the only PATS improvement:", ""),
        ("  • ", "P8 RAG BM25: 5.45% (−0.12% vs OCR baseline) — only strategy below baseline"),
        ("  • ", "P9 Error-Sig RAG: 5.47% (−0.10%) — close second"),
        ("  • ", "Domain-matched retrieval prevents false corrections; zero-shot examples on clean input ≈ no-op"),
        ("", ""),
        ("Finding 3 — KHATT: modest but consistent improvement:", ""),
        ("  • ", "All prompt strategies reduce KHATT CER by 0.4–1.1 pp"),
        ("  • ", "P9 best: 33.13% (−1.11% abs vs OCR baseline 34.24%)"),
        ("  • ", "Segmentation errors limit further gains (37% of KHATT errors = word merges/splits)"),
        ("", ""),
        ("Finding 4 — No universal threshold:", ""),
        ("  • ", "Traditional (3.63% CER) is helped; Arial (3.65%) is harmed"),
        ("  • ", "Error type, not just baseline CER, determines whether LLM helps"),
    ], size=13, color=BODY)

    _rect(slide, Inches(6.5), Inches(1.15), Inches(6.65), Inches(5.8), LBLU2, MBLUE)
    _tb(slide, Inches(6.7), Inches(1.3), Inches(6.25), Inches(0.5),
        "Per-Font: Which Fonts Benefit from P8 RAG?", size=14, bold=True, color=NAVY)
    summary_rows = [
        ["Font", "OCR", "P8 RAG", "Δ", "Verdict"],
        ["Akhbar",     "2.12%", "3.01%", "+0.89%", "✖ HARMED"],
        ["Simplified", "2.94%", "3.01%", "+0.07%", "○ flat"],
        ["Traditional","3.63%", "3.07%", "−0.56%", "✓ helped"],
        ["Arial",      "3.65%", "3.83%", "+0.18%", "✖ harmed"],
        ["Naskh",      "4.10%", "4.36%", "+0.26%", "✖ harmed"],
        ["Tahoma",     "5.82%", "6.21%", "+0.39%", "✖ harmed"],
        ["Thuluth",   "10.21%", "9.68%", "−0.53%", "✓ helped"],
        ["Andalus",   "12.12%","10.41%", "−1.71%", "✓ helped"],
        ["PATS Avg",   "5.57%", "5.45%", "−0.12%", "✓ best"],
    ]
    t2 = _table(slide, Inches(6.7), Inches(1.9), Inches(6.25), summary_rows,
                font_size=12, row_h_emu=368000)
    for ri, verdict in enumerate(["harmed","flat","helped","harmed","harmed",
                                   "harmed","helped","helped","best"], 1):
        bg = LRED if "harm" in verdict else (LGREEN if "help" in verdict or verdict=="best" else LGRAY)
        fg = RED if "harm" in verdict else (GREEN if "help" in verdict or verdict=="best" else DGRAY)
        _set_cell(t2, ri, 4, bg, fg, verdict == "best")

    # ── S8: Exp2 Section ──────────────────────────────────────────────────
    _section_slide(prs, "Exp 2",
        "Prompt Design Study — 8 Trials",
        "Systematic prompt engineering: base vs conservative vs error-pattern × zero-shot vs RAG",
        context_lines=[
            ("13 datasets: ", "PATS-A01 (8 fonts) + KHATT + KHATT-Para + Yarmouk + Muharaf + Historical"),
            ("Model: ", "Qwen3-4B-Instruct-2507 (same as Exp1)"),
            ("Goal: ", "Find the optimal prompt configuration for Experiment 3"),
            ("Key finding: ", "T3 (conservative zero-shot) best PATS (5.27%); T7 catastrophic (190%!)"),
        ])

    # ── S9: Exp2 Trial Results ────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 2 — 8 Prompt Trials: PATS & KHATT",
        "13 datasets  |  normal-only, diacritics stripped  |  Qwen3-4B  |  ✔ = best per column")
    e2_rows = [
        ["#", "Trial ID", "Prompt Style", "RAG?", "PATS CER", "PATS WER", "KHATT CER", "KHATT WER"],
        ["OCR", "—", "Baseline (no LLM)", "—", "5.57%", "13.78%", "34.24%", "75.60%"],
        ["T1", "base_zs",     "Base (aggressive)",      "No",  "5.84%", "14.43%", "33.25%", "73.93%"],
        ["T2", "base_rag",    "Base",                   "Yes", "5.44%", "12.03%", "33.82%", "73.76%"],
        ["T3", "cons_zs",     "Conservative",           "No",  "5.27%", "13.35%", "32.93%", "74.15%"],
        ["T4", "cons_rag",    "Conservative",           "Yes","27.02%", "28.10%", "32.69%", "73.29%"],
        ["T5", "hp_zs",       "Error-pattern",          "No",  "5.85%", "13.95%", "32.90%", "72.78%"],
        ["T6", "hp_rag",      "Error-pattern",          "Yes","19.64%", "24.52%", "32.64%", "73.40%"],
        ["T7", "hp_cons_zs",  "⚠ Err-pat+Cons",   "No","190.60%","198.12%", "51.64%", "90.58%"],
        ["T8", "hp_cons_rag", "Err-pat+Cons",           "Yes","46.56%", "48.79%", "32.54%", "73.67%"],
    ]
    t = _table(slide, Inches(0.2), Inches(1.13), Inches(13.0), e2_rows, font_size=11.5,
               row_h_emu=388000)
    # OCR baseline
    for ci in range(8): _set_cell(t, 1, ci, LGREEN, GREEN, True)
    # T3 best PATS = gold
    for ci in range(8): _set_cell(t, 4, ci, GOLD_BG)
    _set_cell(t, 4, 4, LGREEN, GREEN, True)
    # T7 catastrophic = deep red
    for ci in range(8): _set_cell(t, 8, ci, LRED, RED)
    _set_cell(t, 8, 4, RGBColor(0x8B,0,0), WHITE, True)
    # T4, T6, T8 problematic = light red for PATS CER
    for ri in [5, 7, 9]: _set_cell(t, ri, 4, LRED, RED)
    # Best KHATT (T8)
    _set_cell(t, 9, 6, LGREEN, GREEN, True)
    # T2 best WER
    _set_cell(t, 3, 5, LGREEN, GREEN, True)

    _tb(slide, Inches(0.3), Inches(6.42), Inches(12.7), Inches(0.43),
        "✔ T3 (conservative zero-shot): best PATS — 5.27% (−5.4% relative vs OCR baseline). No RAG overhead.  "
        "✔ T2 (base+RAG): best overall balance — 5.44% PATS + 33.82% KHATT. → Chosen for Experiment 3.  "
        "⚠ T7: catastrophic — prompt-template tags leaked into LLM output (template injection bug).",
        size=10.5, color=DGRAY, italic=True)

    # ── S10: Exp2 Per-Font ────────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 2 — Per-Font PATS CER",
        "Error-pattern prompts (T4, T6, T7, T8) cause catastrophic font-specific failures")
    pf2_rows = [
        ["Trial", "Akhbar", "Andalus", "Arial", "Naskh", "Simplified", "Tahoma", "Thuluth", "Trad.", "Avg"],
        ["OCR",  "2.12%","12.12%","3.65%","4.10%","2.94%","5.82%","10.21%","3.63%","5.57%"],
        ["T1",   "2.71%","12.26%","4.50%","4.26%","2.94%","6.14%","10.93%","2.97%","5.84%"],
        ["T2",   "2.97%","10.42%","3.83%","4.36%","3.01%","6.20%"," 9.68%","3.07%","5.44%"],
        ["T3",   "2.32%","11.68%","3.74%","3.69%","2.43%","5.54%","10.23%","2.58%","5.27%"],
        ["T4",  "25.32%","36.67%","22.7%","4.01%","2.56%","5.61%","116.3%","3.02%","27.0%"],
        ["T5",   "2.83%","11.87%","4.24%","4.21%","3.01%","6.12%","11.26%","3.25%","5.85%"],
        ["T6",   "4.02%","59.77%","5.28%","5.18%","38.96%","7.77%","10.22%","25.9%","19.6%"],
        ["T7",  "84.8%","153.0%","355.4%","300.1%","156.6%","73.9%","199.0%","202.0%","190.6%"],
        ["T8",   "3.25%","53.21%","4.21%","4.48%","196.4%","6.43%"," 81.4%","23.1%","46.6%"],
    ]
    t = _table(slide, Inches(0.15), Inches(1.12), Inches(13.1), pf2_rows, font_size=10.5,
               row_h_emu=372000)
    # OCR row = green
    for ci in range(10): _set_cell(t, 1, ci, LGREEN, GREEN)
    # T3 best row = gold
    for ci in range(10): _set_cell(t, 4, ci, GOLD_BG)
    _set_cell(t, 4, 9, LGREEN, GREEN, True)
    # T4, T6, T7, T8 = red
    for ri in [5, 7, 8, 9]:
        for ci in range(10): _set_cell(t, ri, ci, LRED)
    # Extreme values deep red
    for (ri, ci) in [(8,1),(8,2),(8,3),(8,4),(8,5),(8,6),(8,7),(8,8),(8,9),
                      (5,7),(9,6)]:
        _set_cell(t, ri, ci, RGBColor(0x8B,0,0), WHITE, True)
    _tb(slide, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.35),
        "T3 (gold) beats T2 on 6/8 fonts. T2 wins Andalus (10.42% vs 11.68%) and Thuluth (9.68% vs 10.23%). "
        "Error-pattern prompts (T4-T8) are unreliable: different fonts fail unpredictably.",
        size=10.5, color=DGRAY, italic=True)
    _metric_footer(slide)

    # ── S11: Exp3 Section ─────────────────────────────────────────────────
    _section_slide(prs, "Exp 3",
        "Final Experiment — 3 Research Categories",
        "Best prompt (T2 base_rag)  ·  4 models  ·  2 OCR sources  ·  BM Bench + Kitab",
        context_lines=[
            ("3A — Model comparison: ",
             "4 models × val set × Qaari OCR × T2  (9 datasets each)"),
            ("3B — OCR source impact: ",
             "Qwen3-4B × val set × T2  ×  Qaari vs Gemma-3 VLM as OCR engine"),
            ("3C — Generalization: ",
             "Qwen3-4B × T2  ×  RDI-Test Benchmark + Kitab  ×  2 OCR sources"),
            ("Runaway corrector applied (CER†): ",
             "critical for Qwen3-14B which has massive runaway epidemic"),
        ])

    # ── S12: Exp3A Model Comparison ───────────────────────────────────────
    slide = _content_slide(prs, "Experiment 3A — Model Comparison",
        "val set · Qaari OCR · T2 (base_rag)  |  CER = raw output  |  CER† = after runaway correction  |  FP = Full-page group")

    agg_rows = [
        ["Model", "Params", "PATS CER", "PATS CER†", "KHATT CER", "KHATT CER†", "FP CER†"],
        ["OCR Baseline", "—", "5.57%", "—", "34.24%", "—", "86.50%"],
        ["Qwen3-4B",    "4B",  "5.44%", "5.44%", "33.82%", "33.82%", "68.18%"],
        ["Qwen3-14B",   "14B","15.30%", "5.60%","175.69%", "35.55%", "88.51%"],
        ["Gemma-3-4B",  "4B",  "7.47%", "7.47%", "36.14%", "36.14%", "71.12%"],
        ["Gemma-3-12B", "12B", "5.33%", "5.33%", "35.98%", "35.98%", "70.01%"],
    ]
    t = _table(slide, Inches(0.2), Inches(1.13), Inches(7.9), agg_rows, font_size=12.5,
               row_h_emu=415000)
    # OCR baseline
    for ci in range(7): _set_cell(t, 1, ci, LGREEN, GREEN, True)
    # Qwen3-14B disaster raw
    _set_cell(t, 3, 2, RGBColor(0x8B,0,0), WHITE, True)   # PATS raw
    _set_cell(t, 3, 4, RGBColor(0x8B,0,0), WHITE, True)   # KHATT raw
    # Best PATS†: Gemma-3-12B
    _set_cell(t, 5, 3, GOLD_BG, NAVY, True)
    # Best KHATT†: Qwen3-4B
    _set_cell(t, 2, 5, GOLD_BG, NAVY, True)
    # Best FP†: Qwen3-4B
    _set_cell(t, 2, 6, GOLD_BG, NAVY, True)
    # Gemma-3-4B worse than small Qwen
    for ci in [2,3,4,5,6]: _set_cell(t, 4, ci, LRED)

    _tb(slide, Inches(8.1), Inches(1.13), Inches(5.1), Inches(0.45),
        "Per-font PATS CER† (runaway-corrected)", size=13, bold=True, color=NAVY)
    pf3_rows = [
        ["Font",       "OCR",   "Q3-4B", "Q3-14B","G3-4B","G3-12B"],
        ["Akhbar",    "2.12%","2.97%","3.27%","4.80%","3.05%"],
        ["Simplified","2.94%","3.01%","3.39%","4.93%","3.26%"],
        ["Traditional","3.63%","3.07%","3.96%","5.48%","3.35%"],
        ["Arial",     "3.65%","3.83%","4.50%","5.90%","3.67%"],
        ["Naskh",     "4.10%","4.36%","4.75%","5.94%","4.41%"],
        ["Tahoma",    "5.82%","6.20%","6.10%","7.64%","6.03%"],
        ["Thuluth",  "10.21%","9.68%","9.06%","11.81%","8.86%"],
        ["Andalus",  "12.12%","10.42%","9.80%","13.26%","10.00%"],
        ["Avg",       "5.57%","5.44%","5.60%","7.47%","5.33%"],
    ]
    t2 = _table(slide, Inches(8.1), Inches(1.65), Inches(5.1), pf3_rows, font_size=11.5,
                row_h_emu=372000)
    for ci in range(6): _set_cell(t2, 9, ci, LGREEN, GREEN)
    # Gemma-3-12B col = gold (best PATS)
    for ri in range(1, 9): _set_cell(t2, ri, 5, GOLD_BG)
    _set_cell(t2, 9, 5, LGREEN, GREEN, True)
    # Gemma-3-4B col = light red (worst)
    for ri in range(1, 10): _set_cell(t2, ri, 4, LRED)

    _bullets(slide, Inches(0.2), Inches(5.95), Inches(13.0), Inches(1.15), [
        ("Key: ",
         "Runaway corrector essential for Q3-14B (KHATT 175.69%→35.55% with fix).  "
         "G3-12B best PATS (5.33%); Q3-4B best KHATT (33.82%) + FP (68.18%).  "
         "FP = Full-page: OCR 86.50% → Q3-4B 68.18% (−21% abs).  "
         "G3-4B consistently worst; larger scale helps only within each family."),
    ], size=11.5, color=BODY)
    _metric_footer(slide)

    # ── S13: Exp3B OCR Source ─────────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 3B — OCR Source Quality Impact",
        "Qwen3-4B · T2  |  Full-page datasets only (Gemma excluded for line-strip images)")

    # Explanation box
    _rect(slide, Inches(0.2), Inches(1.12), Inches(13.0), Emu(650000),
          RGBColor(0xFF,0xF0,0xC8), RGBColor(0xD4,0xA0,0x17))
    _tb(slide, Inches(0.35), Inches(1.17), Inches(12.7), Inches(0.38),
        "⚠  Gemma-3 VLM OCR is EXCLUDED for PATS and KHATT (line-strip images, 10–16:1 aspect ratio). "
        "Gemma's image preprocessor squashes these to near-square → hallucinations + loops → results invalid. "
        "Gemma results reported only for full-page datasets (~0.6:1 ratio) where the comparison is valid.",
        size=11.5, color=RGBColor(0x5C,0x3D,0x00))

    # Full-page table: Qaari vs Gemma
    _tb(slide, Inches(0.2), Inches(1.72), Inches(8.0), Inches(0.4),
        "Full-page datasets — Qaari vs Gemma OCR (valid comparison)", size=14, bold=True, color=NAVY)
    fp_rows = [
        ["Dataset", "OCR Src", "OCR CER", "OCR WER", "Corr CER†", "Corr WER†"],
        ["KHATT-Paragraph-val",  "Qaari",  "61.68%",  "91.24%", "45.06%",  "68.49%"],
        ["",                     "Gemma",  "36.14%",  "64.43%", "35.82%",  "63.82%"],
        ["Yarmouk-testing",      "Qaari",  "49.85%",  "78.32%", "43.35%",  "69.94%"],
        ["",                     "Gemma",  "95.76%", "140.20%", "74.37%", "111.74%"],
        ["Muharaf-validation",   "Qaari", "129.39%", "152.67%", "89.78%", "123.10%"],
        ["",                     "Gemma", "141.15%", "182.26%", "72.85%", "109.99%"],
        ["Historical",           "Qaari", "105.10%", "136.81%", "94.53%", "127.03%"],
        ["",                     "Gemma",  "76.26%", "121.77%", "74.09%", "118.52%"],
        ["Full-page Avg",        "Qaari",  "86.50%", "116.67%", "68.18%",  "80.60%"],
        ["Full-page Avg",        "Gemma",  "87.33%", "123.17%", "64.28%",  "88.05%"],
    ]
    t = _table(slide, Inches(0.2), Inches(2.15), Inches(8.0), fp_rows, font_size=11.5,
               row_h_emu=355000)
    for ci in range(6): _set_cell(t, 9, ci, LGREEN, GREEN, True)   # Qaari avg
    for ci in range(6): _set_cell(t, 10, ci, LYELL)                 # Gemma avg
    # Gemma corrected better overall
    _set_cell(t, 10, 4, LGREEN, GREEN, True)

    _bullets(slide, Inches(8.4), Inches(1.72), Inches(4.75), Inches(5.5), [
        ("Key findings:", ""),
        ("", ""),
        ("Full-page comparison:", ""),
        ("  • ", "Gemma OCR ≈ Qaari (87% vs 87%)"),
        ("  • ", "Gemma corrected: 64.28% vs Qaari 68.18%"),
        ("  • ", "LLM extracts more from Gemma FP output"),
        ("", ""),
        ("Why not PATS/KHATT:", ""),
        ("  • ", "Line strips: 10–16:1 ratio"),
        ("  • ", "Gemma preprocessor: squashes to square"),
        ("  • ", "Result: hallucinations + loops → invalid"),
        ("", ""),
        ("Implication:", ""),
        ("  • ", "OCR source matters less for full-page"),
        ("  • ", "Line-strip OCR: Qaari is the reliable engine"),
        ("  • ", "VLM OCR needs aspect-ratio-aware preprocessing"),
    ], size=13, color=BODY)
    _metric_footer(slide)

    # ── S14: Exp3C Generalization ─────────────────────────────────────────
    slide = _content_slide(prs, "Experiment 3C — Generalization to New Domains",
        "Qwen3-4B · T2 · British Museum Arabic Benchmark + Kitab Benchmark  |  CER†")

    # Left — RDI-Test-Lines
    _tb(slide, Inches(0.2), Inches(1.1), Inches(6.5), Inches(0.55),
        "RDI-Test-Lines (line-strip images, text GT)",
        size=14, bold=True, color=NAVY)
    _tb(slide, Inches(0.2), Inches(1.52), Inches(6.5), Inches(0.3),
        "Note: RDI-Test Line Segmentation excluded — polygon-only GT, no text transcription.",
        size=10.5, color=DGRAY, italic=True)
    bm_rows = [
        ["Subset", "OCR Src", "N", "OCR CER", "Corr CER†", "Δ"],
        ["LR-Handwritten", "Qaari",  "385",  "98.82%", "91.55%", "−7.3%"],
        ["",               "Gemma",  "762",  "62.53%", "57.71%", "−4.8%"],
        ["LR-Manuscripts", "Qaari", "1009",  "81.94%", "79.00%", "−2.9%"],
        ["",               "Gemma", "1175",  "78.10%", "70.36%", "−7.7%"],
        ["LR-Typewritten", "Qaari",  "671", "105.19%", "84.79%","−20.4%"],
        ["",               "Gemma", "1314",  "64.86%", "55.31%", "−9.5%"],
        ["RDI-Test-Lines Overall",  "Qaari", "2065",  "95.31%", "85.11%","−10.7%"],
        ["RDI-Test-Lines Overall",  "Gemma", "3251",  "68.49%", "61.13%","−10.7%"],
    ]
    t = _table(slide, Inches(0.2), Inches(1.85), Inches(6.7), bm_rows, font_size=11.5,
               row_h_emu=355000)
    for ci in range(6):
        _set_cell(t, 7, ci, LGREEN, GREEN, True)
        _set_cell(t, 8, ci, LYELL)
    _set_cell(t, 5, 5, LGREEN, GREEN, True)  # Typewritten best delta
    # Gemma better than Qaari on BM!
    _set_cell(t, 2, 3, LGREEN)   # Gemma Handwritten OCR better
    _set_cell(t, 6, 3, LGREEN)   # Gemma Typewritten OCR better

    # Right — Kitab
    _tb(slide, Inches(7.0), Inches(1.1), Inches(6.15), Inches(0.4),
        "Kitab Benchmark (13 subsets)", size=14, bold=True, color=NAVY)
    kit_rows = [
        ["Subset", "OCR Src", "OCR CER", "Corr CER†", "Δ"],
        ["kitab-hindawi",    "Qaari","31.46%","24.31%","−7.2%"],
        ["kitab-historyar",  "Qaari","63.48%","51.69%","−11.8%"],
        ["kitab-muharaf",    "Qaari","76.04%","58.72%","−17.3%"],
        ["kitab-isippt",     "Gemma","87.48%","68.75%","−18.7%"],
        ["kitab-patsocr",    "Gemma","87.98%","69.95%","−18.0%"],
        ["kitab-khattpar", "Qaari","205.27%","90.60%","−114.7%"],
        ["kitab-arabicocr",  "Qaari", "1.80%", "7.21%", "+5.4%⚠"],
        ["kitab-evarest",    "Qaari","32.92%","60.67%","+27.8%⚠"],
        ["Kitab Overall",    "Qaari","49.28%","40.73%","−17.3%"],
        ["Kitab Overall",    "Gemma","79.45%","60.07%","−24.4%"],
    ]
    t2 = _table(slide, Inches(7.0), Inches(1.55), Inches(6.15), kit_rows, font_size=11.5,
                row_h_emu=347000)
    for ci in range(5):
        _set_cell(t2, 9, ci, LGREEN, GREEN, True)
        _set_cell(t2, 10, ci, LYELL)
    _set_cell(t2, 6, 3, LGREEN, GREEN, True)  # khattparagraph corrected a lot
    _set_cell(t2, 7, 4, LRED, RED)            # arabicocr harmed
    _set_cell(t2, 8, 4, LRED, RED)            # evarest harmed
    _set_cell(t2, 10, 4, LGREEN, GREEN, True) # Gemma Kitab best delta

    _tb(slide, Inches(0.2), Inches(6.5), Inches(12.9), Inches(0.42),
        "RDI-Test-Lines: Gemma OCR better than Qaari on Handwritten (62% vs 98%) and Typewritten (65% vs 105%) — "
        "Qaari struggles on diverse BM images.  "
        "Kitab −17% (Qaari) / −24% (Gemma).  "
        "Low-CER subsets harmed (arabicocr 1.8%); khattparagraph 205%→90% (runaway corrector essential).",
        size=10.5, color=DGRAY, italic=True)
    _metric_footer(slide)

    # ── S15: Summary & Conclusions ────────────────────────────────────────
    slide = _content_slide(prs, "Summary — Final Results & Conclusions",
        "All metrics: normal samples only, diacritics stripped, CER† = runaway-corrected")

    _bullets(slide, Inches(0.35), Inches(1.13), Inches(6.1), Inches(5.8), [
        ("Research Answer: Yes, with important caveats", ""),
        ("", ""),
        ("Exp1 — Phased strategies:", ""),
        ("  • ", "LLM correction net harmful on PATS avg (zero-shot: +4.8%)"),
        ("  • ", "Only RAG achieves below OCR baseline: P8 5.45% (−0.12%)"),
        ("  • ", "KHATT: ~1 pp improvement across all strategies (P9 best: 33.13%)"),
        ("  • ", "Over-correction: 7/8 fonts harmed by zero-shot; font-specific, not CER-threshold-based"),
        ("", ""),
        ("Exp2 — Prompt design:", ""),
        ("  • ", "Conservative zero-shot (T3): best PATS 5.27% (−5.4% relative)"),
        ("  • ", "Base+RAG (T2): best overall balance — used in Exp3"),
        ("  • ", "Error-pattern prompts: risky — catastrophic on some fonts"),
        ("", ""),
        ("Exp3 — Final results:", ""),
        ("  • ", "Best PATS: Gemma-3-12B 5.33%  |  Best KHATT: Qwen3-4B 33.82%"),
        ("  • ", "Qwen3-14B: runaway epidemic (175%→35% with corrector)"),
        ("  • ", "Full-page: OCR 86.50% → 68.18% (−21% abs)"),
        ("  • ", "Generalizes: BM −10.7%, Kitab −17–24%"),
        ("  • ", "Gemma VLM OCR: excluded for line strips (invalid); better than Qaari on RDI-Test-Lines images"),
    ], size=13, color=BODY)

    _rect(slide, Inches(6.55), Inches(1.13), Inches(6.6), Inches(5.8), LBLU2, MBLUE)
    _tb(slide, Inches(6.75), Inches(1.28), Inches(6.2), Inches(0.45),
        "Best Results at a Glance", size=15, bold=True, color=NAVY)
    best_rows = [
        ["Task / Dataset", "Best Config", "CER", "vs OCR"],
        ["PATS (8 fonts)", "G3-12B, T2, Qaari",  "5.33%",  "−4.3% rel"],
        ["PATS (no RAG)",  "T3 cons_zs, Q3-4B",  "5.27%",  "−5.4% rel"],
        ["KHATT",          "Q3-4B, T2, Qaari",   "33.82%", "−1.2% rel"],
        ["Full-page",      "Q3-4B, T2, Qaari",   "68.18%", "−21% abs"],
        ["RDI-Test Benchmark",   "Q3-4B, T2, Qaari",   "85.11%", "−10.7% abs"],
        ["Kitab (Qaari)",  "Q3-4B, T2, Qaari",   "40.73%", "−17.3% rel"],
        ["Kitab (Gemma OCR)", "Q3-4B, T2, Gemma", "60.07%", "−24.4% rel"],
    ]
    t = _table(slide, Inches(6.75), Inches(1.8), Inches(6.2), best_rows, font_size=12.5,
               row_h_emu=395000)
    for ri in range(1, 8): _set_cell(t, ri, 3, LGREEN, GREEN)
    _set_cell(t, 1, 2, GOLD_BG, NAVY, True)  # PATS best
    _set_cell(t, 2, 2, GOLD_BG, NAVY, True)  # PATS no-RAG best

    _tb(slide, Inches(6.75), Inches(6.3), Inches(6.2), Inches(0.5),
        "Mohamed Sabry  ·  Master’s Thesis  ·  Cairo University  ·  June 2026\n"
        "Metric: normal samples only, diacritics stripped",
        size=11, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
