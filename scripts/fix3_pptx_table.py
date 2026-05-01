"""
Add Phase 7 and Phase 8 rows to the Experimental Phases Overview table on Slide 18.
Also rename future-work RAG bullet from title-prefix style to numbered style.
"""
import sys
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = "publication/thesis/defense_presentation.pptx"

def safe_print(s):
    sys.stdout.buffer.write((s + "\n").encode("utf-8", errors="replace"))
    sys.stdout.buffer.flush()

def clone_row(tbl_elem, source_row_idx):
    """Deep-clone a <a:tr> element from the table XML."""
    rows = tbl_elem.findall(qn("a:tr"))
    return copy.deepcopy(rows[source_row_idx])

def set_cell_text(row_elem, col_idx, text):
    """Set the text of the col_idx-th cell in a row element."""
    cells = row_elem.findall(qn("a:tc"))
    if col_idx >= len(cells):
        return
    tc = cells[col_idx]
    # Find or create a:txBody > a:p > a:r > a:t
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return
    # Clear existing paragraphs
    for p in txBody.findall(qn("a:p")):
        for r in p.findall(qn("a:r")):
            t = r.find(qn("a:t"))
            if t is not None:
                t.text = ""
        # Just set the first run of the first paragraph
        break
    # Set text in first paragraph's first run
    paras = txBody.findall(qn("a:p"))
    if paras:
        p = paras[0]
        runs = p.findall(qn("a:r"))
        if runs:
            t = runs[0].find(qn("a:t"))
            if t is not None:
                t.text = text
        else:
            # Create a run
            r_new = etree.SubElement(p, qn("a:r"))
            t_new = etree.SubElement(r_new, qn("a:t"))
            t_new.text = text

def add_table_row(tbl_shape, row_data):
    """
    Clone the last data row of the table and set cell texts.
    row_data: list of strings, one per column.
    """
    tbl_elem = tbl_shape._element.find(qn("a:tbl"))
    if tbl_elem is None:
        safe_print("No a:tbl found in shape")
        return
    rows = tbl_elem.findall(qn("a:tr"))
    last_row = rows[-1]
    new_row = copy.deepcopy(last_row)
    for col_idx, text in enumerate(row_data):
        set_cell_text(new_row, col_idx, text)
    # Insert after last row
    tbl_elem.append(new_row)
    safe_print(f"  Added row: {row_data}")

def main():
    prs = Presentation(PPTX_PATH)
    slide18 = prs.slides[17]

    # Find the table shape
    tbl_shape = None
    for shape in slide18.shapes:
        if shape.shape_type == 19:  # TABLE
            tbl_shape = shape
            break

    if tbl_shape is None:
        safe_print("ERROR: No table found on slide 18")
        return

    safe_print("Slide 18 table BEFORE:")
    for row in tbl_shape.table.rows:
        row_texts = [cell.text_frame.text.strip() for cell in row.cells]
        safe_print(f"  {row_texts}")

    # Check if Phase 7 and Phase 8 rows already exist
    all_texts = " ".join(
        cell.text_frame.text for row in tbl_shape.table.rows for cell in row.cells
    )

    if "7" not in [row.cells[0].text_frame.text.strip() for row in tbl_shape.table.rows]:
        add_table_row(tbl_shape, [
            "7",
            "DSPy Prompt Optimization",
            "+Automated prompt search",
            "BootstrapFewShot compiled program",
            "vs Phase 2",
        ])

    if "8" not in [row.cells[0].text_frame.text.strip() for row in tbl_shape.table.rows]:
        add_table_row(tbl_shape, [
            "8",
            "RAG (BM25)",
            "+Retrieved examples",
            "Top-5 similar (OCR,GT) pairs from training",
            "vs Phase 2",
        ])

    safe_print("\nSlide 18 table AFTER:")
    for row in tbl_shape.table.rows:
        row_texts = [cell.text_frame.text.strip() for cell in row.cells]
        safe_print(f"  {row_texts}")

    # ----------------------------------------------------------------
    # Rename future-work RAG bullet (Slide 38) to have a proper number
    # Current: "Improve RAG: replace BM25 with dense Arabic embeddings..."
    # Target:  "8. Improve RAG retrieval: replace BM25 with dense Arabic embeddings; explore cross-dataset retrieval"
    # ----------------------------------------------------------------
    slide38 = prs.slides[37]
    for shape in slide38.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.startswith("Improve RAG"):
                    run.text = run.text.replace(
                        "Improve RAG: replace BM25 with dense Arabic embeddings; explore cross-dataset retrieval",
                        "8. Improve RAG retrieval: replace BM25 with dense Arabic embeddings; explore cross-dataset retrieval",
                    )
                    safe_print(f"Slide 38: renamed RAG future work bullet to #8")

    prs.save(PPTX_PATH)
    safe_print(f"\nSaved: {PPTX_PATH}")

if __name__ == "__main__":
    main()
