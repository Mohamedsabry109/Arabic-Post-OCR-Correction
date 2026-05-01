"""Fix and verify Phase 8 PPTX changes. Must run after update_pptx_phase8.py."""
import sys
import copy
from pptx import Presentation
from pptx.oxml.ns import qn

PPTX_PATH = "publication/thesis/defense_presentation.pptx"

def safe_print(s):
    sys.stdout.buffer.write((s + "\n").encode("utf-8", errors="replace"))
    sys.stdout.buffer.flush()

def print_shape(slide_idx, shape):
    safe_print(f"  Slide {slide_idx+1} shape '{shape.name}':")
    for i, para in enumerate(shape.text_frame.paragraphs):
        t = para.text
        if t.strip():
            safe_print(f"    para[{i}]: {t[:130]!r}")

def remove_paragraphs_containing(text_frame, substring):
    """Remove all paragraphs whose text contains substring."""
    txBody = text_frame._txBody
    to_remove = []
    for p in txBody.findall(qn("a:p")):
        texts = "".join(r.text or "" for r in p.findall(qn("a:r")))
        if substring in texts:
            to_remove.append(p)
    for p in to_remove:
        txBody.remove(p)
    return len(to_remove)

def add_bullet_paragraph(text_frame, text, copy_from_para=None):
    """Append a new paragraph copying format from copy_from_para."""
    if copy_from_para is not None:
        new_para_xml = copy.deepcopy(copy_from_para._p)
        for r in new_para_xml.findall(qn("a:r")):
            new_para_xml.remove(r)
        for br in new_para_xml.findall(qn("a:br")):
            new_para_xml.remove(br)
        ref_runs = copy_from_para._p.findall(qn("a:r"))
        if ref_runs:
            new_r = copy.deepcopy(ref_runs[0])
            t_elem = new_r.find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = text
            new_para_xml.append(new_r)
        text_frame._txBody.append(new_para_xml)
    else:
        para = text_frame.add_paragraph()
        para.text = text

def main():
    prs = Presentation(PPTX_PATH)
    safe_print(f"Loaded: {len(prs.slides)} slides")

    # ----------------------------------------------------------------
    # Slide 7 (idx=6): Fix the duplicate RAG bullet in title TextBox
    # The RAG bullet was incorrectly added to TextBox 5 ("Our Contributions" title).
    # Remove it from there; it belongs only in TextBox 6 (the list).
    # ----------------------------------------------------------------
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text
        # The title box "Our Contributions" should NOT have the RAG bullet
        if "Our Contributions" in full and "BM25 RAG" in full:
            n = remove_paragraphs_containing(shape.text_frame, "BM25 RAG")
            safe_print(f"Slide 7 '{shape.name}': removed {n} erroneous RAG para(s) from title box")

    # ----------------------------------------------------------------
    # Slide 18 (idx=17): Check if Phase 8 was missed; add if absent
    # ----------------------------------------------------------------
    slide18 = prs.slides[17]
    safe_print("\n--- Slide 18 shapes ---")
    for shape in slide18.shapes:
        if shape.has_text_frame:
            print_shape(17, shape)
    # Find the phases list shape (Shape 3 = TextBox 4, shape index 3, should have Phase 1..7)
    for shape in slide18.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text
        if "Phase" in full and ("Phase 1" in full or "Baseline" in full or "Zero-Shot" in full):
            if "Phase 8" not in full and "RAG" not in full:
                paras = shape.text_frame.paragraphs
                ref_para = None
                for p in paras:
                    if "Phase 7" in p.text or "DSPy" in p.text or "Phase 6" in p.text:
                        ref_para = p
                if ref_para is None:
                    for p in paras:
                        if p.text.strip():
                            ref_para = p
                add_bullet_paragraph(
                    shape.text_frame,
                    "Phase 8: RAG -- BM25 retrieval of similar training corrections as few-shot context",
                    copy_from_para=ref_para,
                )
                safe_print(f"Slide 18 '{shape.name}': added Phase 8 RAG")
            elif "Phase 8" in full:
                safe_print(f"Slide 18 '{shape.name}': Phase 8 already present")

    # ----------------------------------------------------------------
    # Slide 19 (idx=18): Check Phase 8 RAG was added
    # ----------------------------------------------------------------
    slide19 = prs.slides[18]
    safe_print("\n--- Slide 19 shapes ---")
    for shape in slide19.shapes:
        if shape.has_text_frame:
            print_shape(18, shape)

    # ----------------------------------------------------------------
    # Check for duplicate RAG entries on slides 38, 40, 41, 42
    # ----------------------------------------------------------------
    TARGET_PHRASES = {
        37: ("dense retrieval", "Slide 38 Future Work"),
        39: ("Phase 8 (RAG, BM25)", "Slide 40 Ranking"),
        40: ("Domain-matched RAG", "Slide 41 Conclusion"),
        41: ("BM25 RAG using domain-matched", "Slide 42 Contributions"),
    }
    for idx, (phrase, label) in TARGET_PHRASES.items():
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape.text_frame.text
            if phrase in full:
                # Count occurrences
                count = full.count(phrase)
                safe_print(f"\n{label} shape '{shape.name}': found phrase {count}x")
                print_shape(idx, shape)
                if count > 1:
                    # Remove all duplicate paragraphs (keep first occurrence only)
                    txBody = shape.text_frame._txBody
                    paras = txBody.findall(qn("a:p"))
                    seen = False
                    for p in paras:
                        texts = "".join(r.text or "" for r in p.findall(qn("a:r")))
                        if phrase in texts:
                            if seen:
                                txBody.remove(p)
                                safe_print(f"  -> removed duplicate paragraph")
                            else:
                                seen = True

    # ----------------------------------------------------------------
    # Final diagnostics
    # ----------------------------------------------------------------
    safe_print("\n=== FINAL STATE ===")
    for idx in [6, 11, 12, 17, 18, 37, 39, 40, 41]:
        slide = prs.slides[idx]
        safe_print(f"\n--- Slide {idx+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print_shape(idx, shape)

    prs.save(PPTX_PATH)
    safe_print(f"\nSaved: {PPTX_PATH}")

if __name__ == "__main__":
    main()
