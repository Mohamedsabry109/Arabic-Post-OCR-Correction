"""Verify and clean up the Phase 8 PPTX changes."""
import sys
from pptx import Presentation

PPTX_PATH = "publication/thesis/defense_presentation.pptx"

def print_shape(slide_idx, shape, label=""):
    print(f"\n  [{label}] Slide {slide_idx+1} shape '{shape.name}':")
    for i, para in enumerate(shape.text_frame.paragraphs):
        t = para.text
        if t.strip():
            print(f"    para[{i}]: {t[:120]!r}")

prs = Presentation(PPTX_PATH)

# Verify key slides
check = {
    6: ["7-phase", "8-phase", "RAG", "BM25"],
    11: ["Research Questions", "Experimental Phases"],
    12: ["RQ8"],
    17: ["Phase 8"],
    18: ["Phase 8", "RAG"],
    37: ["dense retrieval", "RAG"],
    39: ["Phase 8", "RAG"],
    40: ["Phase 8", "RAG"],
    41: ["RAG", "BM25"],
}

for idx, keywords in check.items():
    slide = prs.slides[idx]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text
        if any(kw in full for kw in keywords):
            print_shape(idx, shape, "MATCH")
