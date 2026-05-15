#!/usr/bin/env python3
"""Add detailed analysis slides to the MS Progress Report PPTX.

New slides cover:
  A.  Metric interpretation: CER/WER and the Runaway Effect (expanded with formulas)
  A2. Mathematical derivation + PATS-Akhbar worked example
  B0. Context for slide 3 results (all-samples vs normal-only)
  B.  Detailed Phase 1 Baseline (all 4 metric variants)
  C.  All-Phase KHATT breakdown (all 4 metric variants, P1 n/a fixed)
  D.  Prompt architecture: Phase 2 (Zero-Shot)
  E.  Prompt augmentation: Phase 3 (OCR-Aware) and Phase 4 (Self-Reflective)
  F.  RAG Phase 8: BM25 character n-gram retrieval
  G.  RAG Phase 9: Error-Signature retrieval
  H.  Real correction examples (KHATT and PATS)
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Constants ──────────────────────────────────────────────────────────────

PPTX_IN  = Path(r"D:\Masters\Arabic-Post-OCR-Correction\publication\thesis\MS_Masters_Progress_Report.pptx")
PPTX_OUT = Path(r"D:\Masters\Arabic-Post-OCR-Correction\publication\thesis\MS_Masters_Progress_Report_v3.pptx")

SW = Emu(12188952)   # slide width
SH = Emu(6858000)    # slide height
HH = Emu(1143000)    # header rect height
M  = Emu(457200)     # 0.5-in margin

# Colours
C_NAVY   = RGBColor(0x1F, 0x38, 0x64)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LBLUE  = RGBColor(0xBD, 0xD7, 0xEE)
C_DBLUE  = RGBColor(0x2E, 0x75, 0xB6)
C_BODY   = RGBColor(0x40, 0x40, 0x40)
C_RED    = RGBColor(0xC0, 0x00, 0x00)
C_GREEN  = RGBColor(0x37, 0x86, 0x35)
C_ORANGE = RGBColor(0xED, 0x7D, 0x31)
C_TABLE_HDR = RGBColor(0x1F, 0x38, 0x64)
C_TABLE_ALT = RGBColor(0xD6, 0xE4, 0xF7)
C_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
C_LGREEN = RGBColor(0xE2, 0xEF, 0xDA)
C_LRED   = RGBColor(0xFF, 0xE7, 0xE7)


# ── Low-level helpers ──────────────────────────────────────────────────────

def _blank_layout(prs: Presentation):
    best = min(prs.slide_layouts, key=lambda l: len(l.placeholders))
    return best


def _add_rect(slide, left, top, width, height, fill_rgb: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_rtl(para):
    pPr = para._p.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(para._p, qn('a:pPr'))
        para._p.insert(0, pPr)
    pPr.set('rtl', '1')


def _run(para, text: str, size_pt: int, bold: bool = False,
         color: RGBColor = None, italic: bool = False, font_name: str = "Calibri"):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    run.font.name = font_name
    return run


def _header(slide, title: str, subtitle: str = "", slide_num: int = 0, total: int = 24):
    _add_rect(slide, 0, 0, SW, HH, C_NAVY)
    if slide_num:
        num_box = _add_textbox(slide, SW - Emu(900000), Emu(30000), Emu(860000), Emu(300000))
        tf = num_box.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        _run(para, f"{slide_num} / {total}", 9, color=C_LBLUE)
    t_box = _add_textbox(slide, M, Emu(80000), SW - 2 * M, Emu(600000))
    tf = t_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    _run(para, title, 24, bold=True, color=C_WHITE)
    if subtitle:
        s_box = _add_textbox(slide, M, Emu(750000), SW - 2 * M, Emu(320000))
        tf2 = s_box.text_frame
        tf2.word_wrap = True
        para2 = tf2.paragraphs[0]
        _run(para2, subtitle, 13, color=C_LBLUE)


def _body_box(slide, left, top, width, height):
    box = _add_textbox(slide, left, top, width, height)
    box.text_frame.word_wrap = True
    return box.text_frame


def _add_line(tf, text: str, size: int = 12, bold: bool = False,
              color: RGBColor = None, indent: str = "",
              align: PP_ALIGN = PP_ALIGN.LEFT, rtl: bool = False):
    para = tf.add_paragraph()
    para.alignment = align
    if rtl:
        _set_rtl(para)
    if indent:
        _run(para, indent, size, color=C_BODY)
    _run(para, text, size, bold=bold, color=color or C_BODY)
    return para


def _table(slide, left, top, width, rows, cols):
    tbl = slide.shapes.add_table(rows, cols, left, top, width, Emu(0)).table
    return tbl


def _cell(tbl, row, col, text: str, bold: bool = False, size: int = 10,
          color: RGBColor = None, bg: RGBColor = None, align: PP_ALIGN = PP_ALIGN.CENTER,
          rtl: bool = False):
    cell = tbl.cell(row, col)
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    if rtl:
        _set_rtl(para)
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color if color else C_BODY
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    return cell


def _new_slide(prs: Presentation) -> object:
    layout = _blank_layout(prs)
    return prs.slides.add_slide(layout)


# ── Slide builders ─────────────────────────────────────────────────────────

TOTAL = 24  # 14 original + 10 new

def slide_A_runaway(prs, num):
    """Why CER approx WER? Runaway Effect + Formulas"""
    sl = _new_slide(prs)
    _header(sl,
            "Why CER ≈ WER? — The Runaway Effect Explained",
            "Normal OCR errors: WER > CER  |  Qaari repetition bug: CER ≈ WER  |  Formulas below",
            num, TOTAL)

    panel_top = HH + Emu(160000)
    panel_h   = Emu(2750000)
    W3 = (SW - 3 * M) // 3

    # Panel 1 -- Normal samples
    left1 = M
    _add_rect(sl, left1, panel_top, W3, panel_h, C_LGREEN)
    tf1 = _body_box(sl, left1 + Emu(60000), panel_top + Emu(60000),
                    W3 - Emu(120000), panel_h - Emu(120000))
    _add_line(tf1, "NORMAL Samples", 13, bold=True, color=C_GREEN)
    _add_line(tf1, "(no runaway)", 9, color=C_BODY)
    _add_line(tf1, "", 5)
    _add_line(tf1, "1 char error ruins the whole word.", 10)
    _add_line(tf1, "WER is amplified at word level.", 10)
    _add_line(tf1, "", 5)
    _add_line(tf1, "Result: WER >> CER", 12, bold=True, color=C_GREEN)
    _add_line(tf1, "", 5)
    _add_line(tf1, "PATS-Akhbar normal (no diac):", 10, bold=True, color=C_DBLUE)
    _add_line(tf1, "  CER = 2.12%", 10, color=C_GREEN)
    _add_line(tf1, "  WER = 4.84%", 10, color=C_GREEN)
    _add_line(tf1, "  WER/CER = 2.28x", 10, color=C_GREEN)
    _add_line(tf1, "", 4)
    _add_line(tf1, "KHATT normal (no diac):", 10, bold=True, color=C_DBLUE)
    _add_line(tf1, "  CER = 34.24%", 10, color=C_GREEN)
    _add_line(tf1, "  WER = 75.60%", 10, color=C_GREEN)
    _add_line(tf1, "  WER/CER = 2.21x", 10, color=C_GREEN)

    # Panel 2 -- Runaway samples
    left2 = M + W3 + M
    _add_rect(sl, left2, panel_top, W3, panel_h, C_LRED)
    tf2 = _body_box(sl, left2 + Emu(60000), panel_top + Emu(60000),
                    W3 - Emu(120000), panel_h - Emu(120000))
    _add_line(tf2, "RUNAWAY Samples", 13, bold=True, color=C_RED)
    _add_line(tf2, "(Qaari repetition bug)", 9, color=C_BODY)
    _add_line(tf2, "", 5)
    _add_line(tf2, "OCR repeats entire phrases:", 10)
    _add_line(tf2, '"word word word..." x24 times', 10)
    _add_line(tf2, "Both chars AND words inflate", 10)
    _add_line(tf2, "by exactly the same factor k.", 10)
    _add_line(tf2, "", 5)
    _add_line(tf2, "Result: CER = WER = k-1", 12, bold=True, color=C_RED)
    _add_line(tf2, "(ratio = 1.0x exactly)", 10)
    _add_line(tf2, "", 5)
    _add_line(tf2, "Example (k=24 repeats):", 10, bold=True, color=C_RED)
    _add_line(tf2, "  CER = 24-1 = 23 = 2300%", 10, color=C_RED)
    _add_line(tf2, "  WER = 24-1 = 23 = 2300%", 10, color=C_RED)
    _add_line(tf2, "  Ratio = 1.0x", 10, color=C_RED)

    # Panel 3 -- Mean contamination
    left3 = left2 + W3 + M
    _add_rect(sl, left3, panel_top, W3, panel_h, RGBColor(0xFF, 0xF3, 0xE0))
    tf3 = _body_box(sl, left3 + Emu(60000), panel_top + Emu(60000),
                    W3 - Emu(120000), panel_h - Emu(120000))
    _add_line(tf3, "MEAN (All Samples)", 13, bold=True, color=C_ORANGE)
    _add_line(tf3, "(normal + runaway mixed)", 9, color=C_BODY)
    _add_line(tf3, "", 5)
    _add_line(tf3, "A few extreme runaway values", 10)
    _add_line(tf3, "(CER ~ 2300%) swamp the mean.", 10)
    _add_line(tf3, "=> all-samples CER ~ WER", 10)
    _add_line(tf3, "", 5)
    _add_line(tf3, "Result: CER ~ WER (mean)", 12, bold=True, color=C_ORANGE)
    _add_line(tf3, "", 5)
    _add_line(tf3, "PATS-Akhbar ALL:", 10, bold=True, color=C_ORANGE)
    _add_line(tf3, "  22/459 runaway (4.8%)", 10)
    _add_line(tf3, "  CER = 114.42%", 10, color=C_ORANGE)
    _add_line(tf3, "  WER = 110.82%", 10, color=C_ORANGE)
    _add_line(tf3, "  Ratio = 0.97x", 10, color=C_ORANGE)
    _add_line(tf3, "", 3)
    _add_line(tf3, "KHATT ALL:", 10, bold=True, color=C_ORANGE)
    _add_line(tf3, "  40/226 runaway (17.7%)", 10)
    _add_line(tf3, "  CER = 565.58%", 10, color=C_ORANGE)
    _add_line(tf3, "  WER = 725.03%", 10, color=C_ORANGE)

    # ── Math formulas section ──────────────────────────────────────────────
    math_top = panel_top + panel_h + Emu(100000)
    math_h   = SH - math_top - Emu(150000)
    half_w   = (SW - 3 * M) // 2

    # Left: Normal formula (green tint)
    _add_rect(sl, M, math_top, half_w, math_h, C_LGREEN)
    mf1 = _body_box(sl, M + Emu(80000), math_top + Emu(60000),
                    half_w - Emu(160000), math_h - Emu(120000))
    _add_line(mf1, "Normal OCR Error — Key Formula", 11, bold=True, color=C_GREEN)
    _add_line(mf1, "", 3)
    _add_line(mf1, "CER = edit_dist(ref_chars, ocr_chars) / |ref_chars|", 9, color=C_DBLUE)
    _add_line(mf1, "WER = edit_dist(ref_words, ocr_words) / |ref_words|", 9, color=C_DBLUE)
    _add_line(mf1, "", 3)
    _add_line(mf1, "1 char error => 1 bad word:", 10)
    _add_line(mf1, "  CER weight: 1 / |ref_chars|", 9)
    _add_line(mf1, "  WER weight: 1 / |ref_words|", 9)
    _add_line(mf1, "  WER/CER ratio ~ |ref_chars|/|ref_words|", 9)
    _add_line(mf1, "             = avg chars per word", 9)
    _add_line(mf1, "  Arabic avg ~ 5-8 chars => WER > CER", 9, color=C_GREEN)

    # Right: Runaway formula (red tint)
    right_x = M + half_w + M
    _add_rect(sl, right_x, math_top, half_w, math_h, C_LRED)
    mf2 = _body_box(sl, right_x + Emu(80000), math_top + Emu(60000),
                    half_w - Emu(160000), math_h - Emu(120000))
    _add_line(mf2, "Runaway (k repetitions) — Key Formula", 11, bold=True, color=C_RED)
    _add_line(mf2, "", 3)
    _add_line(mf2, "OCR = GT repeated k times:", 10)
    _add_line(mf2, "  edit_dist = (k-1) x |GT_chars|", 9, color=C_DBLUE)
    _add_line(mf2, "  CER = (k-1)|GT_c| / |GT_c| = k-1", 9, color=C_DBLUE)
    _add_line(mf2, "  WER = (k-1)|GT_w| / |GT_w| = k-1", 9, color=C_DBLUE)
    _add_line(mf2, "", 3)
    _add_line(mf2, "=> CER = WER = k-1  (ratio = 1.0x)", 10, bold=True, color=C_RED)
    _add_line(mf2, "", 3)
    _add_line(mf2, "Mean contamination:", 10)
    _add_line(mf2, "  CER_all = (N_norm x CER_n + N_run x CER_r) / N", 9, color=C_DBLUE)
    _add_line(mf2, "  Few large CER_r values dominate => CER_all ~ WER_all", 9, color=C_RED)


def slide_A2_math_proof(prs, num):
    """Mathematical Proof + PATS-Akhbar Worked Example"""
    sl = _new_slide(prs)
    _header(sl,
            "Mathematical Derivation — Runaway Contamination Proof",
            "Formal derivation of CER=WER=k-1 + PATS-Akhbar back-calculation case study",
            num, TOTAL)

    half = (SW - 3 * M) // 2

    # Left: Formal derivation
    tf1 = _body_box(sl, M, HH + Emu(180000), half, SH - HH - Emu(220000))
    _add_line(tf1, "Formal Derivation", 13, bold=True, color=C_DBLUE)
    _add_line(tf1, "", 4)

    _add_line(tf1, "Setup:", 11, bold=True, color=C_NAVY)
    _add_line(tf1, "Let GT have N_c chars and N_w words.", 10)
    _add_line(tf1, "Qaari outputs GT repeated k times exactly.", 10)
    _add_line(tf1, "", 4)

    _add_line(tf1, "CER calculation:", 11, bold=True, color=C_NAVY)
    _add_line(tf1, "  OCR_chars = k x N_c", 10, color=C_DBLUE)
    _add_line(tf1, "  To transform OCR -> GT:", 10)
    _add_line(tf1, "    delete (k-1) x N_c extra chars", 10)
    _add_line(tf1, "  edit_dist = (k-1) x N_c", 10, color=C_DBLUE)
    _add_line(tf1, "  CER = (k-1) x N_c / N_c = k-1", 10, color=C_RED)
    _add_line(tf1, "", 4)

    _add_line(tf1, "WER calculation:", 11, bold=True, color=C_NAVY)
    _add_line(tf1, "  OCR_words = k x N_w", 10, color=C_DBLUE)
    _add_line(tf1, "  edit_dist = (k-1) x N_w", 10, color=C_DBLUE)
    _add_line(tf1, "  WER = (k-1) x N_w / N_w = k-1", 10, color=C_RED)
    _add_line(tf1, "", 4)

    _add_line(tf1, "Therefore: CER = WER = k-1  (ratio = 1.0x)", 11, bold=True, color=C_RED)
    _add_line(tf1, "", 6)

    _add_line(tf1, "Mean contamination formula:", 11, bold=True, color=C_NAVY)
    _add_line(tf1, "CER_all x N_c_all = CER_n x N_c_n + CER_r x N_c_r", 10, color=C_DBLUE)
    _add_line(tf1, "=> CER_r = (CER_all x N_c_all - CER_n x N_c_n) / N_c_r", 10, color=C_DBLUE)
    _add_line(tf1, "(weighted by reference char count, not sample count)", 9, color=C_BODY)
    _add_line(tf1, "", 6)

    _add_line(tf1, "Key Result:", 11, bold=True, color=C_GREEN)
    _add_line(tf1, "All-samples mean CER ~ WER because runaway samples", 10)
    _add_line(tf1, "drive both metrics equally. Only normal-only separates them.", 10)

    # Right: PATS-Akhbar case study
    right_x = M + half + M
    tf2 = _body_box(sl, right_x, HH + Emu(180000), half, SH - HH - Emu(220000))
    _add_line(tf2, "Case Study: PATS-Akhbar Validation", 13, bold=True, color=C_DBLUE)
    _add_line(tf2, "", 4)

    _add_line(tf2, "Observed data (Phase 1 baseline):", 11, bold=True, color=C_NAVY)
    _add_line(tf2, "  Total samples: 459", 10)
    _add_line(tf2, "  Normal: 437  |  Runaway: 22 (4.8%)", 10)
    _add_line(tf2, "  Ref chars (all): 37,101", 10)
    _add_line(tf2, "  Ref chars (normal): 35,285", 10)
    _add_line(tf2, "  Ref chars (runaway): 37,101-35,285 = 1,816", 10, color=C_ORANGE)
    _add_line(tf2, "", 4)

    _add_line(tf2, "Measured metrics (with diacritics):", 11, bold=True, color=C_NAVY)
    _add_line(tf2, "  CER_all = 114.42%  (= 1.1442 ratio)", 10)
    _add_line(tf2, "  CER_normal = 2.22%  (= 0.0222 ratio)", 10)
    _add_line(tf2, "  WER_all = 110.82%  (= 1.1082 ratio)", 10)
    _add_line(tf2, "  WER_normal = 5.12%  (= 0.0512 ratio)", 10)
    _add_line(tf2, "", 4)

    _add_line(tf2, "Back-calculate CER for runaway samples:", 11, bold=True, color=C_NAVY)
    _add_line(tf2, "  CER_r x 1,816 = 1.1442x37,101 - 0.0222x35,285", 10, color=C_DBLUE)
    _add_line(tf2, "                = 42,449 - 783 = 41,666", 10, color=C_DBLUE)
    _add_line(tf2, "  CER_r = 41,666 / 1,816 = 22.95 = 2295%", 10, bold=True, color=C_RED)
    _add_line(tf2, "", 4)

    _add_line(tf2, "Back-calculate WER for runaway samples:", 11, bold=True, color=C_NAVY)
    _add_line(tf2, "  Ref words (all)=7,631  (normal)=7,255  (runaway)=376", 10, color=C_DBLUE)
    _add_line(tf2, "  WER_r x 376 = 1.1082x7,631 - 0.0512x7,255", 10, color=C_DBLUE)
    _add_line(tf2, "             = 8,457 - 371 = 8,086", 10, color=C_DBLUE)
    _add_line(tf2, "  WER_r = 8,086 / 376 = 21.50 = 2150%", 10, bold=True, color=C_RED)
    _add_line(tf2, "", 4)

    _add_line(tf2, "Interpretation:", 11, bold=True, color=C_GREEN)
    _add_line(tf2, "  CER_r / WER_r = 2295 / 2150 = 1.07x  (theory: 1.0x)", 10, color=C_GREEN)
    _add_line(tf2, "  k-1 = 22.95  =>  avg repetition count k ~ 24x", 10, color=C_GREEN)
    _add_line(tf2, "  22 runaway samples contribute ~98% of all-samples CER!", 10, color=C_RED)
    _add_line(tf2, "  => Normal-only is essential for true quality measurement", 10)


def slide_B0_slide3_justification(prs, num):
    """Context for Phase Results: All-Samples vs Normal-Only"""
    sl = _new_slide(prs)
    _header(sl,
            "Understanding the Results — Metric Context Matters",
            "All-samples CER/WER is inflated by Qaari runaway bug  |  Normal-only = true OCR quality",
            num, TOTAL)

    # Left: comparison table (Phase 1 baseline)
    tbl_w = Emu(6800000)
    tbl_top = HH + Emu(220000)

    # Table: 4 rows x 6 cols
    tbl = _table(sl, M, tbl_top, tbl_w, 4, 6)
    col_widths = [Emu(1700000), Emu(1100000), Emu(1100000), Emu(1100000), Emu(1100000), Emu(700000)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for i in range(4):
        tbl.rows[i].height = Emu(560000)

    # Header row
    hdrs = [
        ("Dataset (Phase 1)", C_WHITE, C_TABLE_HDR),
        ("All+Diac CER", C_WHITE, C_TABLE_HDR),
        ("All NoDiac CER", C_WHITE, C_TABLE_HDR),
        ("Norm+Diac CER", C_WHITE, C_TABLE_HDR),
        ("Norm NoDiac CER *", C_WHITE, C_DBLUE),
        ("Runaway %", C_WHITE, C_TABLE_HDR),
    ]
    for ci, (h, tc, bg) in enumerate(hdrs):
        _cell(tbl, 0, ci, h, bold=True, size=9, color=tc, bg=bg)

    # Data rows
    data = [
        ("PATS avg (8 fonts)", "131.55%", "127.56%", "6.50%", "5.57%", "5.5%"),
        ("KHATT-validation",   "565.58%", "404.64%", "60.59%", "34.24%", "17.7%"),
        ("KHATT / PATS ratio", "4.30x",   "3.17x",   "9.32x",  "6.15x",  "--"),
    ]
    row_bgs = [None, C_TABLE_ALT, C_LGRAY]
    ratio_colors = [C_BODY, C_DBLUE, C_RED]

    for ri, (row_data, bg, rc) in enumerate(zip(data, row_bgs, ratio_colors)):
        name_bold = (ri == 2)
        name_col = C_RED if ri == 2 else (C_DBLUE if ri == 1 else C_BODY)
        _cell(tbl, ri+1, 0, row_data[0], bold=name_bold, size=9,
              color=name_col, bg=bg, align=PP_ALIGN.LEFT)
        for ci in range(1, 6):
            val = row_data[ci]
            vc = rc if ri == 2 else C_BODY
            if ri < 2 and ci == 4:  # star column -- primary metric highlight
                vc = C_DBLUE
            _cell(tbl, ri+1, ci, val, size=9, color=vc, bg=bg)

    # Right: Claim validations
    rx = M + tbl_w + M
    rw = SW - rx - M
    tf2 = _body_box(sl, rx, HH + Emu(180000), rw, SH - HH - Emu(200000))

    _add_line(tf2, "Key Claims Validated", 12, bold=True, color=C_NAVY)
    _add_line(tf2, "", 4)

    _add_line(tf2, "Claim 1: KHATT 4.3x harder than PATS", 11, bold=True, color=C_DBLUE)
    _add_line(tf2, "Context: all-samples + diacritics CER", 10)
    _add_line(tf2, "565.58 / 131.55 = 4.30x  [verified]", 10, color=C_GREEN)
    _add_line(tf2, "Normal-only no-diac: 34.24/5.57 = 6.15x", 10, color=C_ORANGE)
    _add_line(tf2, "(true difficulty gap is larger!)", 9, color=C_BODY)
    _add_line(tf2, "", 6)

    _add_line(tf2, "Claim 2: 51% improvement (P8 KHATT)", 11, bold=True, color=C_DBLUE)
    _add_line(tf2, "Context: all-samples CER", 10)
    _add_line(tf2, "P1=565.58%, P8=277.27%", 10)
    _add_line(tf2, "(565.58-277.27)/565.58 = 51.0%  [verified]", 10, color=C_GREEN)
    _add_line(tf2, "But normal-only no-diac:", 10)
    _add_line(tf2, "P1=34.24%, P8=33.83% => 1.2% gain", 10, color=C_RED)
    _add_line(tf2, "(51% = runaway resolution, not OCR fix)", 9, color=C_BODY)
    _add_line(tf2, "", 6)

    _add_line(tf2, "Why normal-only no-diac is primary:", 11, bold=True, color=C_DBLUE)
    _add_line(tf2, "1. Removes runaway contamination", 10)
    _add_line(tf2, "   (17.7% KHATT / 5.5% PATS runaway)", 10)
    _add_line(tf2, "2. Removes diacritic artifacts", 10)
    _add_line(tf2, "   (Qaari inconsistently adds/drops diacritics)", 10)
    _add_line(tf2, "3. Shows true character correction quality", 10)
    _add_line(tf2, "", 6)

    _add_line(tf2, "* PATS avg = simple mean of 8 fonts (CER %)", 8, color=C_BODY)

    # Bottom note spanning full width
    note_tf = _body_box(sl, M, SH - Emu(520000), SW - 2*M, Emu(480000))
    _add_line(note_tf,
              "Data source: Phase 1 baseline_metrics.json (OCR = Qaari v1, no LLM). "
              "PATS avg is arithmetic mean across 8 fonts. "
              "All-samples includes 22-47 runaway samples per PATS font (avg 5.5%) "
              "and 40 runaway samples for KHATT (17.7%).",
              9, color=C_BODY)


def slide_B_baseline(prs, num):
    """Detailed Phase 1 Baseline: 4 Metric Variants per Dataset"""
    sl = _new_slide(prs)
    _header(sl,
            "Phase 1 Baseline — Complete Metric Breakdown (4 Variants)",
            "Each dataset x {All / Normal-only} x {With diacritics / Without diacritics}",
            num, TOTAL)

    rows = [
        # (Font, all_cer, all_wer, all_nd_cer, all_nd_wer, norm_cer, norm_wer, norm_nd_cer, norm_nd_wer, runaway_pct)
        ("PATS-Akhbar",    "114.42%","110.82%","111.91%","110.53%","2.22%","5.12%","2.12%","4.84%","4.8%"),
        ("PATS-Andalus",   "154.43%","166.20%","143.75%","162.61%","14.61%","37.21%","12.12%","33.67%","6.2%"),
        ("PATS-Arial",     "106.43%","107.61%","106.22%","107.29%","3.78%","7.18%","3.65%","6.84%","4.6%"),
        ("PATS-Naskh",     "85.96%","92.35%","81.97%","88.04%","5.68%","13.77%","4.10%","9.56%","3.5%"),
        ("PATS-Simplified","131.77%","133.07%","131.69%","132.83%","3.02%","6.05%","2.94%","5.81%","5.7%"),
        ("PATS-Tahoma",    "239.25%","235.11%","238.79%","234.99%","5.87%","9.18%","5.82%","9.07%","10.4%"),
        ("PATS-Thuluth",   "102.38%","107.18%","90.24%","104.09%","12.93%","33.96%","10.21%","31.11%","3.9%"),
        ("PATS-Traditional","117.72%","122.20%","115.91%","121.39%","3.87%","10.07%","3.63%","9.35%","5.0%"),
        ("KHATT",          "565.58%","725.03%","404.64%","672.77%","60.59%","89.25%","34.24%","75.60%","17.7%"),
    ]

    n_rows = 10
    n_cols = 10
    tbl = _table(sl, M, HH + Emu(200000), SW - 2*M, n_rows, n_cols)

    col_widths = [Emu(1500000), Emu(900000), Emu(900000), Emu(900000), Emu(900000),
                  Emu(900000), Emu(900000), Emu(900000), Emu(900000), Emu(700000)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for i in range(n_rows):
        tbl.rows[i].height = Emu(490000) if i == 0 else Emu(420000)

    headers = [
        ("Dataset",       C_WHITE, C_TABLE_HDR),
        ("All+Diac CER",  C_WHITE, C_TABLE_HDR),
        ("All+Diac WER",  C_WHITE, C_TABLE_HDR),
        ("All NoDiac CER",C_WHITE, C_TABLE_HDR),
        ("All NoDiac WER",C_WHITE, C_TABLE_HDR),
        ("Norm CER",      C_WHITE, C_TABLE_HDR),
        ("Norm WER",      C_WHITE, C_TABLE_HDR),
        ("Norm NoDiac CER *",C_WHITE, C_DBLUE),
        ("Norm NoDiac WER *",C_WHITE, C_DBLUE),
        ("Runaway",       C_WHITE, C_TABLE_HDR),
    ]
    for ci, (h, tc, bg) in enumerate(headers):
        _cell(tbl, 0, ci, h, bold=True, size=9, color=tc, bg=bg, align=PP_ALIGN.CENTER)

    for ri, row_data in enumerate(rows):
        bg = C_TABLE_ALT if ri % 2 == 1 else None
        font_col = C_DBLUE if "KHATT" in row_data[0] else C_BODY
        _cell(tbl, ri+1, 0, row_data[0], bold=True, size=10, color=font_col, bg=bg, align=PP_ALIGN.LEFT)
        for ci in range(1, 5):
            _cell(tbl, ri+1, ci, row_data[ci], size=9, bg=bg)
        for ci in range(5, 9):
            _cell(tbl, ri+1, ci, row_data[ci], size=9, bg=bg)
        rcolor = C_RED if float(row_data[9].replace('%','')) > 8 else C_BODY
        _cell(tbl, ri+1, 9, row_data[9], size=9, color=rcolor, bg=bg)

    note_tf = _body_box(sl, M, SH - Emu(350000), SW - 2*M, Emu(300000))
    _add_line(note_tf,
              "* Primary metric (starred): Normal-Only + No Diacritics. "
              "All-samples CER ~ WER because runaway inflates both equally. "
              "Normal-only shows true OCR quality: WER > CER as expected.",
              10, color=C_BODY)


def slide_C_phases_khatt(prs, num):
    """All-Phase KHATT Comparison: 4 Metric Variants (P1 all-samples corrected)"""
    sl = _new_slide(prs)
    _header(sl,
            "All-Phase Comparison — KHATT Validation (226 samples, 40 runaway)",
            "Primary metric: Normal-Only + No Diacritics (*)  |  Lower is better",
            num, TOTAL)

    # P1 all-samples data is now correct (no n/a):
    # ocr_all: CER=565.58%, WER=725.03%
    # ocr_all_no_diacritics: CER=404.64%, WER=672.77%
    phases = [
        ("P1 — OCR Baseline",
         "60.59%","89.25%","34.24%","75.60%",
         "565.58%","725.03%","404.64%","672.77%"),
        ("P2 — Zero-Shot (v2)",
         "50.93%","82.64%","33.25%","73.94%",
         "320.73%","453.14%","281.75%","445.32%"),
        ("P3 — OCR-Aware (conf.mat)",
         "50.72%","83.04%","33.26%","74.12%",
         "332.98%","358.75%","276.86%","350.62%"),
        ("P4 — Self-Reflective",
         "47.68%","79.77%","33.24%","72.67%",
         "301.94%","343.97%","258.63%","337.23%"),
        ("P8 — RAG BM25",
         "45.91%","78.86%","33.83%","73.76%",
         "277.27%","327.75%","255.56%","323.06%"),
        ("P9 — Error-Sig RAG",
         "50.47%","81.65%","33.13%","73.06%",
         "317.57%","364.94%","280.24%","357.20%"),
    ]

    n_rows = len(phases) + 2
    n_cols = 9
    tbl = _table(sl, M, HH + Emu(200000), SW - 2*M, n_rows, n_cols)

    col_widths = [Emu(2300000)] + [Emu(1236119)] * 8
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for i in range(n_rows):
        tbl.rows[i].height = Emu(490000) if i <= 1 else Emu(520000)

    # Header row 1
    _cell(tbl, 0, 0, "Phase", bold=True, size=10, color=C_WHITE, bg=C_TABLE_HDR)
    _cell(tbl, 0, 1, "Normal-Only (with diacritics)", bold=True, size=9,
          color=C_WHITE, bg=RGBColor(0x37,0x86,0x35))
    _cell(tbl, 0, 2, "", bg=RGBColor(0x37,0x86,0x35))
    _cell(tbl, 0, 3, "Normal-Only NO diacritics (*)", bold=True, size=9,
          color=C_WHITE, bg=C_DBLUE)
    _cell(tbl, 0, 4, "", bg=C_DBLUE)
    _cell(tbl, 0, 5, "All Samples (with diacritics)", bold=True, size=9,
          color=C_WHITE, bg=RGBColor(0xC0,0x00,0x00))
    _cell(tbl, 0, 6, "", bg=RGBColor(0xC0,0x00,0x00))
    _cell(tbl, 0, 7, "All Samples (NO diacritics)", bold=True, size=9,
          color=C_WHITE, bg=RGBColor(0x7B,0x00,0x00))
    _cell(tbl, 0, 8, "", bg=RGBColor(0x7B,0x00,0x00))

    # Header row 2
    sub_headers = ["Phase / Variant", "CER", "WER", "CER", "WER", "CER", "WER", "CER", "WER"]
    for ci, h in enumerate(sub_headers):
        _cell(tbl, 1, ci, h, bold=True, size=9, color=C_WHITE, bg=C_TABLE_HDR)

    # Data rows
    for ri, row_data in enumerate(phases):
        bg = C_TABLE_ALT if ri % 2 == 1 else None
        is_p1  = ri == 0
        is_p8  = "P8" in row_data[0]
        is_p9  = "P9" in row_data[0]
        name_color = (C_BODY if is_p1 else
                      (RGBColor(0x0D,0x52,0x8A) if is_p8 else
                       (C_RED if is_p9 else C_BODY)))
        _cell(tbl, ri+2, 0, row_data[0], bold=(is_p8 or is_p1), size=10,
              color=name_color, bg=bg, align=PP_ALIGN.LEFT)
        for ci in range(1, 9):
            val = row_data[ci]
            val_color = C_BODY
            # Highlight best in each column group (lower = better), skip P1
            if ri > 0:
                try:
                    v = float(val.replace('%',''))
                    # Best normal-only no-diac CER is P8=33.83% (not lowest but P9=33.13% is)
                    if ci == 3 and v <= 33.25:
                        val_color = C_GREEN
                    if ci == 4 and v <= 72.67:
                        val_color = C_GREEN
                    # Highlight P8 all-samples as best
                    if ci in [5,6,7,8] and is_p8:
                        val_color = C_GREEN
                except:
                    pass
            _cell(tbl, ri+2, ci, val, size=9, color=val_color, bg=bg)

    note_tf = _body_box(sl, M, SH - Emu(430000), SW - 2*M, Emu(400000))
    _add_line(note_tf,
              "* Primary metric for thesis: Normal-Only + No Diacritics. "
              "P8 (RAG BM25) achieves best all-samples CER (277.27% vs 565.58% P1 = 51% improvement). "
              "On normal-only no-diac the improvement is modest: P1=34.24% -> P4=33.24% (best, -1.0 pp). "
              "P1 all-samples data corrected: 565.58%/725.03% (with diac), 404.64%/672.77% (no diac).",
              10, color=C_BODY)


def slide_D_prompt_p2(prs, num):
    """Prompt Structure: Phase 2 Zero-Shot (Crafted Prompt)"""
    sl = _new_slide(prs)
    _header(sl,
            "Prompt Structure — Phase 2: Zero-Shot (Crafted Prompt)",
            "Base prompt shared by all phases  |  Injected XML sections vary per phase",
            num, TOTAL)

    lw = SW // 2 - M - Emu(200000)
    tf1 = _body_box(sl, M, HH + Emu(200000), lw, SH - HH - M - Emu(200000))
    _add_line(tf1, "System Prompt Structure (XML tags)", 13, bold=True, color=C_DBLUE)
    _add_line(tf1, "", 4)
    sections = [
        ("<system>",        "OCR post-correction engine role + output rule"),
        ("<task>",          "Sequential word-by-word alignment instruction"),
        ("<rules priority>","8 absolute rules (detail right)"),
        ("  [Phase 3]",     "=> <confusion_patterns> injected here"),
        ("  [Phase 4]",     "=> <self_analysis> injected here"),
        ("  [Phase 8/9]",   "=> <retrieved_corrections> injected here"),
        ("<examples>",      "Static + dynamic few-shot pairs"),
        ("  [Phase 8/9]",   "=> retrieved pairs appended here"),
        ("<output_format>", "Output is ONLY corrected text, nothing else"),
    ]
    for tag, desc in sections:
        is_inject = "injected" in desc or "appended" in desc
        tag_color = C_ORANGE if is_inject else C_DBLUE
        _add_line(tf1, f"{tag}", 10, bold=True, color=tag_color)
        _add_line(tf1, f"  {desc}", 9, color=C_BODY)
        _add_line(tf1, "", 2)

    rw = SW // 2 - M - Emu(200000)
    rl = SW // 2 + Emu(100000)
    tf2 = _body_box(sl, rl, HH + Emu(200000), rw, SH - HH - M - Emu(200000))
    _add_line(tf2, "Critical Rules (Phase 2 base)", 13, bold=True, color=C_DBLUE)
    _add_line(tf2, "", 4)
    rules = [
        "1. STRICT VISUAL CORRECTION ONLY -- fix misread chars, not grammar",
        "2. COLLAPSE HALLUCINATION LOOPS -- truncate >3 repeats to once",
        "3. NO TASHKEEL -- strip all diacritics from output",
        "4. PRESERVE CLASSICAL SPELLING -- no modernization",
        "5. DELETE LATIN ARTIFACTS -- OCR errors only",
        "6. UNRECOGNIZABLE TEXT -- leave unchanged if unsure",
        "7. PROPER NOUNS -- correct only unambiguous visual errors",
        "8. PRESERVE ALF LAM -- do not add or remove definite article",
    ]
    for r in rules:
        _add_line(tf2, r, 9, color=C_BODY)
        _add_line(tf2, "", 2)
    _add_line(tf2, "", 6)
    _add_line(tf2, "Real Example (PATS-Thuluth):", 11, bold=True, color=C_DBLUE)
    _add_line(tf2, "", 3)
    _add_line(tf2, "OCR input:", 9, bold=True)
    _add_line(tf2, "منه فلما خرج سكعج سكعئين في قبل الحكعت وقال هذة القبلة", 10, rtl=True, color=C_ORANGE)
    _add_line(tf2, "LLM output:", 9, bold=True)
    _add_line(tf2, "منه فلما خرج سكعج سكعئين في قبل الحكمة وقال هذه القبلة", 10, rtl=True, color=C_DBLUE)
    _add_line(tf2, "Ground truth:", 9, bold=True)
    _add_line(tf2, "منه فلما خرج ركع ركعتين في قبل الكعبة وقال هذه القبلة", 10, rtl=True, color=C_GREEN)
    _add_line(tf2, "", 4)
    _add_line(tf2, "LLM partially corrected (fixed typo) but missed", 9, color=C_BODY)
    _add_line(tf2, "systematic Thuluth font distortions throughout.", 9, color=C_BODY)


def slide_E_prompt_p3_p4(prs, num):
    """Prompt Augmentation: Phase 3 (OCR-Aware) and Phase 4 (Self-Reflective)"""
    sl = _new_slide(prs)
    _header(sl,
            "Prompt Augmentation — Phase 3 (OCR-Aware) vs Phase 4 (Self-Reflective)",
            "Both phases inject additional XML sections into the Phase 2 base prompt",
            num, TOTAL)

    half = (SW - 3*M) // 2

    tf1 = _body_box(sl, M, HH + Emu(180000), half, SH - HH - Emu(220000))
    _add_line(tf1, "Phase 3 -- OCR-Aware: <confusion_patterns>", 13, bold=True, color=C_DBLUE)
    _add_line(tf1, "", 4)
    _add_line(tf1, "What is injected:", 11, bold=True)
    _add_line(tf1, "* Top character confusion pairs from Phase 1 confusion matrix", 10)
    _add_line(tf1, "* Cross-referenced with LLM failures from training set", 10)
    _add_line(tf1, "* Format: LEFT (OCR form) -> RIGHT (correct form)", 10)
    _add_line(tf1, "", 6)
    _add_line(tf1, "Prompt section injected (example):", 11, bold=True, color=C_ORANGE)
    _add_line(tf1, "<confusion_patterns>", 10, color=C_ORANGE)
    _add_line(tf1, "This OCR system makes these systematic errors.", 9)
    _add_line(tf1, "When you see LEFT, correct to RIGHT:", 9)
    _add_line(tf1, "", 4)
    _add_line(tf1, "ح -> ج    (same shape, dot difference)", 9, rtl=False)
    _add_line(tf1, "ص -> ض   (very similar forms)", 9, rtl=False)
    _add_line(tf1, "ر -> ز    (dot above/below)", 9, rtl=False)
    _add_line(tf1, "ن -> ي    (confusion in cursive)", 9, rtl=False)
    _add_line(tf1, "", 4)
    _add_line(tf1, "WORD-LEVEL FAILURES from training:", 9)
    _add_line(tf1, "صلاة <- حلطة    (OCR -> correct)", 9, rtl=False)
    _add_line(tf1, "وسلم <- وسله    (OCR -> correct)", 9, rtl=False)
    _add_line(tf1, "</confusion_patterns>", 10, color=C_ORANGE)
    _add_line(tf1, "", 8)
    _add_line(tf1, "Result: KHATT normal CER 60.59%->50.72% (+16.3% rel).", 10)
    _add_line(tf1, "Hurts PATS -- over-triggers on already-correct text.", 10, color=C_RED)

    tf2 = _body_box(sl, M + half + M, HH + Emu(180000), half, SH - HH - Emu(220000))
    _add_line(tf2, "Phase 4 -- Self-Reflective: <self_analysis>", 13, bold=True, color=C_DBLUE)
    _add_line(tf2, "", 4)
    _add_line(tf2, "What is injected:", 11, bold=True)
    _add_line(tf2, "* Systematic weaknesses from Phase 2 training analysis", 10)
    _add_line(tf2, "* Known word-level errors (OCR->GT pairs from training)", 10)
    _add_line(tf2, "* Over-correction warnings (errors LLM INTRODUCED)", 10)
    _add_line(tf2, "", 6)
    _add_line(tf2, "Prompt section injected (example):", 11, bold=True, color=C_ORANGE)
    _add_line(tf2, "<self_analysis>", 10, color=C_ORANGE)
    _add_line(tf2, "Based on analysis of this model's corrections:", 9)
    _add_line(tf2, "", 3)
    _add_line(tf2, "SYSTEMATIC WEAKNESSES:", 9, bold=True)
    _add_line(tf2, "* Confuses h/j, s/d, n/y systematically", 9)
    _add_line(tf2, "* Often misses bin names (bn) in isnad chains", 9)
    _add_line(tf2, "", 3)
    _add_line(tf2, "KNOWN WORD-LEVEL ERRORS:", 9, bold=True)
    _add_line(tf2, "حلطة -> صلاة     وسله -> وسلم", 9, rtl=False)
    _add_line(tf2, "اخنا -> اذا      حادثنا -> حدثنا", 9, rtl=False)
    _add_line(tf2, "", 3)
    _add_line(tf2, "OVER-CORRECTION TRAPS (DO NOT change):", 9, bold=True)
    _add_line(tf2, "علي -> على  (model wrongly normalizes)", 9, rtl=False)
    _add_line(tf2, "مريرة -> هريرة  (wrong proper noun fix)", 9, rtl=False)
    _add_line(tf2, "</self_analysis>", 10, color=C_ORANGE)
    _add_line(tf2, "", 8)
    _add_line(tf2, "Result: Best pure-prompt on KHATT normal: 60.59%->47.68%", 10)
    _add_line(tf2, "(-21.3% rel). Still regresses on PATS (more context = more edits).", 10, color=C_RED)


def slide_F_rag_p8(prs, num):
    """RAG Phase 8: BM25 Character N-gram Retrieval"""
    sl = _new_slide(prs)
    _header(sl,
            "RAG Phase 8 — BM25 Character N-gram Retrieval",
            "Retrieves the 5 most similar (OCR->GT) pairs from Phase 2 training corrections",
            num, TOTAL)

    lw = Emu(4800000)
    rw = SW - 3*M - lw

    tf1 = _body_box(sl, M, HH + Emu(180000), lw, SH - HH - M)
    _add_line(tf1, "Pipeline Architecture", 13, bold=True, color=C_DBLUE)
    _add_line(tf1, "", 4)

    steps = [
        ("BUILD INDEX (once)", [
            "Input: 18,363 training (OCR, GT) pairs from Phase 2",
            "Tokenize each OCR text as character 3-grams",
            "  Example: 'ktab' -> ['kta', 'tab']",
            "Build BM25Okapi index over char n-grams",
            "Also extract word-level error pairs (OCR->GT)",
        ]),
        ("QUERY (per sample)", [
            "For each validation OCR text:",
            "  Tokenize as char 3-grams",
            "  BM25 query -> top-5 most similar training texts",
            "  Also retrieve top-10 word-level fixes",
        ]),
        ("INJECT INTO PROMPT", [
            "Inject retrieved sentences as <retrieved_corrections>",
            "Inject word fixes as <retrieved_word_fixes>",
            "LLM sees concrete examples of similar corrections",
        ]),
    ]

    for step_name, sub_steps in steps:
        _add_line(tf1, f">> {step_name}", 11, bold=True, color=C_DBLUE)
        for s in sub_steps:
            _add_line(tf1, f"  {s}", 9, color=C_BODY)
        _add_line(tf1, "", 5)

    _add_line(tf1, "Why char n-grams?", 11, bold=True, color=C_DBLUE)
    _add_line(tf1, "* Robust to OCR noise (word boundaries corrupted)", 9)
    _add_line(tf1, "* Character-level similarity captures confusable patterns", 9)
    _add_line(tf1, "* No embedding model needed (fast, deterministic)", 9)

    tf2 = _body_box(sl, M + lw + M, HH + Emu(180000), rw, SH - HH - M)
    _add_line(tf2, "Concrete Example -- KHATT Sample", 13, bold=True, color=C_DBLUE)
    _add_line(tf2, "", 4)

    _add_line(tf2, "QUERY (validation OCR):", 10, bold=True, color=C_ORANGE)
    _add_line(tf2,
              "وفد كانت حرة نستهلكة عاسًت أمها هلولبِّا كمايخلِّر من وصف العمهايِّ",
              10, rtl=True, color=C_ORANGE)
    _add_line(tf2, "", 5)

    _add_line(tf2, "TOP-2 RETRIEVED CORRECTIONS:", 10, bold=True, color=C_GREEN)
    _add_line(tf2, "INPUT:  oze cab jnt tfq mys ghby khf dght as fzl hzf", 9, rtl=True)
    _add_line(tf2, "OUTPUT: ezh kab jnt tfq mys ghby khf dght as fzl hzf", 9, rtl=True)
    _add_line(tf2, "", 3)
    _add_line(tf2, "INPUT:  fzl hzf ezh kab jnt tfq mys ghby khf dbght as", 9, rtl=True)
    _add_line(tf2, "OUTPUT: fzl hzf ezh kab jnt tfq mys ghby khf dght as", 9, rtl=True)
    _add_line(tf2, "", 5)

    _add_line(tf2, "TOP-5 WORD-LEVEL FIXES:", 10, bold=True, color=C_GREEN)
    word_fixes = [
        ('ضغت', 'ضغط'),
        ('غضبي', 'غضبى'),
        ('مضغتته', 'مضغته'),
        ('حذف', 'دخل'),
        ('عزيزة', 'هريرة'),
    ]
    for ocr_w, gt_w in word_fixes:
        _add_line(tf2, f'  "{ocr_w}" -> "{gt_w}"', 9, rtl=False, color=C_BODY)
    _add_line(tf2, "", 5)

    _add_line(tf2, "RESULT:", 10, bold=True, color=C_DBLUE)
    _add_line(tf2, "Best all-phases: KHATT CER 60.59%->45.91% (normal+diac)", 10, color=C_GREEN)
    _add_line(tf2, "Normal no-diac: 34.24%->33.83% (modest +1.2% rel)", 10, color=C_ORANGE)
    _add_line(tf2, "All-samples: 565.58%->277.27% (51.0% improvement)", 10, color=C_GREEN)
    _add_line(tf2, "Index: 18,363 training pairs", 9, color=C_BODY)


def slide_G_rag_p9(prs, num):
    """RAG Phase 9: Error-Signature Retrieval"""
    sl = _new_slide(prs)
    _header(sl,
            "RAG Phase 9 — Error-Signature Retrieval",
            "Retrieves by structural error similarity rather than text similarity (Phase 8)",
            num, TOTAL)

    lw = Emu(5200000)
    rw = SW - 3*M - lw

    tf1 = _body_box(sl, M, HH + Emu(180000), lw, SH - HH - M)
    _add_line(tf1, "Phase 9 vs Phase 8 -- Key Difference", 13, bold=True, color=C_DBLUE)
    _add_line(tf1, "", 4)

    rows_cmp = [
        ("Retrieval basis",  "Text similarity (char n-grams)", "Error structural similarity"),
        ("Index features",   "BM25 over OCR char 3-grams",    "Error signature vector"),
        ("Query",            "OCR text tokenized as n-grams",  "Predicted error profile"),
        ("GPU needed?",      "No (BM25 only)",                 "No (CAMeL + confusion matrix)"),
        ("KHATT CER (norm)", "45.91% -- BEST",                 "50.47% (worse than P2)"),
    ]

    compare_tbl = _table(sl, M, HH + Emu(350000), lw, 6, 3)
    cw = [Emu(1800000), Emu(1700000), Emu(1700000)]
    for i, w in enumerate(cw):
        compare_tbl.columns[i].width = w
    for i in range(6):
        compare_tbl.rows[i].height = Emu(500000)

    _cell(compare_tbl, 0, 0, "Dimension", bold=True, size=10, color=C_WHITE, bg=C_TABLE_HDR)
    _cell(compare_tbl, 0, 1, "Phase 8 (BM25)", bold=True, size=10, color=C_WHITE, bg=C_DBLUE)
    _cell(compare_tbl, 0, 2, "Phase 9 (Error-Sig)", bold=True, size=10, color=C_WHITE, bg=C_ORANGE)

    for ri, (dim, p8, p9) in enumerate(rows_cmp):
        bg = C_TABLE_ALT if ri % 2 == 1 else None
        _cell(compare_tbl, ri+1, 0, dim, bold=True, size=9, bg=bg, align=PP_ALIGN.LEFT)
        _cell(compare_tbl, ri+1, 1, p8, size=9, bg=bg,
              color=C_GREEN if "BEST" in p8 else C_BODY)
        _cell(compare_tbl, ri+1, 2, p9, size=9, bg=bg,
              color=C_RED if "worse" in p9 else C_BODY)

    tf1b = _body_box(sl, M, HH + Emu(3600000), lw, SH - HH - Emu(3700000))
    _add_line(tf1b, "Phase 9 Error Signature Pipeline:", 12, bold=True, color=C_DBLUE)
    _add_line(tf1b, "", 4)
    steps9 = [
        ("BUILD INDEX", "For each training pair, compute error signature:"),
        ("",            "  * CAMeL invalid word count (morphological)"),
        ("",            "  * Phase 1 confusion-matrix high-freq confused chars"),
        ("",            "  * CER estimate, error density per word"),
        ("QUERY",       "For new OCR text, predict its error signature:"),
        ("",            "  * Run CAMeL on OCR text -> invalid word rate"),
        ("",            "  * Scan for known confusion-matrix characters"),
        ("RETRIEVE",    "Find training samples with closest error signatures"),
        ("",            "  * Cosine similarity over signature vectors"),
    ]
    for s_name, s_desc in steps9:
        if s_name:
            _add_line(tf1b, f">> {s_name}: {s_desc}", 9, bold=True, color=C_DBLUE)
        else:
            _add_line(tf1b, s_desc, 9, color=C_BODY)

    tf2 = _body_box(sl, M + lw + M, HH + Emu(180000), rw, SH - HH - M)
    _add_line(tf2, "Why Phase 9 Underperforms Phase 8", 13, bold=True, color=C_RED)
    _add_line(tf2, "", 6)
    reasons = [
        ("Signature noise",
         "OCR errors corrupt the features used to build the\nerror signature -- predicted profile is unreliable"),
        ("Wrong proxy",
         "Two texts can share the same error type (e.g., h/j\nconfusion) but need completely different corrections"),
        ("Content mismatch",
         "Structural similarity != useful correction examples.\nRetrieved pairs may have similar error patterns but\nunrelated words => unhelpful context for the LLM"),
        ("Text similarity wins",
         "Phase 8 direct text matching finds examples\nwhere actual Arabic words are similar, making\ncorrections directly applicable"),
    ]
    for title, desc in reasons:
        _add_line(tf2, f">> {title}", 11, bold=True, color=C_DBLUE)
        _add_line(tf2, desc, 9, color=C_BODY)
        _add_line(tf2, "", 5)

    _add_line(tf2, "Verdict:", 11, bold=True, color=C_RED)
    _add_line(tf2,
              "Error-signature retrieval is theoretically appealing\n"
              "but fails in practice: OCR-corrupted text makes accurate\n"
              "error profiling impossible. BM25 text similarity (P8)\n"
              "is simpler and works better.",
              10, color=C_BODY)


def slide_H_examples(prs, num):
    """Real Correction Examples: Phase 2 Zero-Shot"""
    sl = _new_slide(prs)
    _header(sl,
            "Real Correction Examples — Phase 2 (Zero-Shot)",
            "Actual OCR -> LLM corrected -> Ground Truth  |  KHATT (handwritten) + PATS (typewritten)",
            num, TOTAL)

    examples = [
        {
            "domain": "KHATT (Handwritten Arabic -- good correction)",
            "color": C_GREEN,
            "ocr": "وهي مساحات من الأراضين تعلوهارمال حمر في الغالب ، تمّدمن النفود",
            "cor":  "وهي مساحات من الأراضي تعلوها رمال حمر في الغالب تمتد من النفوذ",
            "gt":   "وهي مساحات من الأرضين تعلوها رمال حمر في الغالب، تمتد من النفود",
            "analysis": "Fixed merged word 'taloharaml' -> two words. Fixed 'tmdmn' -> 'tmtd mn'. "
                        "Overcorrected 'lnfod' -> 'lnfoz' (GT keeps original). "
                        "Missed 'laradyn' vs 'lardyn' (minor variant).",
        },
        {
            "domain": "PATS-Andalus (Typewritten -- harmful over-correction)",
            "color": C_RED,
            "ocr": "بابه من نسي حلطة فليجل إخنا ذكَرها ولا يعيد إلا تلك الحلطة",
            "cor":  "بابه من نسي حلاوة فليجل إخنا ذكّرها ولا يعيد إلا تلك الحلاوة",
            "gt":   "باب من نسي صلاة فليصل إذا ذكرها ولا يعيد إلا تلك الصلاة",
            "analysis": "LLM changed 'hlth' -> 'hlawt' (wrong -- GT is 'slah'). "
                        "Missed systematic Andalus font errors throughout. "
                        "LLM 'guessed' semantically rather than correcting visually. CER got worse.",
        },
        {
            "domain": "PATS-Thuluth (Typewritten -- partial correction)",
            "color": C_ORANGE,
            "ocr": "منه فلما خرج سكعج سكعئين في قبل الحكعت وقال هذة القبلة",
            "cor":  "منه فلما خرج سكعج سكعئين في قبل الحكمة وقال هذه القبلة",
            "gt":   "منه فلما خرج ركع ركعتين في قبل الكعبة وقال هذه القبلة",
            "analysis": "Fixed 'hzh' -> 'hzh' (diacritic). Fixed 'lhkot' -> 'lhkmt' (partial). "
                        "Missed 'skoaj' -> 'rko', 'skoine' -> 'rkotin', 'lhkobh' -> 'lkobh'. "
                        "Thuluth extreme shape distortions exceed model's visual decoding.",
        },
    ]

    row_h = (SH - HH - Emu(250000)) // 3
    for i, ex in enumerate(examples):
        top = HH + Emu(150000) + i * row_h
        _add_rect(sl, M, top, SW - 2*M, row_h - Emu(50000), RGBColor(0xF5, 0xF5, 0xF5))
        tf = _body_box(sl, M + Emu(100000), top + Emu(40000),
                       SW - 2*M - Emu(200000), row_h - Emu(100000))
        _add_line(tf, ex["domain"], 11, bold=True, color=ex["color"])
        _add_line(tf, "OCR:  " + ex["ocr"], 9, color=C_ORANGE)
        _set_rtl(tf.paragraphs[-1])
        _add_line(tf, "LLM:  " + ex["cor"], 9, color=C_DBLUE)
        _set_rtl(tf.paragraphs[-1])
        _add_line(tf, "GT:   " + ex["gt"], 9, color=C_GREEN)
        _set_rtl(tf.paragraphs[-1])
        _add_line(tf, ex["analysis"], 8, color=C_BODY)


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(str(PPTX_IN))
    existing = len(prs.slides)
    print(f"Loaded PPTX: {existing} existing slides")

    builders = [
        slide_A_runaway,           # 15: Runaway effect + formulas (expanded)
        slide_A2_math_proof,       # 16: Mathematical derivation + PATS-Akhbar case study (NEW)
        slide_B0_slide3_justification,  # 17: Context for results / claim validation (NEW)
        slide_B_baseline,          # 18: Phase 1 baseline -- all 4 metric variants
        slide_C_phases_khatt,      # 19: All-phase KHATT (P1 n/a fixed)
        slide_D_prompt_p2,         # 20: Prompt structure Phase 2
        slide_E_prompt_p3_p4,      # 21: Phase 3 & 4 augmentation
        slide_F_rag_p8,            # 22: RAG Phase 8
        slide_G_rag_p9,            # 23: RAG Phase 9
        slide_H_examples,          # 24: Real correction examples
    ]

    total = existing + len(builders)
    assert total == TOTAL, f"Expected {TOTAL} slides, got {total}"

    for i, builder in enumerate(builders):
        slide_num = existing + 1 + i
        builder(prs, slide_num)
        doc = builder.__doc__.splitlines()[0].strip()
        doc_ascii = doc.encode('ascii', errors='replace').decode('ascii')
        print(f"  Added slide {slide_num}/{total}: {doc_ascii}")

    prs.save(str(PPTX_OUT))
    print(f"\nSaved to: {PPTX_OUT}")
    print(f"Total slides: {total}")


if __name__ == "__main__":
    main()
