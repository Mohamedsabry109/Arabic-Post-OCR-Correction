#!/usr/bin/env python3
"""Generate a dense, research-grade PowerPoint defense presentation.

Run:
    python publication/thesis/generate_defense_pptx.py

Output:
    publication/thesis/defense_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE   = RGBColor(0x1B, 0x3A, 0x5C)  # slide backgrounds / headers
MED_BLUE    = RGBColor(0x2D, 0x6A, 0x9F)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x17)
ACCENT_RED  = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN= RGBColor(0x27, 0xAE, 0x60)
GREY_BG     = RGBColor(0xF2, 0xF3, 0xF4)
DARK_GREY   = RGBColor(0x5D, 0x6D, 0x7E)
TABLE_HDR   = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ALT   = RGBColor(0xEB, 0xF5, 0xFB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                      color=BLACK, bold_prefix=True, line_spacing=1.3):
    """Add a text frame with bullet points. Items can be str or (bold_part, rest)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)

        if isinstance(item, tuple):
            run_b = p.add_run()
            run_b.text = item[0]
            run_b.font.size = Pt(font_size)
            run_b.font.bold = True
            run_b.font.color.rgb = color
            run_b.font.name = "Calibri"
            run_r = p.add_run()
            run_r.text = item[1]
            run_r.font.size = Pt(font_size)
            run_r.font.bold = False
            run_r.font.color.rgb = color
            run_r.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return tf


def _add_table(slide, left, top, width, height, rows, col_widths=None):
    """Add a formatted table. rows[0] = header row."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = BLACK

            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HDR
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table


def _section_header_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 title, font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1),
                     subtitle, font_size=22, color=ACCENT_GOLD,
                     alignment=PP_ALIGN.CENTER)
    return slide


def _content_slide(prs, title, bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, bg_color)
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
                 title, font_size=28, bold=True, color=WHITE)
    # Slide number placeholder area (bottom right)
    slide_num = len(prs.slides)
    _add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 str(slide_num), font_size=11, color=DARK_GREY,
                 alignment=PP_ALIGN.RIGHT)
    return slide


# ===========================================================================
# BUILD PRESENTATION
# ===========================================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ===================================================================
    # SLIDE 1 — TITLE
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(2.0),
                 "Knowledge-Augmented Large Language Models\nfor Arabic Post-OCR Error Correction",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6),
                 "Master's Thesis Defense",
                 font_size=24, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.6),
                 "Mohamed Sabry",
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
                 "Faculty of Engineering, Cairo University  |  mohamedsabry109@gmail.com",
                 font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                 "Supervisor: Prof. Mohsen Rashwan  |  Dept. of Electronics & Electrical Comm., Cairo University & RDI  |  mrashwan@rdi-eg.ai",
                 font_size=15, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6),
                 "2026",
                 font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # SLIDE 2 — OUTLINE
    # ===================================================================
    slide = _content_slide(prs, "Presentation Outline")
    items = [
        "1.  Motivation & Problem Statement",
        "2.  Background: Arabic Script & OCR Challenges",
        "3.  LLM Foundations: Architecture, Training & Why LLMs for Post-OCR",
        "4.  Research Questions",
        "5.  System Design: Three-Stage Pipeline & All Datasets (13 total)",
        "6.  Experiment 1 — Methodology (Phases 1-9)",
        "7.  Experiment 1 — Phase 1 Results: OCR Error Characterisation",
        "8.  Experiment 1 — Phase 2 Results: Zero-Shot LLM Correction",
        "9.  Experiment 1 — Phases 3-5 Results: Isolated Knowledge Augmentation",
        "10. Experiment 1 — Phase 6: Combinations, Ablation & Statistical Significance",
        "11. The Over-Correction Threshold: Theoretical Model",
        "12. Experiment 2 — Prompt Design Study (8 Trials, 13 Datasets)",
        "13. Experiment 3 — Multi-Model, Multi-OCR-Source, Multi-Domain",
        "14. Why Handwritten & Full-Page Correction Differs",
        "15. System Design Recommendations",
        "16. Limitations & Future Work",
        "17. Conclusion & Contributions",
    ]
    _add_bullet_frame(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(6),
                      items, font_size=18, color=BLACK)

    # ===================================================================
    # SECTION: MOTIVATION
    # ===================================================================
    _section_header_slide(prs, "Motivation & Problem Statement",
                          "Why Arabic Post-OCR Correction Matters")

    # SLIDE — Motivation
    slide = _content_slide(prs, "The Arabic OCR Gap")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
        ("Digitisation bottleneck: ", "Arabic archives, government docs, historical texts all depend on OCR as entry point"),
        ("Arabic script challenges: ", "cursive writing, position-dependent glyph shapes, dot-group confusions (ba/ta/tha/ya/na), diacritics, ligatures"),
        ("Open-source vs. closed-source gap: ", "Qaari (open-source VLM-based OCR) produces systematic errors; closed-source engines are expensive and non-reproducible"),
        ("Downstream impact: ", "OCR errors propagate to search, MT, TTS, and information extraction"),
        ("Research question: ", "Can post-OCR correction via a lightweight LLM close this gap without modifying the recognition engine?"),
    ], font_size=16, color=BLACK)

    _add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(1.0),
                 "Central Hypothesis", font_size=20, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)
    # Hypothesis box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = MED_BLUE
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("A 4B-parameter instruction-tuned LLM (Qwen3-4B) can correct "
              "systematic OCR errors through prompt engineering alone. "
              "Injecting linguistic knowledge (confusion matrices, self-reflective "
              "error analysis, morphological validation) into the prompt further "
              "improves correction quality, but the benefit depends on the "
              "baseline OCR error profile.")
    p.font.size = Pt(15)
    p.font.name = "Calibri"
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.LEFT

    # ===================================================================
    # SECTION: BACKGROUND
    # ===================================================================
    _section_header_slide(prs, "Background",
                          "Arabic Script Properties & Related Work")

    # SLIDE — Arabic Script Challenges
    slide = _content_slide(prs, "Arabic Script Properties & OCR Challenges")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), [
        ("Cursive script: ", "letters connect within words; shape depends on position (isolated/initial/medial/final)"),
        ("Dot-group characters: ", "ba/ta/tha/ya/na share same skeleton, differ only in dot count/placement -- 42% of PATS-A01 errors"),
        ("Similar-shape pairs: ", "fa/qaf (1 vs 2 dots above), ra/zay, sad/dad, sin/shin -- 18% of errors"),
        ("Hamza orthography: ", "6+ hamza-carrier combinations create ambiguity -- 12% of errors"),
        ("Taa marbuta / ha: ", "visually near-identical in many fonts -- 9% of errors"),
        ("Segmentation: ", "no inter-character gaps in cursive; word merge/split errors dominate handwritten text (37% of KHATT errors)"),
        ("Font diversity: ", "same text in 8 fonts produces CER ranging 2.12% (Akhbar) to 12.12% (Andalus) -- a 5.7:1 ratio"),
    ], font_size=15, color=BLACK)

    # Right panel: dot confusion illustration
    _add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.6),
                 "Dot-Group Confusion Example", font_size=18, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    tbl_data = [
        ["GT Char", "OCR Char", "Difference", "Count", "Conf. Rate"],
        ["ya", "nun", "dot below vs above", "1,847", "31.4%"],
        ["ba", "ta", "1 dot below vs 2 above", "1,203", "22.7%"],
        ["qaf", "fa", "2 dots vs 1 dot", "891", "18.9%"],
        ["hamza", "alef", "hamza presence", "743", "24.6%"],
        ["taa-m.", "ha", "dot presence", "612", "15.3%"],
    ]
    _add_table(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5), tbl_data)

    # SLIDE — Related Work
    slide = _content_slide(prs, "Related Work & Research Gap")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
        ("Post-OCR correction paradigms:", ""),
        ("  Rule-based: ", "pattern matching, dictionaries -- limited to known error patterns"),
        ("  Statistical: ", "n-gram LMs, noisy channel models -- language-model-guided substitution"),
        ("  Neural seq2seq: ", "encoder-decoder on (OCR, GT) pairs -- requires large training data"),
        ("  LLM-based: ", "prompt engineering on pre-trained LLMs -- THIS THESIS"),
        ("", ""),
        ("Key gap in literature:", ""),
        ("  1. ", "No systematic multi-font evaluation of LLM post-OCR for Arabic"),
        ("  2. ", "No study of over-correction behaviour (when LLM hurts)"),
        ("  3. ", "No ablation of knowledge augmentation components"),
        ("  4. ", "Limited work on Arabic-specific error taxonomy + correction"),
    ], font_size=15, color=BLACK)

    _add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.6),
                 "Our Contributions", font_size=20, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    _add_bullet_frame(slide, Inches(7.2), Inches(2.2), Inches(5.3), Inches(4.5), [
        "1. Systematic 7-phase experimental framework for Arabic post-OCR correction",
        "2. Over-correction threshold: formal model of when LLM correction helps vs hurts",
        "3. Complementarity proof: confusion-matrix + self-reflective prompting are orthogonal",
        "4. 9-dataset evaluation (8 typewritten fonts + 1 handwritten) with aligned splits",
        "5. Three-stage export-infer-analyse pipeline for reproducible GPU-free local analysis",
        "6. Empirical characterisation of the font-dependent correction boundary (~6-7% CER)",
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SECTION: LLM FOUNDATIONS
    # ===================================================================
    _section_header_slide(prs, "LLM Foundations",
                          "Architecture, Training Pipeline & Why LLMs for Post-OCR")

    # SLIDE — Transformer Architecture & Scaling Laws
    slide = _content_slide(prs, "Transformer Architecture & Scaling Laws")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), [
        ("From N-grams to Transformers:", ""),
        ("  N-gram LMs: ", "P(w_t | w_{t-n+1},...,w_{t-1}) -- sparse, limited context window"),
        ("  Word embeddings: ", "dense vectors in R^d encoding semantic relationships"),
        ("  RNNs/LSTMs: ", "sequential processing, vanishing gradients over long sequences"),
        ("  Transformer (Vaswani et al., 2017): ", "self-attention replaces recurrence entirely"),
        ("", ""),
        ("Self-attention mechanism:", ""),
        ("  ", "Q, K, V projections of input; Attention(Q,K,V) = softmax(QK^T / sqrt(d_k)) V"),
        ("  ", "O(n^2) complexity but fully parallelisable -- enables massive scaling"),
        ("  ", "Multi-head attention captures different linguistic aspects simultaneously"),
        ("", ""),
        ("Scaling laws:", ""),
        ("  Kaplan et al. (2020): ", "loss scales as power law with model size, data, and compute"),
        ("  Hoffmann et al. (2022): ", "compute-optimal training balances parameters and tokens"),
    ], font_size=14, color=BLACK)

    _add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.6),
                 "Emergent Capabilities", font_size=18, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    _add_bullet_frame(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5), [
        ("Wei et al. (2022): capabilities emerge", " at scale that are absent in smaller models"),
        ("", ""),
        ("Relevant emergent abilities:", ""),
        ("  In-context learning: ", "adapt to new tasks from prompt examples alone"),
        ("  Instruction following: ", "execute complex multi-step instructions without fine-tuning"),
        ("  Self-correction: ", "identify and revise own outputs when given feedback"),
        ("", ""),
        ("Why this matters for OCR correction:", ""),
        ("  ", "A 4B-parameter model can follow Arabic-specific correction instructions"),
        ("  ", "Knowledge augmentation via prompting leverages in-context learning"),
        ("  ", "No task-specific fine-tuning needed -- zero-shot generalisation"),
    ], font_size=14, color=BLACK)

    # SLIDE — LLM Training Pipeline
    slide = _content_slide(prs, "LLM Training Pipeline: Three Stages")
    # Stage 1
    _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(3.8), Inches(0.5),
                 "Stage 1: Pre-Training", font_size=18, bold=True, color=DARK_BLUE)
    _add_bullet_frame(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(2.0), [
        ("Autoregressive next-token prediction:", ""),
        ("  L = -sum log P(x_t | x_{<t})", ""),
        ("  ", "Trained on trillions of tokens (web, books, code)"),
        ("  ", "Learns grammar, facts, reasoning patterns"),
        ("  ", "Arabic text included in multilingual corpora"),
    ], font_size=13, color=BLACK)

    # Stage 2
    _add_textbox(slide, Inches(4.6), Inches(1.2), Inches(3.8), Inches(0.5),
                 "Stage 2: Instruction Tuning (SFT)", font_size=18, bold=True, color=DARK_BLUE)
    _add_bullet_frame(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(2.0), [
        ("Supervised fine-tuning on (instruction, response) pairs:", ""),
        ("  ", "Transforms base model into instruction-follower"),
        ("  ", "Teaches format compliance, task decomposition"),
        ("  ", "Qwen3-4B-Instruct uses SFT on multilingual tasks"),
        ("  ", "Enables: 'Correct OCR errors in this text'"),
    ], font_size=13, color=BLACK)

    # Stage 3
    _add_textbox(slide, Inches(8.7), Inches(1.2), Inches(4.3), Inches(0.5),
                 "Stage 3: Alignment (RLHF/DPO)", font_size=18, bold=True, color=DARK_BLUE)
    _add_bullet_frame(slide, Inches(8.7), Inches(1.8), Inches(4.3), Inches(2.0), [
        ("Align outputs with human preferences:", ""),
        ("  RLHF: ", "reward model + PPO optimisation"),
        ("  DPO: ", "direct preference optimisation (no reward model)"),
        ("  ", "Reduces harmful outputs, improves helpfulness"),
        ("  ", "Implicitly teaches conservative behaviour"),
    ], font_size=13, color=BLACK)

    # Connection to OCR correction
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = MED_BLUE
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Connection to Post-OCR Correction"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.LEFT

    conn_items = [
        "Pre-training: gives broad Arabic linguistic knowledge (morphology, syntax, common word patterns)",
        "Instruction tuning: enables the model to follow correction prompts faithfully, return text in expected format",
        "Alignment: teaches conservative behaviour -- but also causes OVER-CORRECTION (the model 'helpfully' changes text that was already correct)",
        "Key tension: alignment makes model want to 'help' by changing text, but OCR correction requires knowing when NOT to change",
    ]
    for item in conn_items:
        p2 = tf.add_paragraph()
        p2.text = "  " + item
        p2.font.size = Pt(13)
        p2.font.name = "Calibri"
        p2.font.color.rgb = BLACK
        p2.space_before = Pt(4)

    # SLIDE — Why LLMs for Post-OCR Correction
    slide = _content_slide(prs, "Why LLMs Over Traditional & Deep Learning Approaches")

    # Left: comparison table
    approach_table = [
        ["Approach", "Strengths", "Limitations"],
        ["Rule-based / Statistical",
         "Precise for known patterns\nLow computational cost\nInterpretable",
         "Cannot handle unseen errors\nManual rule engineering\nNo contextual reasoning"],
        ["Neural Seq2Seq",
         "Learns from (OCR, GT) pairs\nHandles complex patterns\nEnd-to-end training",
         "Requires large parallel data\nFragile to domain shift\nNo zero-shot capability"],
        ["LLM-based (THIS THESIS)",
         "Zero-shot generalisation\nIn-context knowledge injection\nMultilingual pre-training\nNo task-specific training data",
         "Over-correction risk\nComputational cost\nToken-level (not image-level)\nHallucination potential"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5), approach_table)

    # Right: 5 advantages
    _add_textbox(slide, Inches(8.2), Inches(1.3), Inches(4.8), Inches(0.5),
                 "5 Key Advantages of LLMs for Arabic OCR", font_size=16, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
    _add_bullet_frame(slide, Inches(8.2), Inches(2.0), Inches(4.8), Inches(5.0), [
        ("1. Broad Arabic knowledge: ", "pre-trained on billions of Arabic tokens -- knows morphology, common words, syntactic patterns"),
        ("", ""),
        ("2. Contextual reasoning: ", "self-attention considers entire sentence when deciding if a word is correct or erroneous"),
        ("", ""),
        ("3. Zero-shot deployment: ", "no (OCR, GT) parallel corpus needed -- works out of the box with prompt engineering"),
        ("", ""),
        ("4. Knowledge injection via prompt: ", "confusion matrices, error statistics, overcorrection warnings all added at inference time"),
        ("", ""),
        ("5. Arabic morphological awareness: ", "instruction-tuned models understand Arabic root-pattern system, enabling valid-but-wrong error detection"),
    ], font_size=13, color=BLACK)

    _add_bullet_frame(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(1.8), [
        ("Bottom line: ", "Rule-based and seq2seq methods require either manual engineering or large parallel corpora. LLMs provide a middle path: broad linguistic competence + controllable behaviour through prompt engineering. The cost is over-correction risk -- which this thesis characterises formally."),
    ], font_size=15, color=BLACK)

    # ===================================================================
    # SECTION: RESEARCH QUESTIONS
    # ===================================================================
    _section_header_slide(prs, "Research Questions",
                          "7 Research Questions Across 6 Experimental Phases")

    slide = _content_slide(prs, "Research Questions")
    rqs = [
        ("RQ1 (Phase 1): ", "What are the characteristics of Qaari OCR errors across fonts and script types?"),
        ("RQ2 (Phase 2): ", "Can a zero-shot LLM correct Arabic OCR errors? How does it vary across fonts?"),
        ("RQ3 (Phase 3): ", "Does injecting the OCR engine's character confusion matrix improve correction?"),
        ("RQ4 (Phase 4): ", "Does feeding the model its own training-split error statistics (self-reflective prompting) improve correction?"),
        ("RQ5 (Phase 5): ", "Does CAMeL Tools morphological post-processing with known-overcorrection revert help?"),
        ("RQ6 (Phase 6): ", "Which combination of strategies is optimal? What does each contribute in ablation?"),
        ("RQ7 (Overarching): ", "Does the baseline OCR error profile determine correction effectiveness, or can knowledge augmentation overcome the over-correction problem?"),
    ]
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8),
                      rqs, font_size=17, color=BLACK)

    # ===================================================================
    # SECTION: SYSTEM DESIGN
    # ===================================================================
    _section_header_slide(prs, "System Design & Datasets",
                          "Three-Stage Pipeline, PATS-A01, KHATT")

    # SLIDE — Pipeline Architecture
    slide = _content_slide(prs, "Three-Stage Pipeline Architecture")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.5), [
        ("Design constraint: ", "no local GPU -- local machine handles data, cloud GPU handles inference"),
        ("Communication: ", "JSONL files (inference_input.jsonl, corrections.jsonl)"),
        ("Benefit: ", "any cloud platform (Kaggle, Colab, Thunder) can run inference; analysis is reproducible locally"),
    ], font_size=15, color=BLACK)

    # Pipeline stages boxes
    stages = [
        ("EXPORT\n(Local)", "Load OCR + GT\nBuild prompts\nWrite JSONL", Inches(0.8), LIGHT_BLUE),
        ("INFERENCE\n(Cloud GPU)", "Load Qwen3-4B\nRun prompts\nAppend corrections", Inches(4.8), RGBColor(0xD5, 0xF5, 0xE3)),
        ("ANALYZE\n(Local)", "Read corrections\nCompute CER/WER\nGenerate reports", Inches(8.8), RGBColor(0xFD, 0xED, 0xEC)),
    ]
    for title, desc, left, bg in stages:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, Inches(4.0), Inches(3.5), Inches(2.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = DARK_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.name = "Calibri"
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    # Arrows between stages
    for x in [Inches(4.3), Inches(8.3)]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x, Inches(5.0), Inches(0.5), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MED_BLUE
        shape.line.fill.background()

    # SLIDE — Datasets (Experiment 1: line-level)
    slide = _content_slide(prs, "Datasets: PATS-A01 & KHATT (Experiment 1 — 9 datasets)")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.2), [
        ("PATS-A01: ", "8 Arabic typewritten fonts, same underlying text, 2,766 lines/font"),
        ("KHATT: ", "handwritten Arabic line images, real-world variation"),
        ("Aligned 85/15 split: ", "seed=42, all 8 fonts share identical 415 validation samples -- font is the ONLY variable"),
        ("OCR engine: ", "Qaari (NAMAA-Space/Qari-OCR-0.1-VL-2B-Instruct) open-source VLM-based Arabic OCR"),
    ], font_size=15, color=BLACK)

    ds_table = [
        ["Dataset", "Type", "N (normal)", "OCR CER*", "OCR WER*", "Runaway %"],
        ["Akhbar",      "Typewritten", "437", "2.12%",  "4.84%",  "4.8%"],
        ["Simplified",  "Typewritten", "432", "2.94%",  "5.81%",  "5.7%"],
        ["Traditional", "Typewritten", "435", "3.63%",  "9.35%",  "5.0%"],
        ["Arial",       "Typewritten", "436", "3.65%",  "6.84%",  "4.6%"],
        ["Naskh",       "Typewritten", "443", "4.10%",  "9.56%",  "3.5%"],
        ["Tahoma",      "Typewritten", "407", "5.82%",  "9.07%", "10.3%"],
        ["Thuluth",     "Typewritten", "441","10.21%", "31.11%",  "3.9%"],
        ["Andalus",     "Typewritten", "427","12.12%", "33.67%",  "6.2%"],
        ["PATS-A01 Avg","--",         "3458", "5.57%", "13.78%",  "5.5%"],
        ["KHATT",       "Handwritten", "186","34.24%", "75.60%", "17.7%"],
    ]
    _add_table(slide, Inches(0.3), Inches(3.5), Inches(12.5), Inches(3.3), ds_table)
    _add_textbox(slide, Inches(0.4), Inches(6.95), Inches(12), Inches(0.4),
                 "*Normal samples only (OCR/GT ratio <= 5.0), diacritics stripped. "
                 "Runaway = Qaari repetition bug (excluded from metrics).",
                 font_size=11, color=DARK_GREY, alignment=PP_ALIGN.LEFT)

    # SLIDE — Additional Datasets (Experiment 2: page-level)
    slide = _content_slide(prs, "Additional Datasets: Full-Page Corpora (Experiment 2 — 13 datasets total)")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.2), [
        ("KHATT-Paragraph: ", "449 multi-line handwritten paragraph images from KHATT corpus -- paragraph-level context"),
        ("Yarmouk: ", "1,663 contemporary printed encyclopedic pages -- diverse layouts, mixed content"),
        ("Muharaf: ", "116 handwritten Arabic manuscript images -- informal handwriting styles"),
        ("Historical: ", "40 archival manuscript pages from 8 Arabic books -- degraded historical print"),
        ("Challenge: ", "Qaari's runaway bug dominant on page-level input (CER >100% raw) -> runaway corrector essential"),
    ], font_size=14, color=BLACK)

    fp_table = [
        ["Dataset", "Domain", "Samples", "OCR CER (raw)", "T2 CER (corrected)"],
        ["KHATT-Paragraph", "HW paragraph",     "449",   "61.68%",  "45.06%"],
        ["Yarmouk-testing", "Print (encyclop.)", "1,663", "49.85%",  "43.35%"],
        ["Muharaf-val",     "HW document",       "116",  "129.39%", "89.78%"],
        ["Historical",      "HW manuscript",     "40",   "105.10%", "94.53%"],
        ["Full-page Avg",   "---",               "---",   "86.50%",  "68.18%"],
    ]
    _add_table(slide, Inches(0.3), Inches(3.6), Inches(12.5), Inches(2.8), fp_table)
    _add_textbox(slide, Inches(0.4), Inches(6.55), Inches(12), Inches(0.8),
                 "T2 corrected = after runaway correction (threshold 3.0x GT length). "
                 "Full-page gain (-21.2% abs) >> line-level PATS gain (-0.13 pp). "
                 "Gemma VLM OCR excluded for PATS/KHATT (line strips, >10:1 aspect ratio, invalid).",
                 font_size=11, color=DARK_GREY, alignment=PP_ALIGN.LEFT)

    # ===================================================================
    # SECTION: METHODOLOGY
    # ===================================================================
    _section_header_slide(prs, "Experimental Methodology",
                          "Phases 1-6: Isolated Variables, Controlled Comparisons")

    # SLIDE — Phase Overview
    slide = _content_slide(prs, "Experimental Phases Overview")
    phase_table = [
        ["Phase", "Name", "What Changes vs Phase 2", "Knowledge Injected", "Comparison"],
        ["1", "Baseline & Error Taxonomy", "No LLM", "None", "N/A"],
        ["2", "Zero-Shot LLM", "Baseline LLM", "None (conservative v2 prompt)", "vs Phase 1"],
        ["3", "OCR-Aware Prompting", "+Confusion matrix", "Top-10 failure-filtered char pairs", "vs Phase 2"],
        ["4", "Self-Reflective", "+Training-split insights", "Error stats + word-pair failures + overcorrection warning", "vs Phase 2"],
        ["5", "CAMeL Post-Processing", "+Morphological filter", "CAMeL validation + known-overcorrection revert", "vs Phase 2"],
        ["6", "Combinations + Ablation", "Multi-component", "conf_only, self_only, conf_self, best_camel", "vs Phase 2 + ablation"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(12.5), Inches(4.0), phase_table)

    _add_bullet_frame(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), [
        ("Key design: ", "Phases 3-5 are ISOLATED experiments -- each changes exactly ONE variable vs Phase 2. Phase 6 tests combinations."),
        ("Evaluation: ", "CER/WER on validation split, diacritics stripped. PATS-A01 macro-average (equal weight per font). Paired t-tests + Bonferroni correction."),
    ], font_size=15, color=BLACK)

    # SLIDE — Prompt Strategy Details
    slide = _content_slide(prs, "Prompt Engineering Strategies: Detail")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), [
        ("Phase 2 -- Zero-Shot (v2 conservative):", ""),
        ("  ", '"If the text is correct, return it as-is. Only fix clear OCR errors."'),
        ("  ", "Designed to minimise over-correction on clean input"),
        ("", ""),
        ("Phase 3 -- Confusion Matrix Injection:", ""),
        ("  ", "Top-10 char confusion pairs PER FONT, filtered by LLM failures"),
        ("  ", "Only includes pairs the model does NOT handle at zero-shot level"),
        ("  ", "Tells model WHAT to look for"),
        ("", ""),
        ("Phase 4 -- Self-Reflective Prompting:", ""),
        ("  ", "Arabic-language diagnostic derived from training-split Phase 2 errors"),
        ("  ", "Includes: aggregate stats, concrete word-pair failures, overcorrection warning"),
        ("  ", "Tells model HOW to behave (be conservative, avoid specific error types)"),
    ], font_size=14, color=BLACK)

    _add_bullet_frame(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.8), [
        ("Phase 5 -- CAMeL Morphological Post-Processing:", ""),
        ("  ", "Applied AFTER LLM correction (not during prompting)"),
        ("  ", "Step 1: Validate each changed word via CAMeL morphological analyser"),
        ("  ", "Step 2: If LLM output is morphologically invalid AND OCR original was valid, REVERT"),
        ("  ", "Step 3: Known-overcorrection list -- revert specific word pairs known to be over-corrected"),
        ("", ""),
        ("Phase 6 -- Combinations:", ""),
        ("  conf_only: ", "Phase 3 confusion matrix only"),
        ("  self_only: ", "Phase 4 self-reflective only"),
        ("  conf_self: ", "Phase 3 + Phase 4 combined in prompt"),
        ("  best_camel: ", "conf_self + Phase 5 CAMeL post-processing"),
        ("", ""),
        ("Design rationale: ", "4 combinations test whether components are redundant, additive, or synergistic"),
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SECTION: PHASE 1 RESULTS
    # ===================================================================
    _section_header_slide(prs, "Phase 1 Results",
                          "OCR Error Characterisation")

    slide = _content_slide(prs, "Phase 1: Error Type Distribution")
    err_table = [
        ["Error Category", "PATS-A01 (%)", "KHATT (%)"],
        ["Dot confusion (ba/ta/tha/ya/fa/qaf)", "42.1%", "18.4%"],
        ["Similar-shape substitution", "18.3%", "14.7%"],
        ["Hamza / alef variants", "11.6%", "9.2%"],
        ["Taa marbuta / ha confusion", "9.0%", "6.8%"],
        ["Alef maqsura / ya confusion", "7.2%", "5.1%"],
        ["Word merge (segmentation)", "3.8%", "19.6%"],
        ["Word split (segmentation)", "2.6%", "17.3%"],
        ["Other (insertion, deletion, misc.)", "5.4%", "8.9%"],
    ]
    _add_table(slide, Inches(0.5), Inches(1.3), Inches(7), Inches(3.5), err_table)

    _add_bullet_frame(slide, Inches(8.0), Inches(1.3), Inches(5), Inches(5.5), [
        ("Key findings:", ""),
        ("", ""),
        ("PATS-A01: ", "60%+ of errors are dot + similar-shape confusions -- targetable by confusion matrix"),
        ("", ""),
        ("KHATT: ", "37% segmentation errors (merge + split) -- fundamentally harder for text-only LLM"),
        ("", ""),
        ("Morphological validity:", ""),
        ("  PATS-A01: ", "71.3% of OCR errors are valid-but-wrong Arabic words"),
        ("  KHATT: ", "63.8% valid-but-wrong"),
        ("", ""),
        ("Implication: ", "CAMeL morphological filter can only intercept the 28.7% non-word errors -- ceiling on Phase 5 effectiveness"),
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SECTION: PHASE 2 RESULTS
    # ===================================================================
    _section_header_slide(prs, "Phase 2 Results",
                          "Zero-Shot LLM Correction")

    slide = _content_slide(prs, "Phase 2: Per-Font Zero-Shot Correction")
    p2_table = [
        ["Font", "P1 CER", "P2 CER", "Delta (pp)", "Relative Change"],
        ["Akhbar", "2.12%", "2.71%", "+0.59", "+27.8% (HARMED)"],
        ["Simplified", "2.94%", "2.93%", "-0.01", "-0.3%"],
        ["Traditional", "3.63%", "2.96%", "-0.67", "-18.5%"],
        ["Arial", "3.65%", "4.50%", "+0.85", "+23.3% (HARMED)"],
        ["Naskh", "4.10%", "4.27%", "+0.17", "+4.1% (HARMED)"],
        ["Tahoma", "5.82%", "6.13%", "+0.31", "+5.3%"],
        ["Thuluth", "10.21%","10.92%","+0.71", "+7.0%"],
        ["Andalus", "12.12%","12.26%","+0.14", "+1.2%"],
        ["PATS-A01 Avg", "5.57%", "5.84%", "+0.27", "+4.8% (HARMED)"],
        ["KHATT", "34.24%","33.25%", "-0.99", "-2.9%"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(8), Inches(3.8), p2_table)

    _add_bullet_frame(slide, Inches(8.5), Inches(1.3), Inches(4.5), Inches(5.5), [
        ("Critical findings:", ""),
        ("", ""),
        ("Traditional: ", "-18.5% relative CER (3.63% -> 2.96%) -- Qaari errors on this font are 'obvious' to the LLM"),
        ("", ""),
        ("Akhbar/Arial/Naskh: ", "harmed by zero-shot -- over-correction on clean input (CER < 4%)"),
        ("", ""),
        ("KHATT: ", "-2.9% relative (34.24% -> 33.25%) -- segmentation errors limit text-only correction"),
        ("", ""),
        ("Over-correction threshold: ", "between 4.3% (harmed) and 7.9% (helped) -- approx 6-7% CER"),
    ], font_size=14, color=BLACK)

    # SLIDE — Error Introduction Rates
    slide = _content_slide(prs, "Phase 2: Fix Rate vs Introduction Rate")
    intro_table = [
        ["Font", "Fix Rate (%)", "Intro Rate (%)", "Net Effect"],
        ["Akhbar", "62.3%", "25.1%", "Net HARM (few errors to fix)"],
        ["Simplified", "74.6%", "21.8%", "Net benefit"],
        ["Traditional", "91.2%", "8.4%", "Strong benefit (high fix, low intro)"],
        ["Naskh", "80.1%", "20.6%", "Moderate benefit"],
        ["KHATT", "58.4%", "22.1%", "Marginal benefit"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(7.5), Inches(2.8), intro_table)

    _add_bullet_frame(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(2.5), [
        ("Akhbar paradox: ", "62.3% fix rate sounds good, but with only 4.3% baseline CER, the 25.1% introduction rate on the 95.7% correct tokens causes NET harm"),
        ("Traditional success: ", "91.2% fix rate + only 8.4% introduction rate = dramatic improvement. Qaari's errors on this font are systematically 'obvious' to the LLM"),
        ("Introduction rate is ~constant: ", "~20-25% across most fonts -- it depends on model behaviour, not input error rate"),
    ], font_size=15, color=BLACK)

    # ===================================================================
    # SECTION: PHASES 3-5
    # ===================================================================
    _section_header_slide(prs, "Phases 3-5 Results",
                          "Isolated Knowledge Augmentation Experiments")

    slide = _content_slide(prs, "Phases 3-5: Isolated Results Summary")
    iso_table = [
        ["Phase", "Method", "PATS CER", "Delta vs P1", "KHATT CER", "Delta vs P1"],
        ["1", "OCR Baseline (Qaari)", "5.57%", "--", "34.24%", "--"],
        ["2", "Zero-Shot LLM", "5.84%", "+0.27%", "33.25%", "-0.99%"],
        ["3", "OCR-Aware (Confusion)", "5.74%", "-0.83% rel", "33.26%", "-0.98%"],
        ["4", "Self-Reflective (BEST)", "5.59%", "-0.37% rel", "33.24%", "-1.00%"],
        ["5", "CAMeL Post-Proc. (buggy†)", "21.89%", "+293% (broken)", "41.62%", "+21.6%"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(12.5), Inches(2.5), iso_table)

    _add_bullet_frame(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(3.0), [
        ("Phase 3 (Confusion Matrix):", ""),
        ("  ", "Improves 7/9 datasets; only Akhbar degraded (+0.2pp)"),
        ("  ", "Naskh benefits most (-0.8pp) -- most confusion pairs"),
        ("  ", "Failure-filtering prevents injection of already-handled pairs"),
        ("", ""),
        ("Phase 4 (Self-Reflective) -- STRONGEST ISOLATED:", ""),
        ("  ", "Improves ALL 9 datasets including Akhbar (-0.3pp)"),
        ("  ", "ONLY strategy that improves Akhbar -- overcorrection warning works"),
        ("  ", "Naskh: -1.7pp (2x the benefit of Phase 3)"),
    ], font_size=14, color=BLACK)

    _add_bullet_frame(slide, Inches(6.8), Inches(4.2), Inches(5.8), Inches(3.0), [
        ("Phase 5 (CAMeL Post-Processing):", ""),
        ("  ", "Slight degradation on most fonts (+0.1pp avg)"),
        ("  ", "Only 3.4% of changed words reverted (PATS-A01)"),
        ("  ", "71.3% valid-but-wrong errors pass morphological filter"),
        ("  ", "Revert occasionally removes correct LLM fixes"),
        ("", ""),
        ("Key insight: ", "Phase 4's success comes from TWO mechanisms: (1) targeted error guidance, (2) explicit overcorrection warning that shifts model behaviour towards conservatism"),
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SECTION: PHASE 6
    # ===================================================================
    _section_header_slide(prs, "Phase 6 Results",
                          "Combinations, Ablation & Statistical Significance")

    # SLIDE — Combination Results
    slide = _content_slide(prs, "Phase 6: Combination Results")
    combo_table = [
        ["Combination", "PATS CER", "Δ vs P2 (pp)", "KHATT CER", "Δ vs P2 (pp)"],
        ["OCR baseline (P1)", "5.57%", "—", "34.24%", "—"],
        ["P2 (zero-shot)", "5.84%", "—", "33.25%", "—"],
        ["conf_only (P3)", "5.74%", "-0.10", "33.26%", "+0.01"],
        ["self_only (P4)", "5.59%", "-0.25", "33.24%", "-0.01"],
        ["conf_self (P3+P4)", "5.76%", "-0.08", "33.27%", "+0.02"],
        ["P8 RAG BM25 -- BEST PATS", "5.45%", "-0.39", "33.83%", "+0.58"],
        ["P9 Error-Sig RAG -- BEST KHATT", "5.47%", "-0.37", "33.13%", "-0.12"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(8.5), Inches(2.5), combo_table)

    _add_bullet_frame(slide, Inches(9.0), Inches(1.3), Inches(4), Inches(2.5), [
        ("conf_self is best balanced:", ""),
        ("  PATS-A01: ", "5.45% P8 RAG; 5.59% P4 Self-Reflective (both < OCR baseline 5.57%)"),
        ("  KHATT: ", "33.13% P9 Error-Sig RAG (best); −1.11% vs OCR baseline 34.24%"),
        ("", ""),
        ("best_camel: ", "+0.1pp on PATS but -0.1pp on KHATT"),
    ], font_size=14, color=BLACK)

    # Per-font combo table
    pf_combo = [
        ["Font", "P1 (OCR)", "P2 ZS", "P3 OCR-Aware", "P4 Self-Ref", "P6 conf_self", "P8 RAG"],
        ["Akhbar",     "2.12%", "2.71%", "2.57%", "2.59%", "2.46%", "3.01%"],
        ["Simplified", "2.94%", "2.93%", "2.97%", "2.86%", "2.72%", "3.01%"],
        ["Traditional","3.63%", "2.96%", "2.90%", "2.87%", "2.75%", "3.07%"],
        ["Tahoma",     "5.82%", "6.13%", "6.12%", "5.91%", "6.15%", "6.21%"],
        ["Thuluth",   "10.21%","10.92%","10.94%","10.67%","10.56%", "9.68%"],
        ["Andalus",   "12.12%","12.26%","11.75%","11.51%","13.22%","10.41%"],
        ["KHATT",     "34.24%","33.25%","33.26%","33.24%","33.27%","33.83%"],
    ]
    _add_table(slide, Inches(0.3), Inches(4.2), Inches(12.5), Inches(3.0), pf_combo)

    # SLIDE — Ablation + Statistical Significance
    slide = _content_slide(prs, "Phase 6: Ablation & Statistical Significance")
    # Ablation table (left)
    abl_table = [
        ["Configuration", "PATS CER", "vs conf_self"],
        ["conf_self (P3+P4)", "5.76%", "—"],
        ["Remove P3 (= P4 only)", "5.59%", "-0.17 pp (better!)"],
        ["Remove P4 (= P3 only)", "5.74%", "-0.02 pp (better!)"],
        ["P8 RAG (best overall)", "5.45%", "reference best"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(6), Inches(2.0), abl_table)

    _add_bullet_frame(slide, Inches(0.5), Inches(3.5), Inches(5.5), Inches(1.5), [
        ("Key ablation finding: ", "on PATS macro-avg, P4 alone (5.59%) beats conf_self (5.76%). The combination is slightly WORSE than P4 alone."),
        ("Near-additive on specific fonts: ", "Naskh: P3 alone -0.8 + P4 alone -1.7 = -2.5 predicted; actual conf_self Naskh = -2.4 (near-additive)"),
        ("KHATT: ", "conf_self (33.27%) improves over both P3 (33.26%) and P4 (33.24%) alone"),
    ], font_size=13, color=BLACK)

    # Stats table (right)
    stats_table = [
        ["Dataset", "p-value", "Cohen's d", "Significant?"],
        ["Thuluth", "0.002", "0.41", "Yes*"],
        ["Traditional", "<0.001", "1.84", "Yes*"],
        ["Arial", "0.011", "0.31", "Yes*"],
        ["Andalus", "0.009", "0.33", "Yes*"],
        ["Tahoma", "0.008", "0.36", "Yes*"],
        ["Naskh", "<0.001", "0.58", "Yes*"],
        ["Akhbar", "0.038", "0.21", "No"],
        ["Simplified", "0.019", "0.29", "No"],
        ["KHATT", "0.041", "0.19", "No"],
    ]
    _add_table(slide, Inches(6.5), Inches(1.3), Inches(6.3), Inches(3.5), stats_table)

    _add_textbox(slide, Inches(6.5), Inches(5.0), Inches(6.3), Inches(0.5),
                 "Bonferroni-corrected alpha = 0.0125 (4 comparisons). Significant on 6/9 datasets.",
                 font_size=13, color=DARK_GREY, alignment=PP_ALIGN.LEFT)

    # ===================================================================
    # SECTION: EXPERIMENT 2
    # ===================================================================
    _section_header_slide(prs, "Experiment 2: Prompt Design Study",
                          "8 Trials x 13 Datasets — Finding the Best Prompt Configuration")

    # SLIDE — Experiment 2 Trial Definitions
    slide = _content_slide(prs, "Experiment 2: Trial Definitions & PATS/KHATT Results")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(4.5), Inches(2.0), [
        ("Motivation: ", "Experiment 1 used a fixed zero-shot baseline. Experiment 2 tests the space of prompt design choices."),
        ("Design: ", "4 prompt styles x 2 retrieval modes = 8 trials (T1-T8)"),
        ("Scale: ", "13 datasets (PATS-A01 8 fonts + KHATT + 4 full-page)"),
    ], font_size=14, color=BLACK)

    exp2_table = [
        ["Trial", "Style", "RAG", "PATS CER", "KHATT CER", "Note"],
        ["OCR",  "Baseline",      "—",   "5.57%",   "34.24%", "—"],
        ["T1",   "Base (aggress.)", "No", "5.84%",   "33.25%", "= Exp1 P2"],
        ["T2",   "Base",           "Yes", "5.44%",   "33.82%", "Best overall"],
        ["T3*",  "Conservative",   "No",  "5.27%",   "32.93%", "Best PATS!"],
        ["T4",   "Conservative",   "Yes", "27.02%",  "32.69%", "CATASTROPHIC"],
        ["T5",   "Error-pattern",  "No",  "5.85%",   "32.90%", "Marginal"],
        ["T6",   "Error-pattern",  "Yes", "19.64%",  "32.64%", "CATASTROPHIC"],
        ["T7",   "EP + Cons.",     "No",  "190.60%", "51.64%", "TEMPLATE INJECT"],
        ["T8",   "EP + Cons.",     "Yes", "46.56%",  "32.54%", "CATASTROPHIC"],
    ]
    _add_table(slide, Inches(0.3), Inches(3.3), Inches(12.5), Inches(3.8), exp2_table)
    _add_textbox(slide, Inches(0.4), Inches(7.2), Inches(12), Inches(0.25),
                 "T3* best PATS (5.27%, -5.4% rel vs OCR). T2 best overall (5.44% PATS + 33.82% KHATT). T2 selected for Experiment 3.",
                 font_size=12, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.LEFT)

    # SLIDE — Experiment 2 per-font detail
    slide = _content_slide(prs, "Experiment 2: Per-Font PATS Detail (T1, T2, T3)")
    pf_exp2 = [
        ["Font", "OCR", "T1 (base_zs)", "T2 (base_rag)", "T3 (cons_zs)*"],
        ["Akhbar",     "2.12%", "2.71%", "2.97%", "2.32%"],
        ["Simplified", "2.94%", "2.94%", "3.01%", "2.43%"],
        ["Traditional","3.63%", "2.97%", "3.07%", "2.58%"],
        ["Arial",      "3.65%", "4.50%", "3.83%", "3.74%"],
        ["Naskh",      "4.10%", "4.26%", "4.36%", "3.69%"],
        ["Tahoma",     "5.82%", "6.14%", "6.20%", "5.54%"],
        ["Thuluth",   "10.21%","10.93%", "9.68%","10.23%"],
        ["Andalus",   "12.12%","12.26%","10.42%","11.68%"],
        ["PATS Avg",   "5.57%", "5.84%", "5.44%", "5.27%"],
        ["KHATT",     "34.24%","33.25%","33.82%","32.93%"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(8.5), Inches(5.5), pf_exp2)

    _add_bullet_frame(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.5), [
        ("T3 key insight:", ""),
        ("  ", "Conservative framing ('return as-is if correct') beats ALL knowledge augmentation from Experiment 1"),
        ("  ", "T3 achieves 5.27% PATS -- below P8 RAG (5.45%) without any retrieval overhead"),
        ("", ""),
        ("T2 vs T1:", ""),
        ("  ", "RAG reduces Thuluth: 10.93% -> 9.68% (-1.25 pp), Andalus: 12.26% -> 10.42% (-1.84 pp)"),
        ("  ", "RAG costs Akhbar: 2.71% -> 2.97% (+0.26 pp)"),
        ("", ""),
        ("Lesson: ", "Reducing intervention tendency > providing more correct targets for clean typewritten text"),
    ], font_size=13, color=BLACK)

    # SLIDE — Experiment 2 Full-Page Datasets
    slide = _content_slide(prs, "Experiment 2: Full-Page Dataset Results (New Domains)")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.0), [
        ("4 new full-page datasets introduced: ", "KHATT-Paragraph, Yarmouk, Muharaf, Historical"),
        ("Challenge: ", "Qaari's runaway bug severe on page-level input -> OCR CER >100% (Muharaf 129%, Historical 105%)"),
        ("Solution: ", "3-stage runaway corrector (tatweel collapse + phrase-loop detection + hard truncation at 3.0x GT length)"),
        ("Best trial: ", "T2 (base+RAG) selected as the best overall prompt for all domains"),
    ], font_size=14, color=BLACK)

    fp_exp2_table = [
        ["Dataset", "Domain", "Samples", "OCR CER", "T2 CER (corrected)", "Abs. Gain"],
        ["KHATT-Paragraph", "HW paragraph",    "449",   "61.68%",  "45.06%", "-16.62 pp"],
        ["Yarmouk-testing", "Print (encyclop.)","1,663","49.85%",  "43.35%",  "-6.50 pp"],
        ["Muharaf-val",     "HW document",      "116",  "129.39%", "89.78%", "-39.61 pp"],
        ["Historical",      "HW manuscript",    "40",   "105.10%", "94.53%", "-10.57 pp"],
        ["Full-page Avg",   "---",              "---",   "86.50%",  "68.18%", "-18.32 pp"],
    ]
    _add_table(slide, Inches(0.3), Inches(3.4), Inches(12.5), Inches(2.8), fp_exp2_table)
    _add_bullet_frame(slide, Inches(0.5), Inches(6.4), Inches(12), Inches(0.8), [
        ("Full-page gain (-21.2% relative) >> line-level PATS (-2.3%) or KHATT (-1.2%) -- "
         "runaway correction accounts for bulk; LLM adds genuine character-level correction on top. "
         "KHATT-Paragraph benefits most from paragraph context (-16.62 pp genuine correction)."),
    ], font_size=13, color=BLACK)

    # ===================================================================
    # SECTION: EXPERIMENT 3
    # ===================================================================
    _section_header_slide(prs, "Experiment 3: Final Evaluation",
                          "Multi-Model · Multi-OCR-Source · Multi-Domain Benchmarks")

    # SLIDE — Experiment 3A: Model Comparison
    slide = _content_slide(prs, "Experiment 3A: Four Models vs Validation Set (T2 Prompt, Qaari OCR)")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(1.5), [
        ("Setting: ", "T2 (base+RAG), validation set, Qaari OCR, 4 LLMs"),
        ("Full-page group (FP): ", "KHATT-Para + Yarmouk + Muharaf + Historical, OCR baseline 86.50%"),
        ("CER†: ", "after runaway correction (threshold 3.0x GT length)"),
    ], font_size=14, color=BLACK)

    m3a_table = [
        ["Model", "Params", "PATS CER†", "KHATT CER†", "FP CER†", "Notes"],
        ["OCR Baseline", "—",  "5.57%",  "34.24%", "86.50%", "Qaari raw OCR"],
        ["Qwen3-4B",     "4B", "5.44%",  "33.82%", "68.18%", "Best KHATT + FP"],
        ["Qwen3-14B",   "14B", "5.60%",  "35.55%", "88.51%", "Runaway epidemic (175% raw KHATT)"],
        ["Gemma-3-4B",   "4B", "7.47%",  "36.14%", "71.12%", "Underperforms consistently"],
        ["Gemma-3-12B", "12B", "5.33%",  "35.98%", "70.01%", "Best PATS"],
    ]
    _add_table(slide, Inches(0.3), Inches(2.9), Inches(12.5), Inches(2.8), m3a_table)

    _add_bullet_frame(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(1.3), [
        ("Key finding: ", "Larger != better within Qwen3. Qwen3-14B (raw 175.69% KHATT) needs runaway corrector to reach 35.55% -- still worse than Qwen3-4B (33.82%). Gemma-3-12B best PATS (5.33%); Qwen3-4B best KHATT + Full-page."),
        ("Gemma excluded for PATS/KHATT: ", "line-strip images (>10:1 aspect ratio) squashed by Gemma preprocessor -> hallucinations (CER >170%, invalid). Only valid for full-page (~0.6:1) comparison."),
    ], font_size=13, color=BLACK)

    # SLIDE — Experiment 3B: OCR Source Quality
    slide = _content_slide(prs, "Experiment 3B: OCR Source Impact (Qaari vs Gemma-3 VLM)")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.2), [
        ("Key design: ", "Gemma-3 VLM OCR (Gemma-3-27B-IT vision encoder) vs Qaari on same images"),
        ("Critical exclusion: ", "Gemma excluded for PATS/KHATT -- line strips >10:1 aspect ratio squashed to near-square by preprocessor -> hallucinations (CER >170%)"),
        ("Valid comparison: ", "Full-page documents (~0.6:1 aspect ratio) only"),
        ("Finding: ", "OCR source quality matters less after LLM correction for full-page documents"),
    ], font_size=14, color=BLACK)

    m3b_table = [
        ["Dataset", "Qaari OCR", "Qaari Corr†", "Gemma OCR", "Gemma Corr†"],
        ["KHATT-Para-val",  "61.68%",  "45.06%",  "36.14%",  "35.82%"],
        ["Yarmouk-testing", "49.85%",  "43.35%",  "95.76%",  "74.37%"],
        ["Muharaf-val",    "129.39%",  "89.78%", "141.15%",  "72.85%"],
        ["Historical",     "105.10%",  "94.53%",  "76.26%",  "74.09%"],
        ["Full-page Avg",   "86.50%",  "68.18%",  "87.33%",  "64.28%"],
    ]
    _add_table(slide, Inches(0.3), Inches(3.4), Inches(12.5), Inches(2.5), m3b_table)
    _add_bullet_frame(slide, Inches(0.5), Inches(6.1), Inches(12), Inches(1.0), [
        ("Gemma OCR slightly better after correction (64.28% vs 68.18%). "
         "Qaari better on KHATT-Para and Yarmouk; Gemma better on Muharaf and Historical. "
         "For full-page docs, OCR source doesn't dominate after LLM correction."),
    ], font_size=13, color=BLACK)

    # SLIDE — Experiment 3C: Generalization
    slide = _content_slide(prs, "Experiment 3C: Generalization to Unseen Benchmarks")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(1.5), [
        ("Two held-out benchmarks: ", "RDI-Test-Lines (handwritten/manuscript/typewritten) + Kitab (12 diverse Arabic text collections)"),
        ("Model: ", "Qwen3-4B, T2 prompt, runaway-corrected"),
        ("Surprising finding: ", "Gemma OCR outperforms Qaari on RDI-Test-Lines (68% vs 95% CER) -- Qaari struggles on diverse document styles"),
    ], font_size=14, color=BLACK)

    m3c_left = [
        ["RDI-Test-Lines", "OCR Source", "OCR CER", "Corr CER†"],
        ["LR-Handwritten",  "Qaari", "98.82%", "91.55%"],
        ["LR-Handwritten",  "Gemma", "62.53%", "57.71%"],
        ["LR-Manuscripts",  "Qaari", "81.94%", "79.00%"],
        ["LR-Manuscripts",  "Gemma", "78.10%", "70.36%"],
        ["LR-Typewritten",  "Qaari","105.19%", "84.79%"],
        ["LR-Typewritten",  "Gemma", "64.86%", "55.31%"],
        ["OVERALL",         "Qaari", "95.31%", "85.11%"],
        ["OVERALL",         "Gemma", "68.49%", "61.13%"],
    ]
    _add_table(slide, Inches(0.3), Inches(2.9), Inches(6.5), Inches(3.8), m3c_left)

    m3c_right = [
        ["Kitab Subset", "Qaari OCR", "Qaari Corr†", "Note"],
        ["kitab-arabicocr",   "1.80%",    "7.21%",  "OVER-CORRECT"],
        ["kitab-patsocr",     "4.79%",    "5.98%",  "Over-correct"],
        ["kitab-hindawi",    "31.46%",   "24.31%",  "-7.2 pp"],
        ["kitab-muharaf",    "76.04%",   "58.72%",  "-17.3 pp"],
        ["kitab-historyar",  "63.48%",   "51.69%",  "-11.8 pp"],
        ["khattparagraph",  "205.27%",   "90.60%",  "-114.7 pp!"],
        ["Kitab OVERALL",    "49.28%",   "40.73%",  "-17.3% rel"],
    ]
    _add_table(slide, Inches(7.0), Inches(2.9), Inches(6.0), Inches(3.8), m3c_right)
    _add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                 "Over-correction confirmed on new domains: low-CER subsets (kitab-arabicocr 1.80%) are HARMED. "
                 "RDI-Test-Lines -10.7% abs. Kitab -17% rel (Qaari) to -24% rel (Gemma).",
                 font_size=12, color=ACCENT_RED, alignment=PP_ALIGN.LEFT)

    # SLIDE — Experiment 3 Summary
    slide = _content_slide(prs, "Experiment 3: Key Takeaways")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5), [
        ("1. Model size != always better (Qwen3 family):", ""),
        ("   ", "Qwen3-14B underperforms Qwen3-4B on KHATT (35.55% vs 33.82%) despite 3.5x more parameters"),
        ("   ", "Gemma-3-12B best PATS (5.33%); Qwen3-4B best for KHATT + full-page"),
        ("", ""),
        ("2. Runaway corrector is critical for large models:", ""),
        ("   ", "Without it: Qwen3-14B = 175.69% KHATT CER (unusable). With it: 35.55% (still worse than 4B)"),
        ("   ", "Full-page datasets: OCR 86.50% -> 68.18% (large share is runaway correction)"),
        ("", ""),
        ("3. OCR source matters less after correction on full-page:", ""),
        ("   ", "Qaari 86.50% vs Gemma 87.33% -> after correction: Qaari 68.18% vs Gemma 64.28%"),
        ("   ", "Gemma excluded for line-strip PATS/KHATT (preprocessing artifact, invalid comparison)"),
        ("", ""),
        ("4. Generalisation confirmed on unseen benchmarks:", ""),
        ("   ", "RDI-Test-Lines: -10.7% abs. Kitab: -17% to -24% relative CER reduction"),
        ("   ", "Over-correction threshold holds in new domains: low-CER inputs (1.80%, 4.79%) are harmed"),
        ("   ", "Gemma OCR actually BETTER than Qaari on RDI-Test-Lines (68% vs 95%) -- Qaari is domain-specialised"),
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SECTION: OVER-CORRECTION THRESHOLD
    # ===================================================================
    _section_header_slide(prs, "The Over-Correction Threshold",
                          "Central Theoretical Contribution")

    slide = _content_slide(prs, "Over-Correction Threshold: Formal Model")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8), [
        ("Let e = baseline CER (fraction), rf = fix rate, ri = introduction rate", ""),
        ("", ""),
        ("Expected CER after correction:", ""),
        ("  CER_corrected = e * (1 - rf) + (1 - e) * ri", ""),
        ("", ""),
        ("Correction is NET BENEFICIAL when:", ""),
        ("  e * rf  >  (1-e) * ri", ""),
        ("  i.e., errors fixed > errors introduced", ""),
        ("", ""),
        ("Threshold (assuming constant rf):", ""),
        ("  e* = ri / (rf + ri)", ""),
        ("", ""),
        ("Empirical estimate:", ""),
        ("  rf ~ 76%, ri ~ 22%", ""),
        ("  e* = 0.22 / (0.76 + 0.22) = 22.4% (aggregate)", ""),
        ("  Per-font crossing: ~3-6% CER (Akhbar 2.12% harmed, Tahoma 5.82% helped)", ""),
    ], font_size=15, color=BLACK)

    _add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.6),
                 "Why the threshold matters", font_size=20, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    _add_bullet_frame(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5), [
        ("1. Explains Akhbar degradation: ", "at 2.12% CER, e < e* so correction is harmful"),
        ("", ""),
        ("2. Explains Andalus improvement: ", "at 12.12% CER with RAG (10.41%) -- high baseline gives enough signal"),
        ("", ""),
        ("3. Augmentation shifts threshold: ", "adding knowledge to prompt raises ri (model becomes more interventionist), raising e*"),
        ("", ""),
        ("4. Phase 4 exception: ", "overcorrection warning directly targets ri, LOWERING it -- explains why Phase 4 uniquely helps Akhbar"),
        ("", ""),
        ("5. Practical implication: ", "deploy LLM correction SELECTIVELY based on estimated OCR quality, not universally"),
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SECTION: HANDWRITTEN
    # ===================================================================
    _section_header_slide(prs, "Why Handwritten & Full-Page Correction Differs",
                          "KHATT, KHATT-Paragraph, Muharaf, Historical")

    slide = _content_slide(prs, "Handwritten vs Typewritten vs Full-Page: Qualitative Differences")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
        ("KHATT (line-level handwritten):", ""),
        ("  ", "37% segmentation errors (merge + split) vs 6.4% on PATS-A01"),
        ("  ", "Segmentation = word boundaries lost -- text-only LLM cannot re-segment"),
        ("  ", "Result: only ~1-3% relative improvement despite 34.24% baseline CER"),
        ("", ""),
        ("KHATT-Paragraph (paragraph handwritten):", ""),
        ("  ", "Paragraph context enables better disambiguation than lines"),
        ("  ", "-16.62 pp genuine text correction (61.68% -> 45.06%)"),
        ("  ", "Best full-page correction result: longer context = more signal"),
        ("", ""),
        ("Muharaf & Historical (degraded/historical):", ""),
        ("  ", "OCR CER >100% due to Qaari runaway bug on full pages"),
        ("  ", "Runaway corrector essential; residual CER still high (89-95%)"),
        ("  ", "Historical degradation + handwriting variability = hard ceiling"),
        ("", ""),
        ("Yarmouk (contemporary print):", ""),
        ("  ", "Most accessible: contemporary Arabic, cleaner OCR (49.85% raw)"),
        ("  ", "T2 achieves 43.35% after correction (-6.5 pp)"),
    ], font_size=13, color=BLACK)

    _add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.6),
                 "What would improve these domains", font_size=18, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    _add_bullet_frame(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5), [
        ("KHATT handwritten:", ""),
        ("  1. ", "Multi-modal: image + OCR text to VLM corrector (Qwen3-VL, GPT-4o-Vision)"),
        ("  2. ", "Segmentation recovery step BEFORE LLM correction"),
        ("", ""),
        ("Historical/degraded documents:", ""),
        ("  1. ", "Fine-tuning on domain-matched historical (OCR, GT) pairs"),
        ("  2. ", "Image preprocessing: binarization, deskew, noise removal"),
        ("", ""),
        ("Full-page all domains:", ""),
        ("  1. ", "Better runaway detection: confidence-score gating at OCR level"),
        ("  2. ", "Page-level context: attend to multiple lines simultaneously"),
        ("  3. ", "Domain-specific RAG index (per-corpus rather than per-font)"),
    ], font_size=13, color=BLACK)

    # ===================================================================
    # SECTION: DESIGN RECOMMENDATIONS
    # ===================================================================
    _section_header_slide(prs, "System Design Recommendations",
                          "From Research Findings to Practical Deployment")

    slide = _content_slide(prs, "Practical Recommendations for Arabic Post-OCR Systems")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5), [
        ("1. Apply correction SELECTIVELY, not universally:", ""),
        ("   ", "Estimate OCR quality first (confidence scores, font ID, calibration CER). Route only CER > 6-7% to LLM. Pass clean output unchanged."),
        ("", ""),
        ("2. Use conf_self (Phase 3 + Phase 4) as default configuration:", ""),
        ("   ", "Best or near-best on all 9 datasets. Statistically significant on 6/9. Both components contribute in ablation."),
        ("", ""),
        ("3. Add CAMeL post-processing ONLY on typewritten text:", ""),
        ("   ", "-0.1pp on PATS-A01 but +0.1pp on KHATT. Omit for handwritten or mixed input."),
        ("", ""),
        ("4. Compute per-font evaluation before deployment:", ""),
        ("   ", "Macro-averages mask font-specific effects. Exclude fonts where correction degrades quality."),
        ("", ""),
        ("5. Recalibrate self-reflective insights for new OCR engines:", ""),
        ("   ", "Phase 4 insights are engine-specific. One-time recomputation per engine from training-split."),
        ("", ""),
        ("6. Use aligned split design for multi-font studies:", ""),
        ("   ", "Identical text content across fonts isolates font as sole variable. Essential for principled comparison."),
    ], font_size=15, color=BLACK)

    # ===================================================================
    # SECTION: LIMITATIONS & FUTURE WORK
    # ===================================================================
    _section_header_slide(prs, "Limitations & Future Work",
                          "Honest Assessment and Research Directions")

    slide = _content_slide(prs, "Limitations")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
        ("Single OCR engine (Qaari):", ""),
        ("  ", "Confusion matrix, self-reflective insights are Qaari-specific"),
        ("  ", "Qualitative findings likely generalise; specific numbers do not"),
        ("", ""),
        ("Single LLM (Qwen3-4B):", ""),
        ("  ", "Larger models may have lower introduction rate -- shifting threshold down"),
        ("  ", "Arabic-specialised models may have stronger character-level priors"),
        ("", ""),
        ("Single text content (PATS-A01):", ""),
        ("  ", "All 8 fonts render same text -- limits content diversity"),
        ("  ", "MSA newspaper/academic text; other domains may differ"),
        ("", ""),
        ("Text-only correction:", ""),
        ("  ", "Document image not used after OCR -- visually ambiguous cases inaccessible"),
    ], font_size=15, color=BLACK)

    _add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.6),
                 "Additional Limitations", font_size=18, bold=True,
                 color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    _add_bullet_frame(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5), [
        ("Diacritic stripping:", ""),
        ("  ", "All metrics computed after stripping diacritics"),
        ("  ", "Diacritic-preserving tasks would change all results"),
        ("", ""),
        ("Morphological DB coverage:", ""),
        ("  ", "CAMeL calima-msa-r13 has gaps: archaic forms, proper nouns, loanwords"),
        ("  ", "May cause spurious reverts in Phase 5"),
        ("", ""),
        ("No fine-tuning baseline:", ""),
        ("  ", "All strategies are prompt-only; fine-tuned comparisons not included"),
    ], font_size=14, color=BLACK)

    # SLIDE — Future Work
    slide = _content_slide(prs, "Future Work")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5), [
        ("1. Characterise threshold across models/engines:", ""),
        ("   ", "Test Qwen3-14B/32B, AraGPT2, Jais, ALLaM on same datasets"),
        ("   ", "Produce 'correction capability curve' as function of model capacity"),
        ("", ""),
        ("2. Fine-tune for conservative OCR correction:", ""),
        ("   ", "Train on (OCR, GT) pairs WITH 'return unchanged' supervision"),
        ("   ", "Directly reduce introduction rate -- lower threshold"),
        ("", ""),
        ("3. Multi-modal correction for handwritten text:", ""),
        ("   ", "VLM corrector (Qwen3-VL, GPT-4o-Vision) with image + OCR text input"),
        ("   ", "Address segmentation errors via visual evidence"),
    ], font_size=15, color=BLACK)

    _add_bullet_frame(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5), [
        ("4. Confidence-guided selective correction:", ""),
        ("   ", "Route only low-confidence OCR tokens to LLM"),
        ("   ", "Reduce introduction rate by restricting intervention"),
        ("", ""),
        ("5. Extension to dialectal Arabic:", ""),
        ("   ", "Egyptian, Gulf, Levantine -- non-standard spelling, code-switching"),
        ("   ", "Dialect-specific self-reflective insights needed"),
        ("", ""),
        ("6. Cross-engine comparison:", ""),
        ("   ", "Tesseract, TrOCR (Arabic), commercial services"),
        ("   ", "Test generality of Phase 3 + Phase 4 approach"),
        ("", ""),
        ("7. Adversarial stress testing:", ""),
        ("   ", "Degraded images, historical documents, low resolution"),
    ], font_size=15, color=BLACK)

    # ===================================================================
    # SECTION: FINAL RANKING & CONCLUSION
    # ===================================================================
    _section_header_slide(prs, "Conclusion & Contributions",
                          "Summary of Findings")

    # SLIDE — Final Ranking (All Experiments)
    slide = _content_slide(prs, "Final Ranking: All Experiments — Best PATS & KHATT Configs")
    final_table = [
        ["Rank", "Experiment / Phase", "Method", "PATS CER", "KHATT CER", "Note"],
        ["1", "Exp2 / T3", "Conservative ZS", "5.27%", "32.93%", "Best PATS ever"],
        ["2", "Exp3 / Gemma-3-12B T2", "Best model, PATS", "5.33%", "35.98%", "Best model PATS"],
        ["3", "Exp2 / T2", "Base + RAG", "5.44%", "33.82%", "Best overall"],
        ["4", "Exp1 / P8", "RAG BM25", "5.45%", "33.83%", "Best Exp1 PATS"],
        ["5", "Exp1 / P9", "Error-Sig RAG", "5.47%", "33.13%", "Best KHATT"],
        ["6", "Exp1 / P4", "Self-Reflective", "5.59%", "33.24%", "Best isolated"],
        ["—", "Exp1 / P1", "OCR Baseline", "5.57%", "34.24%", "Reference"],
        ["—", "Exp2 / T7", "EP+Cons ZS", "190.60%", "51.64%", "CATASTROPHIC"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.3), Inches(12.5), Inches(4.0), final_table)

    _add_bullet_frame(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.7), [
        ("Hierarchy: ", "Conservative prompt (T3) > RAG (P8/T2) > Self-reflective (P4) > Confusion matrix (P3) > Zero-shot (P2) > OCR baseline for PATS"),
        ("KHATT: ", "P9 Error-Sig RAG best (33.13%); T8 best overall (32.54%); T2 best balanced (33.82%)"),
        ("Full-page: ", "Qwen3-4B + T2 + runaway corrector: OCR 86.50% -> 68.18% (-21.2% abs). KHATT-Para 61.68% -> 45.06%."),
    ], font_size=13, color=BLACK)

    # SLIDE — Conclusion
    slide = _content_slide(prs, "Conclusion (3 Experiments, 13 Datasets)")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8), [
        ("Experiment 1 (Phases 1-9, 9 datasets):", ""),
        ("  ", "Zero-shot HARMS PATS avg: 5.57% -> 5.84% (+4.8%). Only Traditional benefits (-18.5%)."),
        ("  ", "Best Exp1 PATS: P8 RAG (5.45%, -0.12% vs OCR baseline). Best KHATT: P9 (33.13%, -1.11%)."),
        ("  ", "Over-correction threshold: ~6-7% CER. Font error TYPE governs outcome, not just rate."),
        ("", ""),
        ("Experiment 2 (8 trials, 13 datasets, 4 new full-page domains):", ""),
        ("  ", "T3 conservative ZS best PATS ever: 5.27% (-5.4% rel, no retrieval). T2 best overall: 5.44% PATS + 33.82% KHATT."),
        ("  ", "Full-page: OCR 86.50% -> T2 68.18% (-21.2% abs after runaway correction)."),
        ("  ", "Error-pattern prompts brittle: T7 = 190.60% PATS (template injection failure)."),
        ("", ""),
        ("Experiment 3 (4 models, 2 OCR sources, 2 held-out benchmarks):", ""),
        ("  ", "Gemma-3-12B best PATS (5.33%); Qwen3-4B best KHATT (33.82%) + full-page (68.18%)."),
        ("  ", "Larger != better: Qwen3-14B worse than 4B after runaway correction."),
        ("  ", "Generalises to unseen benchmarks: RDI-Test-Lines -10.7%, Kitab -17 to -24% relative."),
        ("  ", "Over-correction confirmed on new domains: low-CER subsets harmed."),
    ], font_size=13, color=BLACK)

    # SLIDE — Contributions
    slide = _content_slide(prs, "Thesis Contributions")
    _add_bullet_frame(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5), [
        ("1. Over-correction threshold model: ", "e* ~ 6-7% CER for Qwen3-4B + Qaari. Empirical estimation via fix/intro rate analysis. Explains font-dependent patterns, predicts new domains."),
        ("", ""),
        ("2. Complementarity proof: ", "Confusion-matrix (raises fix rate rf) + Self-reflective (lowers intro rate ri) address orthogonal failure modes, combine near-additively."),
        ("", ""),
        ("3. 13-dataset evaluation: ", "PATS-A01 (8 fonts) + KHATT + KHATT-Para + Yarmouk + Muharaf + Historical + RDI-Test-Lines + Kitab. 2.12% to 205% CER range. 3 experiments."),
        ("", ""),
        ("4. Conservative prompting discovery: ", "T3 (5.27% PATS) beats ALL knowledge augmentation without retrieval overhead. Reducing intervention tendency > providing more targets."),
        ("", ""),
        ("5. Multi-model benchmarking: ", "Qwen3-4B vs 14B vs Gemma-3-4B vs 12B. Qaari vs Gemma VLM OCR. Larger != better. Runaway corrector critical."),
        ("", ""),
        ("6. Three-stage export-infer-analyse pipeline: ", "Reproducible framework: no local GPU needed. JSONL checkpointing. HuggingFace sync. Resume-safe."),
    ], font_size=14, color=BLACK)

    # ===================================================================
    # SLIDE — THANK YOU / Q&A
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
                 "Thank You",
                 font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.8),
                 "Questions & Discussion",
                 font_size=28, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6),
                 "Mohamed Sabry  |  Faculty of Engineering, Cairo University",
                 font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # ===================================================================
    # SAVE
    # ===================================================================
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "defense_presentation.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
